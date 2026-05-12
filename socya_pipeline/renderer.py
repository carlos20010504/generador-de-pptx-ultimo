"""Pure execute renderer with editorial-grade visual design.

Layout system inspired by SIG_Revision_Direccion_2025 (the 'good' reference deck):
- 16:9 widescreen (13.33"x7.5")
- Forest green + gold palette
- Georgia for titles and big KPIs, Calibri for body text
- Eyebrow + accent bar + title + subtitle pattern on each slide
- Footer band with pagination on each slide
- KPI cards with white fill + colored top header band
- No fallbacks. No template fillers.
"""
import io
from pathlib import Path
from typing import List, Optional, Tuple
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ──────────────────── Design system ────────────────────

# Palette (matches reference deck: forest green + gold)
FOREST_DARK = RGBColor(0x1B, 0x3B, 0x2F)   # primary, titles, footer
FOREST_MID  = RGBColor(0x2C, 0x5F, 0x2D)   # KPI accent A
FOREST_LITE = RGBColor(0x4A, 0x7C, 0x3A)   # KPI accent B
GOLD        = RGBColor(0xC9, 0xA2, 0x27)   # eyebrow, KPI accent C, accent line
TERRACOTTA  = RGBColor(0xB8, 0x50, 0x42)   # warning/alert accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_DARK   = RGBColor(0x33, 0x33, 0x33)
PAGE_BG     = RGBColor(0xFA, 0xFA, 0xF5)   # off-white page background
ZEBRA_BG    = RGBColor(0xF2, 0xF0, 0xE6)   # subtle warm gray for zebra rows

KPI_COLORS = [FOREST_MID, FOREST_LITE, GOLD, TERRACOTTA]
CHART_PALETTE = ["#2C5F2D", "#4A7C3A", "#C9A227", "#B85042", "#1B3B2F", "#84B59F"]

HEADING_FONT = "Georgia"
BODY_FONT = "Calibri"

# Layout constants (16:9 = 13.33 x 7.50 inches)
SLIDE_W = 13.33
SLIDE_H = 7.50
MARGIN = 0.5
CONTENT_W = SLIDE_W - 2 * MARGIN  # 12.33

EYEBROW_Y = 0.45
ACCENT_Y = 0.78
TITLE_Y = 0.90
SUBTITLE_Y = 1.85
CONTENT_TOP = 2.45
CONTENT_BOTTOM = 6.30
PAGINATION_Y = 7.05

# ── Chart layout heuristics (named so the intent is clear in code review)
_DOMINANT_RATIO        = 5.0   # leader / second-place — single-dominant pivot
_DONUT_PIE_DOMINANT    = 0.50  # one slice >50% → donut becomes barh
_LONG_TAIL_LEADER_X    = 10    # leader / median(rest) — aggressive grouping
_LONG_TAIL_HEAD_TIGHT  = 5     # top-N kept when leader dominates
_LONG_TAIL_HEAD_LOOSE  = 6     # top-N kept normally
_LONG_TAIL_MAX_RATIO   = 0.20  # tail ≤20% of total → group into "Otros"


def _hex_to_rgb(hex_str) -> Optional[RGBColor]:
    """Convert "#RRGGBB" or "RRGGBB" to RGBColor; return None if invalid."""
    if not hex_str:
        return None
    s = str(hex_str).strip().lstrip("#")
    if len(s) != 6:
        return None
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return None


def _apply_theme_overrides(theme: Optional[dict]) -> dict:
    """Patch the module-level color palette from a theme dict and return
    the originals so the caller can restore them with `_restore_theme`.
    Empty dict means no overrides applied. Theme schema:
      {primary_hex, accent_hex, text_hex, bg_hex}
    Only the colors we know how to map are touched — extras are ignored.
    """
    if not isinstance(theme, dict):
        return {}
    g = globals()
    saved: dict = {}

    primary = _hex_to_rgb(theme.get("primary_hex"))
    accent = _hex_to_rgb(theme.get("accent_hex"))
    text = _hex_to_rgb(theme.get("text_hex"))
    bg = _hex_to_rgb(theme.get("bg_hex"))

    if primary is not None:
        saved["FOREST_DARK"] = g["FOREST_DARK"]
        g["FOREST_DARK"] = primary
        # Derive a slightly lighter mid for KPIs (visual depth without a
        # second hex from the theme).
        saved["FOREST_MID"] = g["FOREST_MID"]
        g["FOREST_MID"] = primary
    if accent is not None:
        saved["GOLD"] = g["GOLD"]
        g["GOLD"] = accent
        # Refresh derived palettes that hold the accent
        saved["KPI_COLORS"] = g["KPI_COLORS"]
        g["KPI_COLORS"] = [g["FOREST_MID"], g["FOREST_LITE"], accent, g["TERRACOTTA"]]
        # CHART_PALETTE uses string hex (matplotlib), not RGBColor
        saved["CHART_PALETTE"] = g["CHART_PALETTE"]
        accent_hex = "#{:02X}{:02X}{:02X}".format(accent[0], accent[1], accent[2])
        primary_hex = "#{:02X}{:02X}{:02X}".format(
            (primary or g["FOREST_DARK"])[0],
            (primary or g["FOREST_DARK"])[1],
            (primary or g["FOREST_DARK"])[2],
        )
        g["CHART_PALETTE"] = [primary_hex, accent_hex,
                              "#4A7C3A", "#B85042", primary_hex, "#84B59F"]
    if text is not None:
        saved["TEXT_DARK"] = g["TEXT_DARK"]
        g["TEXT_DARK"] = text
    if bg is not None:
        saved["PAGE_BG"] = g["PAGE_BG"]
        g["PAGE_BG"] = bg
    return saved


def _restore_theme(saved: dict) -> None:
    """Undo what `_apply_theme_overrides` patched, so the module returns to
    its built-in palette between renders (relevant when running multiple
    renders in the same Python process — tests, dev hot-reload, etc.)."""
    if not saved:
        return
    g = globals()
    for key, value in saved.items():
        g[key] = value


def render_pptx(slides: List[dict], presentation_meta: dict,
                template_path: Path, output_path: Path,
                theme: Optional[dict] = None) -> None:
    """Render slides to a PPTX. presentation_meta drives the title slide + footer.

    `theme` (optional): {primary_hex, accent_hex, text_hex, bg_hex} from the
    UI selector. When provided, swaps the brand palette before rendering and
    restores it after (so the next render can pick a different theme cleanly).
    """
    saved_palette = _apply_theme_overrides(theme)
    try:
        prs = Presentation(str(template_path))
        # Force 16:9 widescreen
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        # Strip pre-existing slides from the template — we control the deck
        while len(prs.slides._sldIdLst) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        blank_layout = _pick_blank_layout(prs)
        deck_title = (presentation_meta.get("title") or "").strip()
        deck_subtitle = (presentation_meta.get("subtitle") or "").strip()
        footer_tagline = _build_footer_tagline(deck_title)

        total = len(slides)
        for i, slide_def in enumerate(slides):
            stype = slide_def.get("type")
            slide = prs.slides.add_slide(blank_layout)
            _add_page_background(slide)

            if stype == "title":
                _add_title_slide(slide, slide_def, deck_title, deck_subtitle,
                                 i + 1, total)
                continue

            eyebrow = slide_def.get("eyebrow") or _eyebrow_from_type(stype, i)
            _add_eyebrow(slide, eyebrow)
            _add_title_block(slide, slide_def.get("title", ""),
                             _slide_subtitle(slide_def))

            if stype == "kpi_row":
                _add_kpi_row(slide, slide_def)
            elif stype == "chart":
                _add_chart_slide(slide, slide_def)
            elif stype == "table":
                _add_table_slide(slide, slide_def)
            elif stype == "text_bullets":
                _add_bullets_slide(slide, slide_def)

            _add_footer(slide, footer_tagline, i + 1, total)

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
    finally:
        _restore_theme(saved_palette)


# ──────────────────── Slide builders ────────────────────

def _add_title_slide(slide, slide_def: dict, deck_title: str, deck_subtitle: str,
                     page: int, total: int) -> None:
    """Cover slide: dark forest BG, gold left accent strip, top eyebrow,
    big serif title, subtitle, divider line, and a date band at the foot.
    Designed to feel like a published report cover."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W), Inches(SLIDE_H))
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = FOREST_DARK

    # Left accent strip — a thin gold column that anchors the layout
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0),
                                     Inches(0.18), Inches(SLIDE_H))
    strip.line.fill.background()
    strip.fill.solid(); strip.fill.fore_color.rgb = GOLD

    eyebrow = slide_def.get("eyebrow") or "PRESENTACIÓN EJECUTIVA"
    _put_text(slide, eyebrow.upper(), Inches(MARGIN + 0.4), Inches(0.75),
              Inches(CONTENT_W), Inches(0.4),
              font=BODY_FONT, font_size=12, bold=True, color=GOLD,
              letter_spacing_em=0.18)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(MARGIN + 0.4), Inches(1.30),
                                   Inches(0.8), Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD

    title_text = slide_def.get("data", {}).get("title") or deck_title or "Presentación"
    # Title size shrinks for long titles so they never overflow
    title_size = 56 if len(title_text) <= 38 else (46 if len(title_text) <= 56 else 38)
    _put_text(slide, title_text, Inches(MARGIN + 0.4), Inches(2.0),
              Inches(CONTENT_W - 0.4), Inches(3.0),
              font=HEADING_FONT, font_size=title_size, bold=True, color=WHITE,
              line_spacing=1.05)

    subtitle_text = slide_def.get("data", {}).get("subtitle") or deck_subtitle
    if subtitle_text:
        _put_text(slide, subtitle_text, Inches(MARGIN + 0.4), Inches(5.2),
                  Inches(CONTENT_W * 0.85), Inches(1.0),
                  font=BODY_FONT, font_size=18,
                  color=RGBColor(0xCF, 0xD7, 0xCB),
                  line_spacing=1.35)

    # Bottom hairline divider — gives the cover that 'report' rhythm
    hr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(MARGIN + 0.4), Inches(SLIDE_H - 0.95),
                                  Inches(CONTENT_W - 0.4), Inches(0.02))
    hr.line.fill.background()
    hr.fill.solid(); hr.fill.fore_color.rgb = RGBColor(0x35, 0x55, 0x47)

    # Date — today's date, computed at render time
    import datetime as _dt
    today = _dt.date.today().strftime("%d.%m.%Y").upper()
    _put_text(slide, today,
              Inches(MARGIN + 0.4), Inches(SLIDE_H - 0.75),
              Inches(3.0), Inches(0.4),
              font=BODY_FONT, font_size=11, bold=True, color=GOLD,
              letter_spacing_em=0.12)

    _put_text(slide, f"{page:02d} / {total:02d}",
              Inches(SLIDE_W - 1.8), Inches(SLIDE_H - 0.75),
              Inches(1.3), Inches(0.4),
              font=BODY_FONT, font_size=11, bold=True, color=GOLD,
              align="right")


def _add_kpi_row(slide, slide_def: dict) -> None:
    """KPI cards. 1-3 KPIs → single row big cards. 4 KPIs → 2x2 grid (each
    card wider, numbers stay readable). Numbers shrink-to-fit so they NEVER wrap."""
    kpis = slide_def.get("data", {}).get("kpis", [])[:6]
    if not kpis:
        return

    n = len(kpis)
    gap = 0.28
    available_h = CONTENT_BOTTOM - CONTENT_TOP - 0.15

    if n <= 3:
        cols = n
        rows = 1
    else:
        cols = 2
        rows = (n + 1) // 2  # 4→2 rows, 5→3, 6→3

    card_w = (CONTENT_W - gap * (cols - 1)) / cols
    card_h = (available_h - gap * (rows - 1)) / rows
    card_h = min(card_h, 3.6)
    number_size = 88 if (cols <= 3 and card_w >= 3.5) else 64

    for i, k in enumerate(kpis):
        col_i = i % cols
        row_i = i // cols
        x = MARGIN + col_i * (card_w + gap)
        y = CONTENT_TOP + row_i * (card_h + gap)
        accent = KPI_COLORS[i % len(KPI_COLORS)]
        _add_kpi_card(slide, x, y, card_w, card_h,
                      header=k.get("label", ""),
                      value=k.get("value", ""),
                      description=k.get("description", ""),
                      accent_color=accent,
                      number_size=number_size)


def _add_kpi_card(slide, x: float, y: float, w: float, h: float,
                  header: str, value: str, description: str,
                  accent_color: RGBColor, number_size: int) -> None:
    """KPI card: white body + colored top band + huge number + (optional)
    description. Layout is computed proportionally so it works for both tall
    cards (1 row × N cards, h~3.6") and short ones (2-row grid, h~1.78")."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    card.adjustments[0] = 0.04
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xDF)
    card.line.width = Pt(0.5)

    # Header band: proportional to card height, capped
    header_h = max(0.36, min(0.5, h * 0.22))
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y),
                                    Inches(w), Inches(header_h))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = accent_color

    if header:
        _put_text(slide, header.upper(),
                  Inches(x + 0.2), Inches(y),
                  Inches(w - 0.4), Inches(header_h),
                  font=BODY_FONT, font_size=10 if h < 2.0 else 11, bold=True,
                  color=WHITE, align="left", anchor="middle",
                  letter_spacing_em=0.08)

    # Decide whether description fits. For short cards (grid 2x2), use a tight
    # footer line; for tall cards (single-row, n≤3), use a roomier description.
    show_desc = bool(description) and h >= 1.5
    if show_desc:
        desc_h = 0.45 if h < 2.4 else 0.7
    else:
        desc_h = 0.0

    # Value box gets all remaining vertical space below the header.
    value_y = y + header_h + 0.05
    value_h = max(0.5, h - header_h - desc_h - 0.1)
    fitted_size = _fit_kpi_font(str(value), max_pt=number_size,
                                 card_w_in=w - 0.4)
    # Also cap by available height: an em ~ 1.0 line height. Don't exceed value_h.
    height_cap_pt = int(value_h * 72 * 0.85)
    fitted_size = min(fitted_size, height_cap_pt)
    fitted_size = max(28, fitted_size)
    _put_text(slide, str(value),
              Inches(x + 0.2), Inches(value_y),
              Inches(w - 0.4), Inches(value_h),
              font=HEADING_FONT, font_size=fitted_size, bold=True,
              color=accent_color, align="center", anchor="middle",
              line_spacing=1.0)

    if show_desc:
        desc_font = 10 if h < 2.4 else 11
        _put_text(slide, description,
                  Inches(x + 0.3), Inches(y + h - desc_h - 0.05),
                  Inches(w - 0.6), Inches(desc_h),
                  font=BODY_FONT, font_size=desc_font, color=TEXT_GRAY,
                  align="center", anchor="middle", line_spacing=1.2)


def _fit_kpi_font(value: str, max_pt: float, card_w_in: float) -> int:
    """Find the largest integer font size (in pt) for `value` (Georgia bold)
    that fits in `card_w_in` inches without wrapping. Georgia bold uppercase
    'M' and digits average ~0.78em — generous so we never overflow."""
    char_em_avg = 0.78
    available_pt = card_w_in * 72.0
    n_chars = max(1, len(value))
    fit = available_pt / (char_em_avg * n_chars)
    return int(max(28, min(int(max_pt), int(fit))))


def _add_chart_slide(slide, slide_def: dict) -> None:
    """Chart on the left, narrative card on the right."""
    data = slide_def.get("data", {})
    narrative = (slide_def.get("narrative") or "").strip()

    chart_w = 7.7
    chart_h = CONTENT_BOTTOM - CONTENT_TOP - 0.1
    chart_x = MARGIN
    chart_y = CONTENT_TOP

    # Polished PNG includes a transparent padding ring for the drop shadow
    # (matplotlib dpi=160 → padding in inches = pad_px / 160). We expand the
    # picture box to match so the chart itself stays at its original size.
    img = _build_chart_png(data, width_in=chart_w, height_in=chart_h)
    pad_in = CHART_POLISH_PADDING_PX / 160.0
    slide.shapes.add_picture(
        img,
        Inches(chart_x - pad_in), Inches(chart_y - pad_in),
        Inches(chart_w + 2 * pad_in), Inches(chart_h + 2 * pad_in),
    )

    if narrative:
        card_x = MARGIN + chart_w + 0.3
        card_w = CONTENT_W - chart_w - 0.3
        card_y = CONTENT_TOP
        card_h = chart_h
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(card_x), Inches(card_y),
                                        Inches(card_w), Inches(card_h))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xDF)
        card.line.width = Pt(0.5)
        header_h = 0.55
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(card_x), Inches(card_y),
                                        Inches(card_w), Inches(header_h))
        band.line.fill.background()
        band.fill.solid(); band.fill.fore_color.rgb = FOREST_MID
        _put_text(slide, "LECTURA EJECUTIVA",
                  Inches(card_x + 0.25), Inches(card_y),
                  Inches(card_w - 0.5), Inches(header_h),
                  font=BODY_FONT, font_size=11, bold=True, color=WHITE,
                  anchor="middle", letter_spacing_em=0.08)
        _put_text(slide, narrative,
                  Inches(card_x + 0.3), Inches(card_y + header_h + 0.2),
                  Inches(card_w - 0.6), Inches(card_h - header_h - 0.4),
                  font=BODY_FONT, font_size=14, color=TEXT_DARK,
                  line_spacing=1.4, anchor="top")


def _add_table_slide(slide, slide_def: dict) -> None:
    """Table with colored header band, zebra-striped rows, opcional totals
    row al final, y heatmap-style cell coloring para columnas numéricas."""
    data = slide_def.get("data", {})
    headers = data.get("headers", [])
    rows = data.get("rows", [])
    totals_row = data.get("totals_row")
    numeric_col_indices = set(data.get("numeric_col_indices", []) or [])
    numeric_col_ranges = data.get("numeric_col_ranges", {}) or {}
    if not headers or not rows:
        return

    max_cols = 6
    max_rows = 12
    headers = [_clean_header(h) for h in headers[:max_cols]]
    rows = [list(r)[:max_cols] for r in rows[:max_rows]]
    n_cols = len(headers)
    has_totals = bool(totals_row)
    body_rows_n = len(rows)
    n_rows = body_rows_n + 1 + (1 if has_totals else 0)

    table_x = MARGIN
    table_y = CONTENT_TOP
    table_w = CONTENT_W
    available_h = CONTENT_BOTTOM - CONTENT_TOP - 0.2
    row_h = min(0.4, available_h / n_rows)
    table_h = row_h * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                           Inches(table_x), Inches(table_y),
                                           Inches(table_w), Inches(table_h))
    tbl = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = FOREST_DARK
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        # Numeric columns alineadas a derecha
        p.alignment = PP_ALIGN.RIGHT if j in numeric_col_indices else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(h).upper()
        run.font.name = BODY_FONT
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            # Heatmap shading para columnas numéricas (override del zebra).
            # Bajo intensidad — la idea es ayudar a escanear, no gritar.
            heat_color = None
            if j in numeric_col_indices and j in numeric_col_ranges:
                heat_color = _heatmap_color(val, numeric_col_ranges[j])
            cell.fill.solid()
            if heat_color is not None:
                cell.fill.fore_color.rgb = heat_color
            else:
                cell.fill.fore_color.rgb = ZEBRA_BG if i % 2 == 0 else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j in numeric_col_indices else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "" if val is None else str(val)
            run.font.name = BODY_FONT
            run.font.size = Pt(10)
            run.font.color.rgb = TEXT_DARK

    # Totals row
    if has_totals:
        row_idx = body_rows_n + 1
        for j, val in enumerate(totals_row[:n_cols]):
            cell = tbl.cell(row_idx, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = FOREST_DARK
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j in numeric_col_indices else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "" if val is None else str(val)
            run.font.name = BODY_FONT
            run.font.size = Pt(10.5)
            run.font.bold = True
            run.font.color.rgb = WHITE


def _heatmap_color(display_val, value_range) -> Optional[RGBColor]:
    """Mapea un valor a un tono pastel verde→amarillo→rojo según su posición
    en [min, max]. Devuelve None si no se puede parsear el valor o si el
    rango es degenerado.

    Convención semántica neutra (no asume "alto = bueno"): usamos verde
    para valores ALTOS (típicamente magnitudes mayores) y blanco-cremoso
    para valores BAJOS. Si el contexto requiere otra convención (ej.
    "tasa de rechazo" donde alto = malo), eso necesitaría una hint del
    inventory — fuera de scope.
    """
    try:
        # display_val viene formateado ($1.5M, 23%, etc.). Extraer el numérico.
        import re as _re
        s = str(display_val or "")
        m = _re.search(r"-?\d+(?:[.,]\d+)?", s.replace(".", "").replace(",", "."))
        # Mejor: limpiar puntos de miles y coma decimal — heurística simple
        cleaned = _re.sub(r"[^\d.\-]", "", s.replace(",", "."))
        if not cleaned or cleaned in ("-", "."):
            return None
        # Si tenía sufijo K/M/B aplicarlo
        v = float(cleaned)
        upper = s.upper()
        if "B" in upper:
            v *= 1_000_000_000
        elif "M" in upper and "MIL" not in upper:
            v *= 1_000_000
        elif "K" in upper:
            v *= 1_000
        lo, hi = float(value_range[0]), float(value_range[1])
        if hi <= lo:
            return None
        # Posición [0, 1]
        t = max(0.0, min(1.0, (v - lo) / (hi - lo)))
        # Pastel desde crema (#FAFAF5 ≈ PAGE_BG) hasta verde Socya (#69BE28)
        # con muy baja saturación para no gritar.
        # cream → soft green
        # cream  rgb(250, 250, 245)
        # green  rgb(105, 190, 40)  — usamos sólo 25% de intensidad
        # Mezcla: dest = cream + t * (green - cream) * 0.30
        cr, cg, cb = 250, 250, 245
        gr, gg, gb = 105, 190, 40
        intensity = 0.30  # mantener sutil
        r = int(cr + (gr - cr) * t * intensity)
        g = int(cg + (gg - cg) * t * intensity)
        b = int(cb + (gb - cb) * t * intensity)
        return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
    except Exception:
        return None


def _add_bullets_slide(slide, slide_def: dict) -> None:
    """Bullets rendered as cards in a single column with gold left accent."""
    bullets = slide_def.get("data", {}).get("bullets", [])[:6]
    if not bullets:
        return

    card_w = CONTENT_W
    available_h = CONTENT_BOTTOM - CONTENT_TOP - 0.3
    n = len(bullets)
    gap = 0.18
    card_h = (available_h - gap * (n - 1)) / n
    card_h = min(card_h, 1.0)
    card_h = max(card_h, 0.5)

    for i, bullet in enumerate(bullets):
        y = CONTENT_TOP + i * (card_h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(MARGIN), Inches(y),
                                        Inches(card_w), Inches(card_h))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xDF)
        card.line.width = Pt(0.5)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(MARGIN), Inches(y),
                                          Inches(0.12), Inches(card_h))
        accent.line.fill.background()
        accent.fill.solid(); accent.fill.fore_color.rgb = GOLD
        _put_text(slide, f"{i + 1:02d}",
                  Inches(MARGIN + 0.3), Inches(y),
                  Inches(0.7), Inches(card_h),
                  font=HEADING_FONT, font_size=22, bold=True, color=FOREST_DARK,
                  align="left", anchor="middle")
        _put_text(slide, bullet,
                  Inches(MARGIN + 1.1), Inches(y),
                  Inches(card_w - 1.4), Inches(card_h),
                  font=BODY_FONT, font_size=14, color=TEXT_DARK,
                  align="left", anchor="middle", line_spacing=1.3)


# ──────────────────── Common chrome ────────────────────

def _add_page_background(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W), Inches(SLIDE_H))
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = PAGE_BG


def _add_eyebrow(slide, text: str) -> None:
    if text:
        _put_text(slide, text.upper(),
                  Inches(MARGIN), Inches(EYEBROW_Y),
                  Inches(CONTENT_W), Inches(0.3),
                  font=BODY_FONT, font_size=10, bold=True, color=GOLD,
                  letter_spacing_em=0.12)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(MARGIN), Inches(ACCENT_Y),
                                   Inches(0.4), Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD


def _add_title_block(slide, title: str, subtitle: str) -> None:
    if title:
        _put_text(slide, _ellipsize(title, 110),
                  Inches(MARGIN), Inches(TITLE_Y),
                  Inches(CONTENT_W), Inches(0.95),
                  font=HEADING_FONT, font_size=28, bold=True, color=FOREST_DARK,
                  line_spacing=1.1)
    if subtitle:
        _put_text(slide, _ellipsize(subtitle, 220),
                  Inches(MARGIN), Inches(SUBTITLE_Y),
                  Inches(CONTENT_W), Inches(0.55),
                  font=BODY_FONT, font_size=12, color=TEXT_GRAY,
                  line_spacing=1.3)


def _ellipsize(text, max_chars: int) -> str:
    """Trim long titles/subtitles with a visible ellipsis instead of letting
    python-pptx truncate them silently to whatever fits the placeholder.
    Honest truncation > stealth clipping."""
    if text is None:
        return ""
    s = str(text).strip()
    if len(s) <= max_chars:
        return s
    # Cut on a word boundary if reasonable
    cutoff = s.rfind(" ", 0, max_chars - 1)
    if cutoff < max_chars - 20:
        cutoff = max_chars - 1
    return s[:cutoff].rstrip(",.;:- ") + "…"


def _add_footer(slide, tagline: str, page: int, total: int) -> None:
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(MARGIN), Inches(PAGINATION_Y - 0.05),
                                   Inches(0.4), Inches(0.04))
    sep.line.fill.background()
    sep.fill.solid(); sep.fill.fore_color.rgb = GOLD

    _put_text(slide, tagline,
              Inches(MARGIN + 0.6), Inches(PAGINATION_Y),
              Inches(CONTENT_W - 1.6), Inches(0.3),
              font=BODY_FONT, font_size=9, color=TEXT_GRAY,
              letter_spacing_em=0.05)

    _put_text(slide, f"{page} / {total}",
              Inches(SLIDE_W - 1.3), Inches(PAGINATION_Y),
              Inches(0.8), Inches(0.3),
              font=BODY_FONT, font_size=9, bold=True, color=FOREST_DARK,
              align="right")


# ──────────────────── Charts (matplotlib) ────────────────────

def _add_mean_reference_line(ax, values: list, *, axis: str = "y",
                                color: str = "#6B6B6B",
                                show_threshold: float = 0.10) -> None:
    """Dibuja una línea horizontal (axis='y') o vertical (axis='x') al
    promedio cuando la varianza es significativa (≥ `show_threshold` * mean).
    Si la serie es plana (todo similar al promedio) no aporta — skip.
    Mantiene la línea sutil para no robar atención al data principal.
    """
    try:
        nums = [float(v) for v in values if v is not None]
        if len(nums) < 4:
            return
        mean = sum(nums) / len(nums)
        if mean <= 0:
            return
        max_v = max(nums)
        min_v = min(nums)
        rng = max_v - min_v
        # Skip si la serie es prácticamente plana
        if rng < mean * show_threshold:
            return
        if axis == "y":
            ax.axhline(y=mean, color=color, linestyle="--", linewidth=1.2,
                        alpha=0.55, zorder=1)
            # Etiqueta a la derecha del eje
            xlim = ax.get_xlim()
            ax.text(xlim[1], mean, f" prom: {_format_value_short(mean)}",
                      va="center", ha="left", fontsize=8.5,
                      fontfamily="Calibri", color=color, alpha=0.85)
        else:
            ax.axvline(x=mean, color=color, linestyle="--", linewidth=1.2,
                        alpha=0.55, zorder=1)
    except Exception:
        # Polish, nunca rompe.
        pass


def _add_variance_band(ax, values: list, color: str = "#1B3B2F") -> None:
    """En time-series largos (>=8 puntos), sombrea ±1σ alrededor del
    promedio. Visualiza qué tan volátil es la serie respecto a su nivel
    típico — útil cuando la línea zigzaguea pero el AI dice 'estable'.
    """
    try:
        nums = [float(v) for v in values if v is not None]
        if len(nums) < 8:
            return
        mean = sum(nums) / len(nums)
        var = sum((v - mean) ** 2 for v in nums) / max(1, (len(nums) - 1))
        sd = var ** 0.5
        if sd <= 0 or sd / max(abs(mean), 1) < 0.05:
            return  # serie estable — band aporta poco
        x = list(range(len(nums)))
        upper = mean + sd
        lower = max(0.0, mean - sd)
        ax.fill_between(x, [lower] * len(x), [upper] * len(x),
                          color=color, alpha=0.06, zorder=0)
    except Exception:
        pass


def _annotate_peak_valley_line(ax, values: list) -> None:
    """Highlight peak y valley visualmente en una línea cuando NO son los
    endpoints (que ya tienen anotación) y la diferencia con el endpoint es
    significativa (≥10% del rango). Marca con punto dorado y label corto.
    """
    if not values or len(values) < 4:
        return
    try:
        max_v = max(values)
        min_v = min(values)
        rng = max_v - min_v
        if rng <= 0:
            return
        max_idx = values.index(max_v)
        min_idx = values.index(min_v)
        last_idx = len(values) - 1

        def _is_significant(idx: int) -> bool:
            # No re-anotamos el endpoint
            if idx == last_idx:
                return False
            # Si el extremo está "cerca" del endpoint (mismo orden de
            # magnitud que el endpoint y a 1-2 puntos de distancia), skip.
            return abs(values[idx] - values[last_idx]) >= rng * 0.10

        if _is_significant(max_idx):
            ax.plot([max_idx], [max_v], marker="o", markersize=11,
                     markerfacecolor="#C9A227", markeredgecolor="#1B3B2F",
                     markeredgewidth=1.5, zorder=5)
            ax.annotate(f"▲ {_format_value_short(max_v)}",
                          xy=(max_idx, max_v),
                          xytext=(0, 14), textcoords="offset points",
                          ha="center", fontsize=9, fontfamily="Calibri",
                          fontweight="bold", color="#1B3B2F")
        if _is_significant(min_idx) and min_idx != max_idx:
            ax.plot([min_idx], [min_v], marker="o", markersize=11,
                     markerfacecolor="#B85042", markeredgecolor="#1B3B2F",
                     markeredgewidth=1.5, zorder=5)
            ax.annotate(f"▼ {_format_value_short(min_v)}",
                          xy=(min_idx, min_v),
                          xytext=(0, -22), textcoords="offset points",
                          ha="center", fontsize=9, fontfamily="Calibri",
                          fontweight="bold", color="#1B3B2F")
    except Exception:
        # Anotación es polish — nunca debe romper el render.
        pass


def _build_chart_png(data: dict, width_in: float = 7.5,
                       height_in: float = 4.0) -> io.BytesIO:
    chart_type = data.get("chart_type", "bar")
    raw_labels = [str(l) for l in data.get("labels", [])]
    values = list(data.get("values", []))
    is_chronological = bool(data.get("chronological"))

    # Detect single-dominant on RAW values (before "Otros" grouping inflates the
    # 2nd bar). This is the key signal for pivoting layout.
    nums_for_check = [float(v) for v in values if v is not None]
    raw_dominant = (chart_type in ("bar", "pie") and not is_chronological
                     and _is_single_dominant(nums_for_check))

    # A pie/donut where one slice dominates (>50%) is visually unhelpful —
    # the small slices become slivers. Convert to horizontal bar where each
    # category gets its own readable line.
    if chart_type == "pie" and nums_for_check:
        total = sum(nums_for_check) or 1
        if max(nums_for_check) / total > _DONUT_PIE_DOMINANT:
            chart_type = "barh"
            raw_dominant = False  # already pivoted; skip the bar→barh path

    if not is_chronological:
        raw_labels, values = _group_long_tail(raw_labels, values, chart_type)

    if raw_dominant:
        chart_type = "barh"
    if chart_type == "bar" and is_chronological:
        # Chronological "bar": only 2-4 periods → bar (single points read poorly
        # on a line chart). 5+ periods → line trend.
        if len(values) >= 5:
            chart_type = "line"

    # Truncate long labels (less aggressive for horizontal: more horizontal room)
    label_max = 36 if chart_type == "barh" else 18
    labels = [(l if len(l) <= label_max else l[:label_max - 1] + "…") for l in raw_labels]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=160,
                             facecolor="#FAFAF5")
    ax.set_facecolor("#FAFAF5")

    n = len(labels)
    colors = _build_chart_colors(n, chart_type, chronological=is_chronological)

    if chart_type == "pie":
        wedges, _texts = ax.pie(
            values, labels=None,
            colors=colors, startangle=90, counterclock=False,
            wedgeprops={"linewidth": 2, "edgecolor": "#FAFAF5", "width": 0.42},
        )
        total = sum(values) or 1
        legend_labels = [f"{lbl}  ·  {_format_value_short(v)}  ·  {v/total*100:.0f}%"
                          for lbl, v in zip(labels, values)]
        ax.legend(wedges, legend_labels, loc="center left",
                  bbox_to_anchor=(1.05, 0.5), frameon=False,
                  fontsize=11, labelcolor="#333333",
                  prop={"family": "Calibri"})
        # Center total in the donut hole
        ax.text(0, 0, _format_value_short(total),
                ha="center", va="center",
                fontsize=22, fontfamily="Georgia", fontweight="bold",
                color="#1B3B2F")
        ax.text(0, -0.18, "TOTAL",
                ha="center", va="center",
                fontsize=8.5, fontfamily="Calibri",
                color="#6B6B6B")

    elif chart_type == "line":
        # Variance band PRIMERO (zorder=0) para que quede DETRÁS de la línea.
        _add_variance_band(ax, values, color=CHART_PALETTE[0])
        ax.plot(labels, values, marker="o", color=CHART_PALETTE[0],
                 linewidth=3.0, markersize=8, markerfacecolor=CHART_PALETTE[2],
                 markeredgecolor=CHART_PALETTE[0], markeredgewidth=2)
        # Soft fill under the line for visual weight
        ax.fill_between(range(len(labels)), values, alpha=0.08,
                          color=CHART_PALETTE[0])
        _style_axes(ax, integer_y=_all_integers(values))
        # Reference line al promedio — sólo si la serie tiene varianza
        # suficiente para que aporte (heurística dentro del helper).
        _add_mean_reference_line(ax, values)
        ax.tick_params(axis="x", rotation=25, labelsize=10)
        # Endpoint annotation
        if values:
            ax.annotate(_format_value_short(values[-1]),
                          xy=(len(values) - 1, values[-1]),
                          xytext=(8, 8), textcoords="offset points",
                          fontsize=10, fontfamily="Calibri",
                          color="#1B3B2F", fontweight="bold")
        # Peak / valley callouts — sólo si NO coinciden con el endpoint
        # (que ya tiene anotación) y si la diferencia es significativa.
        _annotate_peak_valley_line(ax, values)

    elif chart_type == "barh":
        # Horizontal bars: top of chart = leader, bottom = smallest / "Otros".
        ypos = list(range(len(labels)))
        bars = ax.barh(ypos, values, color=colors, edgecolor="#FAFAF5",
                          linewidth=1.5, height=0.7)
        ax.set_yticks(ypos)
        ax.set_yticklabels(labels, fontfamily="Calibri", fontsize=10,
                            color="#333333")
        ax.invert_yaxis()  # so labels[0] (the leader) shows at the top
        # Move y-axis to right? No — keep left. Hide x ticks/labels — values shown next to bars.
        ax.tick_params(axis="x", colors="#FAFAF5", labelcolor="#FAFAF5",
                          labelsize=1, length=0)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["bottom"].set_visible(False)
        ax.spines["left"].set_color("#D5D5CD")
        max_v = max(values) if values else 1
        for b, v in zip(bars, values):
            ax.text(b.get_width() + max_v * 0.012,
                      b.get_y() + b.get_height() / 2,
                      _format_value_short(v),
                      ha="left", va="center", fontsize=11,
                      fontfamily="Calibri", color="#1B3B2F", fontweight="bold")
        # Add headroom on x for labels
        ax.set_xlim(0, max_v * 1.18)

    elif chart_type == "histogram":
        # Bins ya pre-calculados en extractor (_build_histogram_data).
        # Renderizamos como barras adyacentes (sin gap) — visualmente "es"
        # un histograma. Mostramos línea de mediana sobreimpuesta cuando
        # los stats están disponibles.
        x_pos = list(range(len(labels)))
        ax.bar(x_pos, values, color=CHART_PALETTE[0], edgecolor="#FAFAF5",
                linewidth=0.8, width=1.0, align="edge")
        ax.set_xticks([i + 0.5 for i in x_pos])
        # Sólo mostrar etiquetas alternas si hay >8 bins (espacio)
        step = max(1, len(labels) // 8)
        tick_labels = [l if i % step == 0 else "" for i, l in enumerate(labels)]
        ax.set_xticklabels(tick_labels, fontfamily="Calibri", fontsize=9,
                            color="#333333", rotation=30, ha="right")
        _style_axes(ax, integer_y=True)
        # Mediana visible
        stats = data.get("_stats") or {}
        median = stats.get("median")
        if median is not None:
            try:
                # Find which bin the median falls into for x positioning
                med_x = None
                for i, lbl in enumerate(labels):
                    parts = lbl.split("–")
                    if len(parts) == 2:
                        # Convert compact format back is hard; usar índice
                        # como aproximación por mitad de bins.
                        med_x = len(labels) / 2
                        break
                if med_x is not None:
                    ax.axvline(x=med_x, color="#C9A227", linestyle="--",
                                linewidth=1.5, alpha=0.85, zorder=3)
                    ax.text(med_x, max(values) * 0.95,
                              f" Mediana: {_format_value_short(median)}",
                              ha="left", va="top", fontsize=9,
                              fontfamily="Calibri", color="#1B3B2F",
                              fontweight="bold")
            except Exception:
                pass
        ax.set_ylabel("Frecuencia", fontsize=10, color="#6B6B6B",
                       fontfamily="Calibri")

    else:  # vertical bar
        x_pos = list(range(len(labels)))
        bars = ax.bar(x_pos, values, color=colors, edgecolor="#FAFAF5",
                       linewidth=1.5, width=0.62)
        ax.set_xticks(x_pos)
        ax.set_xticklabels(labels, fontfamily="Calibri", fontsize=10,
                            color="#333333")
        rotate = 25 if any(len(l) > 8 for l in labels) else 0
        ax.tick_params(axis="x", rotation=rotate, labelsize=10)
        _style_axes(ax, integer_y=_all_integers(values))
        # Reference line al promedio (sólo si aporta — ver heurística).
        # Útil para que el lector ubique cada barra vs la media de un vistazo.
        _add_mean_reference_line(ax, values)
        max_v = max(values) if values else 1
        peak_idx = values.index(max_v) if values else -1
        for i, (b, v) in enumerate(zip(bars, values)):
            ax.text(b.get_x() + b.get_width() / 2,
                      b.get_height() + max_v * 0.015,
                      _format_value_short(v),
                      ha="center", va="bottom", fontsize=11,
                      fontfamily="Calibri", color="#1B3B2F", fontweight="bold")
            # Highlight visual del peak — corona pequeña y borde reforzado.
            # Sólo cuando hay ≥4 barras (en menos no aporta).
            if i == peak_idx and len(values) >= 4:
                b.set_edgecolor("#C9A227")
                b.set_linewidth(2.5)
                ax.text(b.get_x() + b.get_width() / 2,
                          b.get_height() + max_v * 0.085,
                          "★ MÁX",
                          ha="center", va="bottom", fontsize=8.5,
                          fontfamily="Calibri", fontweight="bold",
                          color="#C9A227")
        ax.set_ylim(0, max_v * 1.18)

    try:
        plt.tight_layout(pad=0.6)
    except Exception:
        pass
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                 facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    # Apply Pillow polish (rounded corners + soft drop shadow). If Pillow
    # fails for any reason we just use the raw matplotlib PNG.
    polished, _padding_px = _apply_chart_polish(buf)
    return polished


# ──── Pillow post-processing ────

# Padding (in matplotlib output px) added around the chart for the shadow.
# Charts are saved at dpi=160 inside _build_chart_png. 30 px ≈ 0.19" of slack.
CHART_POLISH_PADDING_PX = 30


def _apply_chart_polish(buf: io.BytesIO,
                          radius_px: int = 14,
                          shadow_blur: int = 18,
                          shadow_offset=(0, 8),
                          shadow_opacity: float = 0.18,
                          padding_px: int = CHART_POLISH_PADDING_PX
                          ) -> Tuple[io.BytesIO, int]:
    """Take a chart PNG (matplotlib output), return a polished version with
    rounded corners and a soft drop shadow on a transparent canvas.

    The output canvas is `padding_px` larger on every side to fit the shadow
    spill. Returns (buffer, padding_px) so the caller can adjust the picture
    insertion box accordingly. If Pillow fails, returns the original buffer
    with padding_px=0 so downstream geometry stays correct."""
    try:
        from PIL import Image, ImageFilter, ImageDraw
    except Exception:
        buf.seek(0)
        return buf, 0
    try:
        img = Image.open(buf).convert("RGBA")
        w, h = img.size

        # Rounded-corner mask applied to the chart's alpha channel
        mask = Image.new("L", (w, h), 0)
        ImageDraw.Draw(mask).rounded_rectangle(
            (0, 0, w, h), radius=radius_px, fill=255)
        img.putalpha(mask)

        # Final canvas accommodates shadow spill in every direction
        canvas_w = w + 2 * padding_px
        canvas_h = h + 2 * padding_px
        canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))

        # Build shadow on its own layer so we can blur it without touching the chart
        shadow = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
        sd = ImageDraw.Draw(shadow)
        sx = padding_px + shadow_offset[0]
        sy = padding_px + shadow_offset[1]
        alpha = max(0, min(255, int(255 * shadow_opacity)))
        sd.rounded_rectangle(
            (sx, sy, sx + w, sy + h),
            radius=radius_px,
            fill=(0, 0, 0, alpha),
        )
        shadow = shadow.filter(ImageFilter.GaussianBlur(shadow_blur))
        canvas = Image.alpha_composite(canvas, shadow)

        # Paste the rounded chart on top
        canvas.paste(img, (padding_px, padding_px), img)

        out = io.BytesIO()
        canvas.save(out, format="PNG", optimize=True)
        out.seek(0)
        return out, padding_px
    except Exception:
        buf.seek(0)
        return buf, 0


# ──── Chart helpers ────

def _is_single_dominant(values) -> bool:
    """Leader is ≥_DOMINANT_RATIO× the second-place. Vertical bar would
    show invisible slivers; horizontal layout reads better."""
    nums = [float(v) for v in values if v is not None]
    if len(nums) < 2:
        return False
    s = sorted(nums, reverse=True)
    if s[1] <= 0:
        return True
    return s[0] / s[1] >= _DOMINANT_RATIO


def _all_integers(values) -> bool:
    try:
        return all(float(v).is_integer() for v in values if v is not None)
    except (TypeError, ValueError):
        return False


def _group_long_tail(labels, values, chart_type):
    """Collapse the long tail of a distribution into 'Otros' so the chart
    stays readable.

    Two strategies:
      - Standard: >7 cats and tail ≤20% of total → keep top 6, group rest
      - Aggressive: top value >10x the others' median → keep top 5 plus
        'Otros'. Without this, charts with one massive bar show the others
        as invisible slivers (e.g. CONTABILIZADO=1500 vs ESPERA=8).
    """
    if len(labels) <= 5:
        return labels, values
    nums = [float(v) for v in values]
    total = sum(nums) or 1

    # Aggressive grouping when the leader dwarfs the rest
    rest = sorted(nums[1:], reverse=True)
    median_rest = rest[len(rest) // 2] if rest else 0
    leader_dominates = (median_rest > 0
                          and nums[0] / median_rest >= _LONG_TAIL_LEADER_X)

    head_n = _LONG_TAIL_HEAD_TIGHT if (leader_dominates and len(labels) > _LONG_TAIL_HEAD_TIGHT) else _LONG_TAIL_HEAD_LOOSE
    if len(labels) <= head_n:
        return labels, values

    head_labels = labels[:head_n]
    head_vals = nums[:head_n]
    tail_vals = nums[head_n:]
    if not tail_vals:
        return labels, values
    tail_total = sum(tail_vals)
    tail_pct = tail_total / total
    if leader_dominates or tail_pct < _LONG_TAIL_MAX_RATIO or chart_type == "pie":
        return head_labels + ["Otros"], head_vals + [tail_total]
    return labels, values


def _build_chart_colors(n: int, chart_type: str, chronological: bool = False):
    """For categorical (ranked) bars: highlight #1, fade the rest. For chronological
    bars (months/quarters): use the same primary color for ALL — there's no
    'leader' in a time series. For pies/lines: rotate full palette."""
    if chart_type in ("bar", "barh") and n >= 2 and not chronological:
        primary = CHART_PALETTE[0]
        secondary = CHART_PALETTE[1]
        muted = "#A9B7A1"
        return [primary] + [secondary if i == 1 else muted for i in range(1, n)]
    if chart_type in ("bar", "barh") and chronological:
        return [CHART_PALETTE[0]] * n
    return (CHART_PALETTE * ((n // len(CHART_PALETTE)) + 1))[:n]


def _style_axes(ax, integer_y: bool = False) -> None:
    ax.tick_params(axis="x", labelsize=10, colors="#6B6B6B")
    ax.tick_params(axis="y", labelsize=10, colors="#6B6B6B")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#D5D5CD")
    ax.spines["bottom"].set_color("#D5D5CD")
    ax.grid(axis="y", linestyle="--", linewidth=0.5, color="#D5D5CD", alpha=0.6)
    if integer_y:
        from matplotlib.ticker import MaxNLocator
        ax.yaxis.set_major_locator(MaxNLocator(integer=True))
    # Suppress matplotlib's "1e7" offset/scientific notation; format K/M/B.
    from matplotlib.ticker import FuncFormatter
    ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: _format_value_short(v)))
    try:
        ax.ticklabel_format(axis="y", style="plain", useOffset=False)
    except Exception:
        pass
    ax.yaxis.get_offset_text().set_visible(False)


def _format_value_short(v) -> str:
    """Compact axis/annotation formatter (delegates to insights.format_compact)."""
    from socya_pipeline import insights
    return insights.format_compact(v, fallback="0")


# ──────────────────── Helpers ────────────────────

def _pick_blank_layout(prs):
    """Find a truly empty layout for our hand-built slides. We CANNOT fall
    back to the last layout blindly — many templates put rich placeholders on
    the last layout, which then bleed onto our generated slides as ghost
    titles/footers. Strategy:
      1. Layout named "Blank"/"blanco" → use it.
      2. Else the layout with the FEWEST placeholders (closest to blank).
    """
    for layout in prs.slide_layouts:
        name = (getattr(layout, "name", "") or "").lower()
        if "blank" in name or "blanco" in name or "vacío" in name or "vacio" in name:
            return layout
    # Fallback: pick the layout with the smallest placeholder count (often
    # zero or just a title placeholder we'll override anyway).
    try:
        return min(prs.slide_layouts, key=lambda lay: len(list(lay.placeholders)))
    except Exception:
        return prs.slide_layouts[-1]


def _eyebrow_from_type(stype: str, idx: int) -> str:
    return {
        "kpi_row": "INDICADORES CLAVE",
        "chart": "ANÁLISIS VISUAL",
        "table": "DETALLE OPERATIVO",
        "text_bullets": "HALLAZGOS Y CONCLUSIONES",
    }.get(stype, f"SECCIÓN {idx + 1:02d}")


def _slide_subtitle(slide_def: dict) -> str:
    return (slide_def.get("subtitle")
            or slide_def.get("data", {}).get("subtitle")
            or "")


def _build_footer_tagline(deck_title: str) -> str:
    parts = []
    if deck_title:
        parts.append(deck_title)
    parts.append("PRESENTACIÓN GENERADA AUTOMÁTICAMENTE")
    return "  ·  ".join(parts).upper()


def _clean_header(value) -> str:
    s = str(value or "").strip()
    if s.startswith("Unnamed:") or not s or s.lower() in ("nan", "none", "null"):
        return ""
    return s


def _put_text(slide, text, x, y, w, h, *, font: str = BODY_FONT,
              font_size: float = 12, bold: bool = False,
              color: Optional[RGBColor] = None,
              align: str = "left", anchor: str = "top",
              line_spacing: float = 1.2,
              letter_spacing_em: float = 0.0) -> None:
    """Add a text box. align: left/center/right. anchor: top/middle/bottom."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    p.line_spacing = line_spacing

    run = p.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing_em:
        try:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing_em * font_size * 100)))
        except Exception:
            pass
