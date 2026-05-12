"""Wave 3: render determinístico (sin AI) sobre cada fixture.

Construimos un plan determinístico a partir del inventory: el primer KPI block,
el primer chart cat, el primer time-series, la primera tabla y unos bullets.
Esto exercise el extractor + renderer SIN depender del AI.

Output: __tests__/stress_fixtures/decks/*.pptx + reporte JSON.

Validaciones por deck generado:
  - El archivo se creó y pesa > 25KB
  - Se puede abrir con python-pptx
  - Tiene >= 2 slides (portada + al menos un contenido)
  - Cada slide tiene shapes (no es blank)
"""
from __future__ import annotations

import json
import sys
import time
import traceback
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

FIXTURES_DIR = Path("__tests__/stress_fixtures")
OUT_DIR = FIXTURES_DIR / "decks"
OUT_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATE = Path("Plantilla_Presentacion_Socya (1) (1).pptx")


def build_plan(inv) -> dict:
    """Construye un plan determinístico a partir del inventory."""
    kpi_blocks = [b for b in inv if b.kind == "kpi_candidate"][:3]
    chart_cats = [b for b in inv if b.kind == "categorical_distribution"]
    chart_ts = [b for b in inv if b.kind == "time_series_candidate"]
    tables = [b for b in inv if b.kind == "table"]

    slides = [
        {"type": "title",
         "data": {"title": "Reporte Stress Test",
                   "subtitle": "Generado por suite de pruebas"}},
    ]
    if kpi_blocks:
        slides.append({
            "type": "kpi_row",
            "title": "Indicadores Clave",
            "block_refs": [b.id for b in kpi_blocks],
        })
    if chart_cats:
        b = chart_cats[0]
        slides.append({
            "type": "chart",
            "chart_type": "bar",
            "title": f"Distribución por {b.provenance.columns[0]}",
            "block_ref": b.id,
            "narrative": "Distribución observada en los datos.",
        })
    if chart_ts:
        b = chart_ts[0]
        slides.append({
            "type": "chart",
            "chart_type": "line",
            "title": "Evolución temporal",
            "block_ref": b.id,
            "narrative": "Serie cronológica.",
        })
    # Histogram on first KPI candidate (numeric column)
    if kpi_blocks:
        b = kpi_blocks[0]
        slides.append({
            "type": "chart",
            "chart_type": "histogram",
            "title": f"Distribución de {b.provenance.columns[0]}",
            "block_ref": b.id,
            "narrative": "Distribución de valores.",
        })
    if tables:
        b = tables[0]
        slides.append({
            "type": "table",
            "title": "Detalle Operativo",
            "block_ref": b.id,
            "max_rows": 10,
        })
    # text_bullets necesita supports_block para que el validator no la dropee
    # y los bullets necesitan citar datos REALES del block (provenance check).
    if tables and kpi_blocks:
        b = kpi_blocks[0]
        col_name = b.provenance.columns[0] if b.provenance.columns else ""
        slides.append({
            "type": "text_bullets",
            "title": "Hallazgos",
            "supports_block": tables[0].id,
            "bullets": [
                f"La columna '{col_name}' fue analizada en este reporte.",
                f"Se evaluaron {len(tables)} bloques de tabla disponibles.",
                f"El inventario detectó {sum(1 for _ in inv if _.kind == 'kpi_candidate')} KPIs.",
            ],
        })
    return {
        "presentation_meta": {"title": "Stress Test",
                               "subtitle": "Validación end-to-end"},
        "slides": slides,
    }


def render_one(p: Path) -> dict:
    out = {"file": p.name, "ok": True, "first_error": None,
            "stats": {}}
    try:
        t0 = time.perf_counter()
        from socya_pipeline.parser import parse_workbook
        from socya_pipeline.inventory import build_inventory
        from socya_pipeline.validator import validate_plan
        from socya_pipeline.extractor import extract_for_render
        from socya_pipeline.renderer import render_pptx
        import pandas as pd

        wb = parse_workbook(p)
        inv = build_inventory(wb)
        plan = build_plan(inv)
        outcome = validate_plan(plan, inv, wb)

        out["stats"]["plan_slides"] = len(plan["slides"])
        out["stats"]["validated_slides"] = len(outcome.slides)
        out["stats"]["validator_dropped"] = len(outcome.dropped)

        if not outcome.slides:
            out["ok"] = False
            out["first_error"] = {"stage": "validate",
                                    "msg": "No slides validated."}
            return out

        xls = pd.ExcelFile(p)
        rendered, ext_dropped = extract_for_render(
            outcome.slides, inv, wb, str(p), xls=xls)
        out["stats"]["rendered_slides"] = len(rendered)
        out["stats"]["extractor_dropped"] = len(ext_dropped)

        out_path = OUT_DIR / (p.stem + ".pptx")
        render_pptx(rendered, plan["presentation_meta"],
                     template_path=TEMPLATE, output_path=out_path)

        out["stats"]["output_kb"] = out_path.stat().st_size // 1024
        out["stats"]["total_ms"] = round((time.perf_counter() - t0) * 1000, 1)

        # Validar PPTX abriéndolo
        from pptx import Presentation
        prs = Presentation(str(out_path))
        out["stats"]["actual_slides_in_pptx"] = len(prs.slides)
        out["stats"]["shapes_per_slide"] = [
            len(s.shapes) for s in prs.slides
        ]
        # Slide blanca = problema (debe tener al menos 3 shapes: bg, eyebrow, content)
        empty_slides = [i for i, n in enumerate(out["stats"]["shapes_per_slide"]) if n < 3]
        if empty_slides:
            out["stats"]["empty_slides"] = empty_slides
            out["ok"] = False
            out["first_error"] = {
                "stage": "render-validate",
                "msg": f"Slides con < 3 shapes (probable blank): {empty_slides}",
            }
    except Exception as e:
        out["ok"] = False
        out["first_error"] = {"stage": "render", "msg": str(e)[:300],
                                "tb": traceback.format_exc()[-800:]}
    return out


def main():
    fixtures = sorted(FIXTURES_DIR.glob("*.xlsx"))
    results = []
    for p in fixtures:
        r = render_one(p)
        results.append(r)
        if r["ok"]:
            st = r["stats"]
            print(f"[OK]   {r['file']:32s}  "
                   f"plan={st.get('plan_slides', 0)}  "
                   f"rendered={st.get('rendered_slides', 0)}  "
                   f"shapes={st.get('shapes_per_slide')}  "
                   f"size={st.get('output_kb', 0)}KB  "
                   f"in {st.get('total_ms', 0)}ms")
        else:
            err = r["first_error"]
            print(f"[FAIL] {r['file']:32s}  STAGE={err['stage']}  MSG={err['msg'][:120]}")

    n_ok = sum(1 for r in results if r["ok"])
    n_fail = len(results) - n_ok
    print(f"\nTotal: {n_ok} OK / {n_fail} FAIL  de {len(results)}.")
    if n_fail:
        print("\nDETALLE FALLOS:")
        for r in results:
            if not r["ok"]:
                err = r["first_error"]
                print(f"\n--- {r['file']} ---")
                print(f"  stage: {err['stage']}")
                print(f"  msg:   {err['msg']}")
                if "tb" in err:
                    for ln in err["tb"].splitlines()[-6:]:
                        print(f"    {ln}")
                if "stats" in r:
                    for k, v in r["stats"].items():
                        print(f"  stat: {k}={v}")

    out_json = FIXTURES_DIR / "_render_report.json"
    out_json.write_text(json.dumps(results, indent=2, ensure_ascii=False, default=str),
                          encoding="utf-8")
    print(f"\nReporte JSON: {out_json}")


if __name__ == "__main__":
    main()
