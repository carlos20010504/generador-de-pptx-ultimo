import json
import math
import os
import re
import sys
import tempfile
import unicodedata
from datetime import datetime
from datetime import UTC
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib import ticker as mticker
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from organizer import preparar_datos_para_slides, parse_user_request_context, DEFAULT_PRESENTATION_THEME


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


BASE_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = BASE_DIR / "Plantilla_Presentacion_Socya (1) (1).pptx"

DEFAULT_COLOR_BLUE = RGBColor(0, 51, 102)
DEFAULT_COLOR_ORANGE = RGBColor(255, 102, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(34, 34, 34)
COLOR_SOFT = RGBColor(102, 102, 102)
COLOR_LINE = RGBColor(220, 226, 232)
DEFAULT_COLOR_GREEN = RGBColor(67, 160, 71)
COLOR_YELLOW = RGBColor(237, 167, 0)
COLOR_RED = RGBColor(198, 40, 40)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_PANEL = RGBColor(248, 249, 252)
COLOR_BLUE = DEFAULT_COLOR_BLUE
COLOR_ORANGE = DEFAULT_COLOR_ORANGE
COLOR_GREEN = DEFAULT_COLOR_GREEN


def rgb_components(color):
    return tuple(int(component) for component in color)


def blend_rgb(base_color, target_color, factor):
    base = rgb_components(base_color)
    target = rgb_components(target_color)
    ratio = max(0.0, min(1.0, factor))
    return RGBColor(
        *(
            int(round(base[index] + (target[index] - base[index]) * ratio))
            for index in range(3)
        )
    )

SLIDE3_ZONE = {
    "x": Inches(5.95),
    "y": Inches(0.42),
    "w": Inches(6.65),
    "h": Inches(6.35),
}

SLIDE4 = {
    "title_x": Inches(1.55),
    "title_y": Inches(0.38),
    "title_w": Inches(10.85),
    "title_h": Inches(0.42),
    "content_x": Inches(1.55),
    "content_y": Inches(1.12),
    "content_w": Inches(10.95),
    "content_h": Inches(5.85),
    "footer_x": Inches(1.55),
    "footer_y": Inches(6.95),
    "footer_w": Inches(10.95),
    "footer_h": Inches(0.18),
}

ROWS_PER_TABLE_PAGE = 10
MAX_TABLE_COLS = 6
MAX_CONTENT_BLOCKS = 14
SLIDE3_MAX_BULLETS = 3
MAX_DASHBOARD_BARS = 4
MAX_BOARDROOM_ITEMS = 5
MAX_DERIVED_CHARTS = 3
MAX_CHART_BLOCKS = 3
MAX_KPI_DASHBOARDS = 5

REQUESTED_VISUAL_MODE = str(os.getenv("SOCYA_PRESENTATION_MODE", "boardroom")).strip().lower()
PRESENTATION_VISUAL_MODE = REQUESTED_VISUAL_MODE if REQUESTED_VISUAL_MODE in {"charts", "tables", "mixed", "boardroom"} else "mixed"
IS_BOARDROOM_MODE = PRESENTATION_VISUAL_MODE == "boardroom"

NUMBER_WORDS = {
    "un": 1,
    "una": 1,
    "uno": 1,
    "dos": 2,
    "tres": 3,
    "cuatro": 4,
    "cinco": 5,
}

DEFAULT_PROMPT_PREFERENCES = {
    "raw": "",
    "normalized": "",
    "chart_slides": 0,
    "table_slides": 0,
    "kpi_slides": 0,
    "conclusion_slides": 0,
    "prefer_short_tables": False,
    "prefer_bar_charts": False,
    "prefer_executive_tone": False,
    "prefer_dark_blue_green_palette": False,
    "explicit_structure": False,
    "max_boardroom_blocks": 4,
    "max_total_content_blocks": MAX_CONTENT_BLOCKS,
    "priority_tokens": [],
    "focus_terms": [],
}

PROMPT_STOPWORDS = {
    "de", "la", "el", "los", "las", "un", "una", "unos", "unas", "para", "con", "sin", "por",
    "del", "al", "y", "o", "u", "que", "se", "en", "sobre", "desde", "hacia", "como", "quiero",
    "muestre", "mostrar", "presentacion", "presentación", "powerpoint", "ppt", "slide", "slides",
    "datos", "dato", "excel", "archivo", "base", "solo", "real", "reales", "tipo", "tema",
    "colores", "tono", "narrativa", "grafica", "grafico", "graficas", "graficos", "tabla", "tablas",
    "kpi", "kpis", "muy", "mas", "más", "del", "las", "los", "una", "uno", "quiera", "necesito",
    "enfocada", "enfocado", "centrada", "centrado", "decision", "decisiones", "solicitado", "solicitada",
    "financieras", "financiera", "acciones", "control", "mayor", "mayores", "altos", "altas",
    "genera", "generar", "basado", "basada", "basados", "basadas", "prioritario", "prioritaria",
    "prioritarios", "prioritarias", "ejecutiva", "ejecutivo", "ejecutivos", "enfasis", "foco",
}
STRICT_VISUAL_CURATION_MODE = True


TEXT_ARTIFACT_REPLACEMENTS = {
    "├í": "á",
    "├⌐": "é",
    "├®": "í",
    "├│": "ó",
    "├║": "ú",
    "├▒": "ñ",
    "├ü": "Ü",
    "├ä": "Ä",
    "├ô": "Ó",
    "├Ü": "Ñ",
    "Ã¡": "á",
    "Ã©": "é",
    "Ã­": "í",
    "Ã³": "ó",
    "Ãº": "ú",
    "Ã±": "ñ",
    "Ã": "A",
    "â€™": "'",
    "â€œ": '"',
    "â€\x9d": '"',
    "â€“": "-",
    "â€”": "-",
    "â€¢": "-",
    "’": "'",
    "‘": "'",
    "“": '"',
    "”": '"',
    "¾": "ó",
    "ß": "á",
    "Ý": "í",
    "Ú": "É",
    "Ë": "Ó",
    "═": "Í",
    "³": "",
    "┐": "¿",
    "┬┐": "¿",
    "À": "-",
}

MOJIBAKE_HINT_CHARS = "ÃÂ├┤┐╔╗═¾ÝËÐ�"


def repair_text_artifacts(text):
    cleaned = str(text or "")
    cleaned = repair_mojibake_text(cleaned)
    for source, target in TEXT_ARTIFACT_REPLACEMENTS.items():
        cleaned = cleaned.replace(source, target)
    cleaned = unicodedata.normalize("NFKC", cleaned)
    cleaned = re.sub(r"[\u200b-\u200f\u202a-\u202e]", "", cleaned)
    cleaned = cleaned.replace("\xa0", " ").replace("_x000d_", " ")
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def mojibake_score(text):
    return sum(text.count(char) for char in MOJIBAKE_HINT_CHARS)


def attempt_redecode(text, source_encoding, target_encoding):
    try:
        return text.encode(source_encoding, errors="ignore").decode(target_encoding, errors="ignore")
    except Exception:
        return text


def repair_mojibake_text(text):
    current = str(text or "")
    if not current:
        return current
    best = current
    best_score = mojibake_score(current)
    candidates = [
        attempt_redecode(current, "latin1", "utf-8"),
        attempt_redecode(current, "cp1252", "utf-8"),
        attempt_redecode(current, "latin1", "cp1252"),
    ]
    for candidate in candidates:
        if not candidate:
            continue
        score = mojibake_score(candidate)
        if score < best_score and sum(char.isalnum() for char in candidate) >= max(2, sum(char.isalnum() for char in current) // 2):
            best = candidate
            best_score = score
    return best


def clean_text(value, max_len=None):
    if value is None:
        return ""
    text = str(value).replace("_x000d_", " ").replace("\r", " ").replace("\n", " ")
    text = repair_text_artifacts(text)
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = " ".join(text.split()).strip()
    if text.lower() in {"nan", "none", "null", "n/a", "nd", "s/d"}:
        return ""
    return text


def get_request_context_from_value(user_instructions):
    return parse_user_request_context(user_instructions)


def get_request_context(data):
    context = data.get("presentation_request")
    if isinstance(context, dict):
        return context
    return parse_user_request_context(data.get("user_instructions"))


def get_theme_context(data):
    context = get_request_context(data)
    theme = context.get("theme")
    if isinstance(theme, dict):
        return theme
    return dict(DEFAULT_PRESENTATION_THEME)


def get_presentation_design_plan(data):
    plan = data.get("presentation_design_ia")
    if not isinstance(plan, dict):
        plan = data.get("prompt_fallback_design_plan")
    return plan if isinstance(plan, dict) else {}


def get_ai_curation_status(data):
    status = data.get("ai_curation_status")
    if isinstance(status, dict):
        return status
    return {
        "strict_visual_mode": STRICT_VISUAL_CURATION_MODE,
        "visual_curation_ready": False,
        "visual_plan_received": False,
        "selected_chart_ids": [],
        "selected_table_ids": [],
        "reason": "ai_status_missing",
    }


def is_visual_curation_ready(data):
    status = get_ai_curation_status(data)
    if not status.get("strict_visual_mode", STRICT_VISUAL_CURATION_MODE):
        return bool(data.get("visual_plan_ia"))
    return bool(status.get("visual_curation_ready")) and bool(data.get("visual_plan_ia"))


def count_request_near_keywords(normalized_prompt, keywords):
    if not normalized_prompt:
        return 0
    keyword_pattern = "|".join(re.escape(keyword) for keyword in keywords)
    pattern = re.compile(
        rf"\b(\d+|un|una|uno|dos|tres|cuatro|cinco)\b(?:\s+(?:diapositiva|diapositivas|slide|slides|grafica|graficas|grafico|graficos|tabla|tablas|kpi|kpis|resumen|cierres?))?[^.:\n]{{0,45}}\b(?:{keyword_pattern})\b"
    )
    matches = pattern.findall(normalized_prompt)
    if not matches:
        return 0
    raw_count = matches[-1].lower()
    return NUMBER_WORDS.get(raw_count, int(raw_count) if raw_count.isdigit() else 0)


def parse_prompt_preferences(user_instructions):
    request_context = get_request_context_from_value(user_instructions)
    raw_prompt = clean_text(request_context.get("prompt"))
    if not raw_prompt:
        return dict(DEFAULT_PROMPT_PREFERENCES)

    normalized = normalize_header_key(raw_prompt)
    chart_slides = count_request_near_keywords(normalized, ("grafica", "graficas", "grafico", "graficos", "barras", "comparativa", "comparativas"))
    table_slides = count_request_near_keywords(normalized, ("tabla", "tablas", "hallazgos"))
    kpi_slides = count_request_near_keywords(normalized, ("kpi", "kpis", "indicador", "indicadores"))
    conclusion_slides = count_request_near_keywords(normalized, ("conclusion", "conclusiones", "recomendacion", "recomendaciones", "cierre"))

    priority_tokens = []
    topic_aliases = {
        "hallazgos": ("hallazgo", "hallazgos"),
        "costos": ("costo", "costos", "monto", "montos", "valor", "valores"),
        "estado": ("estado", "estados"),
        "coso": ("coso", "control interno", "controles"),
        "riesgos": ("riesgo", "riesgos", "alerta", "alertas"),
    }
    for token, aliases in topic_aliases.items():
        if any(alias in normalized for alias in aliases):
            priority_tokens.append(token)
    focus_terms = extract_prompt_focus_terms(raw_prompt)

    explicit_structure = any(value > 0 for value in (chart_slides, table_slides, kpi_slides, conclusion_slides))
    max_total_content_blocks = MAX_CONTENT_BLOCKS
    max_boardroom_blocks = 4
    if explicit_structure:
        requested_total = sum(value for value in (chart_slides, table_slides, kpi_slides, conclusion_slides) if value > 0)
        max_total_content_blocks = max(6, min(9, requested_total + 3))
        max_boardroom_blocks = 2

    return {
        "raw": raw_prompt,
        "normalized": normalized,
        "request_context": request_context,
        "chart_slides": chart_slides,
        "table_slides": table_slides,
        "kpi_slides": kpi_slides,
        "conclusion_slides": conclusion_slides,
        "prefer_short_tables": any(token in normalized for token in ("tabla corta", "tablas cortas", "evita tablas extensas", "resume solo", "resumir solo")),
        "prefer_bar_charts": any(token in normalized for token in ("barras", "bar chart", "barra comparativa", "graficas comparativas", "comparativa")),
        "prefer_executive_tone": any(token in normalized for token in ("ejecutiva", "ejecutivo", "comite", "comité", "boardroom", "directivo", "sobria")),
        "prefer_dark_blue_green_palette": any(
            token in normalized
            for token in ("azul oscuro y verde", "paleta azul", "azul oscuro", "verde")
        ),
        "explicit_structure": explicit_structure,
        "max_boardroom_blocks": max_boardroom_blocks,
        "max_total_content_blocks": max_total_content_blocks,
        "priority_tokens": priority_tokens,
        "focus_terms": focus_terms,
    }


def apply_prompt_palette(preferences):
    global COLOR_BLUE, COLOR_ORANGE, COLOR_GREEN
    COLOR_BLUE = DEFAULT_COLOR_BLUE
    COLOR_ORANGE = DEFAULT_COLOR_ORANGE
    COLOR_GREEN = DEFAULT_COLOR_GREEN
    theme = (preferences.get("request_context") or {}).get("theme") or {}
    primary_hex = clean_text(theme.get("primary_hex"))
    accent_hex = clean_text(theme.get("accent_hex"))
    if re.fullmatch(r"#[0-9A-Fa-f]{6}", primary_hex or ""):
        COLOR_BLUE = RGBColor.from_string(primary_hex[1:])
    if re.fullmatch(r"#[0-9A-Fa-f]{6}", accent_hex or ""):
        COLOR_ORANGE = RGBColor.from_string(accent_hex[1:])
    if preferences.get("prefer_dark_blue_green_palette"):
        COLOR_BLUE = RGBColor(15, 23, 42)
        COLOR_ORANGE = RGBColor(14, 165, 140)
        COLOR_GREEN = RGBColor(34, 139, 94)


def is_noise_text(value):
    text = clean_text(value)
    if not text:
        return True
    normalized = text.casefold().strip(" _-./")
    if not normalized:
        return True
    alnum = "".join(char for char in normalized if char.isalnum())
    if len(alnum) <= 1:
        return True
    return normalized in {"ù", "û", "n/d", "nd", "na", "s/d", "none", "null", "nan", "x", "xx", "si/no"}


def basename_label(excel_path):
    return clean_text(Path(excel_path).stem.replace("_", " ").replace("-", " "), 80)


def extract_prompt_focus_terms(raw_prompt, limit=8):
    prompt = clean_text(raw_prompt, 240)
    if not prompt:
        return []
    normalized = normalize_header_key(prompt)
    candidates = re.findall(r"[a-z0-9]{4,}", normalized)
    terms = []
    seen = set()
    for token in candidates:
        if token in PROMPT_STOPWORDS:
            continue
        if token.isdigit():
            continue
        if token in seen:
            continue
        seen.add(token)
        terms.append(token)
        if len(terms) >= limit:
            break
    return terms


def build_prompt_focus_phrase(data, excel_path, max_len=72):
    request_context = get_request_context(data)
    raw_prompt = clean_text(request_context.get("prompt"), 220)
    topic = compact_line(infer_excel_topic(data, excel_path), 42)
    preferences = parse_prompt_preferences(raw_prompt)
    tokens = preferences.get("priority_tokens") or []
    focus_terms = preferences.get("focus_terms") or []

    base = ""
    if ("hallazgos" in tokens or "riesgos" in tokens) and ("costos" in tokens or "estado" in tokens):
        base = "Hallazgos y concentraciones clave"
    elif "hallazgos" in tokens or "riesgos" in tokens:
        base = "Hallazgos y decisiones prioritarias"
    elif "costos" in tokens or "estado" in tokens:
        base = "Concentraciones y decisiones prioritarias"
    elif len(focus_terms) >= 2:
        base = compact_line(f"{focus_terms[0]} y {focus_terms[1]}", 54)
    elif focus_terms:
        base = compact_line(f"Decision ejecutiva sobre {focus_terms[0]}", 54)
    else:
        base = compact_line(raw_prompt, 54)

    if base:
        normalized_base = normalize_header_key(base)
        normalized_topic = normalize_header_key(topic)
        if topic and normalized_topic and normalized_topic not in normalized_base and len(base) + len(topic) + 3 <= max_len:
            return compact_line(f"{base} · {topic}", max_len)
        return compact_line(base, max_len)
    return topic or basename_label(excel_path) or "datos reales del Excel"


def safe_list(values):
    return [item for item in (values or []) if item]


def is_generic_header(value):
    text = clean_text(value)
    if not text:
        return True
    normalized = text.casefold().strip()
    alnum = "".join(char for char in normalized if char.isalnum())
    if len(alnum) <= 2:
        return True
    return bool(re.match(r"^(col(?:umna)?[_ -]?\d+|unnamed:? ?\d*|field[_ -]?\d+)$", normalized))


def meaningful_cell_count(row):
    return sum(1 for cell in row or [] if not is_noise_text(cell))


def filter_informative_rows(rows, header_count):
    if not rows:
        return []
    min_meaningful = 1 if header_count <= 2 else 2
    filtered = []
    fallback = []
    for row in rows:
        current = list(row[:header_count]) if row else []
        if len(current) < header_count:
            current.extend([""] * (header_count - len(current)))
        score = meaningful_cell_count(current)
        if score >= min_meaningful:
            filtered.append(current)
        elif score >= 1:
            fallback.append(current)
    return filtered or fallback


def trim_sparse_columns(headers, rows):
    if not headers or not rows:
        return headers, rows

    row_count = len(rows)
    keep_indexes = []
    for index, header in enumerate(headers):
        values = [row[index] if index < len(row) else "" for row in rows]
        meaningful = sum(1 for value in values if not is_noise_text(value))
        threshold = max(1, math.ceil(row_count * (0.08 if index == 0 else 0.12)))
        if is_generic_header(header):
            threshold = max(threshold, 3 if row_count >= 8 else 2)
        if meaningful >= threshold or not is_generic_header(header):
            keep_indexes.append(index)

    if len(keep_indexes) < 2:
        return headers, rows

    trimmed_headers = [headers[index] for index in keep_indexes]
    trimmed_rows = [
        [row[index] if index < len(row) else "" for index in keep_indexes]
        for row in rows
    ]
    return trimmed_headers, trimmed_rows


def densify_table_matrix(headers, rows, min_rows=2):
    if not headers or not rows:
        return [], []

    normalized_rows = []
    for row in rows:
        current = [clean_text(row[index]) if index < len(row) else "" for index in range(len(headers))]
        normalized_rows.append(current)

    keep_indexes = []
    for index, _header in enumerate(headers):
        column_values = [row[index] for row in normalized_rows]
        if column_values and all(not is_noise_text(value) for value in column_values):
            keep_indexes.append(index)

    if len(keep_indexes) < 2:
        return [], []

    dense_headers = [headers[index] for index in keep_indexes]
    dense_rows = []
    for row in normalized_rows:
        candidate = [row[index] for index in keep_indexes]
        if all(not is_noise_text(value) for value in candidate):
            dense_rows.append(candidate)

    if len(dense_rows) < min_rows:
        return [], []
    return dense_headers, dense_rows


def looks_like_sentence(value):
    text = clean_text(value)
    if not text:
        return False
    words = [part for part in text.split() if part]
    if len(words) >= 8:
        return True
    return len(text) >= 55


def count_valid_table_rows(rows):
    valid = 0
    for row in rows or []:
        current = row.values() if isinstance(row, dict) else row
        if any(clean_text(cell) for cell in current):
            valid += 1
    return valid


def extract_table_payload(table):
    if not table:
        return [], []
    headers = [clean_text(header, 72) for header in safe_list(table.get("encabezados") or table.get("columnas") or [])]
    rows = table.get("filas_limpias") or table.get("filas") or []
    normalized_rows = []
    for row in rows:
        if isinstance(row, dict):
            normalized = [clean_text(row.get(header, "")) for header in headers]
        else:
            normalized = [clean_text(cell) for cell in list(row)[: len(headers)]]
            if len(normalized) < len(headers):
                normalized.extend([""] * (len(headers) - len(normalized)))
        if any(normalized):
            normalized_rows.append(normalized)

    normalized_rows = filter_informative_rows(normalized_rows, len(headers))
    headers, normalized_rows = trim_sparse_columns(headers, normalized_rows)
    return headers, normalized_rows


def is_valid_chart(chart):
    labels = [clean_text(item, 36) for item in safe_list(chart.get("labels"))]
    values = []
    for value in safe_list(chart.get("valores")):
        values.append(try_number(value) or 0)
    pairs = [(label, value) for label, value in zip(labels, values) if label and value > 0]
    return len(pairs) >= 2


def is_time_like_label(value):
    normalized = normalize_header_key(value)
    if not normalized:
        return False
    month_tokens = {
        "ene", "enero", "feb", "febrero", "mar", "marzo", "abr", "abril", "may", "mayo", "jun", "junio",
        "jul", "julio", "ago", "agosto", "sep", "sept", "septiembre", "oct", "octubre", "nov", "noviembre",
        "dic", "diciembre", "month", "mes", "semana", "trimestre", "periodo",
    }
    if any(token in normalized for token in month_tokens):
        return True
    return bool(re.search(r"\b20\d{2}\b", normalized) or re.search(r"\b\d{1,2}[/-]\d{1,2}([/-]\d{2,4})?\b", normalized))


def resolve_derived_chart_type(table_name, label_header, labels, grouped_length):
    normalized = " ".join(filter(None, [clean_text(table_name), clean_text(label_header)])).lower()
    if "estado" in normalized:
        return "doughnut"
    if is_time_like_label(label_header) or any(is_time_like_label(label) for label in labels[:4]):
        return "line"
    if any(token in normalized for token in ("tipo", "categoria", "categoría", "mes", "ciudad", "destino")) and grouped_length <= 6:
        return "doughnut"
    if grouped_length <= 4:
        return "pie"
    return "bar"


def score_chart_dimension(header, rows, col_index):
    values = [clean_text(row[col_index]) for row in rows if col_index < len(row)]
    non_empty = [value for value in values if not is_noise_text(value)]
    if len(non_empty) < 3:
        return -1

    unique_values = []
    seen = set()
    for value in non_empty:
        key = value.casefold()
        if key not in seen:
            seen.add(key)
            unique_values.append(value)
    unique_count = len(unique_values)
    unique_ratio = unique_count / max(1, len(non_empty))
    if unique_count < 2 or unique_count > 18:
        return -1

    avg_len = sum(len(value) for value in non_empty) / max(1, len(non_empty))
    if avg_len > 30 or sum(1 for value in non_empty if looks_like_sentence(value)) / max(1, len(non_empty)) > 0.18:
        return -1

    normalized_header = normalize_header_key(header)
    score = 0
    if is_time_like_label(header) or any(is_time_like_label(value) for value in unique_values[:4]):
        score += 7
    if "fecha" in normalized_header and unique_count > 8:
        return -1
    if is_person_like_header(header):
        return -1
    if is_executive_dimension_header(header):
        score += 5
    if is_identifier_header(header):
        score -= 8
    if unique_ratio >= 0.72:
        score -= 5
    elif unique_ratio <= 0.35:
        score += 3
    if 2 <= unique_count <= 6:
        score += 4
    elif unique_count <= 10:
        score += 2
    score += min(len(non_empty), 12) / 4
    return score


def unique_texts(items, limit=None, min_len=12):
    seen = set()
    result = []
    for item in items or []:
        text = clean_text(item)
        if not text or len(text) < min_len:
            continue
        key = text.casefold()
        if key in seen:
            continue
        seen.add(key)
        result.append(text)
        if limit and len(result) >= limit:
            break
    return result


def compact_line(text, max_len=120):
    cleaned = clean_text(text)
    if not cleaned:
        return ""
    cleaned = cleaned.replace("•", "").strip()
    parts = [part.strip() for part in cleaned.replace(";", ". ").split(". ") if part.strip()]
    if parts:
        cleaned = parts[0]
    return clean_text(cleaned)


def build_ranked_metric_summary(labels, values, kind="number"):
    pairs = []
    for label, value in zip(safe_list(labels), safe_list(values)):
        label_text = clean_text(label, 48)
        numeric_value = try_number(value)
        if not label_text or numeric_value is None:
            continue
        pairs.append((label_text, float(numeric_value)))
    if not pairs:
        return None
    pairs.sort(key=lambda item: item[1], reverse=True)
    total = sum(max(value, 0) for _, value in pairs)
    top_label, top_value = pairs[0]
    share = (top_value / total) if total > 0 else None
    return {
        "items": pairs,
        "total": total,
        "top_label": top_label,
        "top_value": top_value,
        "share": share,
        "top_text": format_metric_full(top_value, kind),
        "share_text": format_value(share, "percent", compact=False) if share is not None else None,
    }


def looks_like_code_value(value):
    text = clean_text(value)
    compact = re.sub(r"[\s\-_/]", "", text)
    if len(compact) < 4 or len(compact) > 14:
        return False
    if " " in text:
        return False
    return bool(re.fullmatch(r"[A-Z]*\d+[A-Z0-9]*", compact))


def present_story_label(value, fallback_label):
    text = clean_text(value, 44)
    if not text:
        return fallback_label
    if looks_like_code_value(text):
        return fallback_label
    if text == text.upper() and not looks_like_code_value(text):
        text = " ".join(part.capitalize() for part in text.split())
    return text


def is_management_noise_line(value):
    normalized = normalize_header_key(value)
    if not normalized:
        return True
    blocked_tokens = (
        "hoja", "hojas", "tabla", "tablas", "grafica", "graficas",
        "auditoria y control", "archivo corresponde", "hoja principal",
        "filas", "columnas", "funcionales", "coso", "procedimiento",
        "distribucion", "oportunidades",
    )
    return any(token in normalized for token in blocked_tokens)


def try_number(value):
    if value is None:
        return None
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = clean_text(value)
    if not text:
        return None
    lowered = (
        text.replace("\xa0", " ")
        .replace("−", "-")
        .replace("–", "-")
        .replace("—", "-")
        .strip()
    )
    negative = False
    if lowered.startswith("(") and lowered.endswith(")"):
        negative = True
        lowered = lowered[1:-1].strip()
    if lowered.startswith("-"):
        negative = True
        lowered = lowered[1:].strip()

    multiplier = 1.0
    lowered_case = lowered.lower()
    if re.search(r"\b(mil\s*mm|mil\s*millones|bn|billones?)\b", lowered_case):
        multiplier = 1_000_000_000.0
    elif re.search(r"\b(mm|millones?)\b", lowered_case) or re.search(r"\d\s*m\b", lowered_case):
        multiplier = 1_000_000.0
    elif re.search(r"\b(k|mil)\b", lowered_case):
        multiplier = 1_000.0

    lowered = re.sub(r"(?i)\b(cop|usd|eur|pesos?|dolares?|moneda|aprox|aproximado|estimado)\b", "", lowered)
    normalized = re.sub(r"[^0-9,.\-]", "", lowered).strip()
    if not normalized or not re.search(r"\d", normalized):
        return None
    normalized = re.sub(r"(?<!^)-", "", normalized)

    if "," in normalized and "." in normalized:
        if normalized.rfind(",") > normalized.rfind("."):
            normalized = normalized.replace(".", "").replace(",", ".")
        else:
            normalized = normalized.replace(",", "")
    elif "," in normalized:
        parts = normalized.split(",")
        if len(parts) == 2:
            left, right = parts
            if len(right) <= 2:
                normalized = f"{left}.{right}"
            elif len(right) == 3:
                normalized = f"{left}{right}"
            else:
                normalized = "".join(parts)
        else:
            normalized = "".join(parts)
    elif "." in normalized:
        parts = normalized.split(".")
        if len(parts) == 2:
            left, right = parts
            if len(right) <= 2:
                normalized = f"{left}.{right}"
            elif len(right) == 3:
                normalized = f"{left}{right}"
        elif len(parts) > 2:
            if len(parts[-1]) <= 2:
                normalized = "".join(parts[:-1]) + "." + parts[-1]
            else:
                normalized = "".join(parts)
    if negative and not normalized.startswith("-"):
        normalized = f"-{normalized}"
    try:
        return float(normalized) * multiplier
    except Exception:
        return None


def format_spanish_number(value, decimals=0):
    formatted = f"{value:,.{decimals}f}"
    return formatted.replace(",", "_").replace(".", ",").replace("_", ".")


def infer_metric_kind(label=None, values=None):
    normalized = normalize_header_key(label)
    if any(token in normalized for token in ("fecha", "date", "mes", "month", "periodo", "period", "ano", "anio", "year", "semana", "week", "dia", "day")):
        return "temporal"
    if "%" in clean_text(label) or any(
        token in normalized
        for token in ("porcentaje", "ratio", "participacion", "share", "margen", "cumplimiento", "avance")
    ):
        return "percent"
    if any(
        token in normalized
        for token in ("valor", "monto", "total", "costo", "gasto", "ingreso", "venta", "precio", "presupuesto", "saldo", "importe", "cop", "peso", "pesos", "moneda", "tarifa")
    ):
        return "currency"
    numeric_values = [abs(float(value)) for value in safe_list(values) if value is not None]
    if numeric_values and max(numeric_values) <= 1.2 and any(not float(value).is_integer() for value in numeric_values):
        return "percent"
    return "number"


def get_compact_scale_parts(value):
    absolute_value = abs(float(value))
    if absolute_value >= 1_000_000_000:
        return float(value) / 1_000_000_000, " mil MM"
    if absolute_value >= 1_000_000:
        return float(value) / 1_000_000, " M"
    if absolute_value >= 1_000:
        return float(value) / 1_000, " mil"
    return float(value), ""


def get_compact_decimals(value, compact):
    if not compact:
        return 0 if float(value).is_integer() else 2
    absolute_value = abs(float(value))
    if absolute_value >= 100:
        return 0
    if absolute_value >= 10:
        return 1
    return 1


def format_value(value, kind="number", compact=False):
    if value is None:
        return "N/D"
    numeric_value = float(value)
    absolute_value = abs(numeric_value)
    if kind == "percent":
        percent_value = numeric_value * 100 if absolute_value <= 1.2 else numeric_value
        decimals = 0 if compact and abs(percent_value) >= 10 else 1
        return f"{format_spanish_number(percent_value, decimals)}%"
    scaled, suffix = get_compact_scale_parts(numeric_value) if compact else (numeric_value, "")
    if kind == "currency":
        decimals = 0 if not compact and absolute_value >= 1000 else get_compact_decimals(scaled, compact)
        if not compact and absolute_value < 100 and not numeric_value.is_integer():
            decimals = 2
        return f"COP {format_spanish_number(scaled, decimals)}{suffix}"
    if compact:
        decimals = get_compact_decimals(scaled, compact)
        return f"{format_spanish_number(scaled, decimals)}{suffix}"
    if absolute_value >= 1000:
        return format_spanish_number(numeric_value, 0)
    if numeric_value.is_integer():
        return str(int(numeric_value))
    decimals = 2 if absolute_value < 1 else 1
    return format_spanish_number(numeric_value, decimals)


def format_metric(value, kind="number"):
    return format_value(value, kind=kind, compact=True)


def format_metric_full(value, kind="number"):
    return format_value(value, kind=kind, compact=False)


def format_axis_value(value, kind="number"):
    if value is None:
        return ""
    numeric_value = float(value)
    if kind == "percent":
        percent_value = numeric_value * 100 if abs(numeric_value) <= 1.2 else numeric_value
        decimals = 0 if abs(percent_value) >= 10 else 1
        return f"{format_spanish_number(percent_value, decimals)}%"
    if kind == "currency":
        scaled, suffix = get_compact_scale_parts(numeric_value)
        decimals = 0 if abs(scaled) >= 10 else 1
        return f"{format_spanish_number(scaled, decimals)}{suffix}".strip()
    return format_value(numeric_value, kind=kind, compact=True)


def format_chart_callout_value(value, kind="number", prefer_compact=False):
    full_text = format_value(value, kind=kind, compact=False)
    if prefer_compact or len(full_text) > 16:
        return format_value(value, kind=kind, compact=True)
    return full_text


def build_source_basis_text(dimension_label=None, metric_label=None, aggregation=None):
    dimension = clean_text(dimension_label, 28) or "categoria"
    metric = clean_text(metric_label, 32) or "registros"
    if metric == "%":
        metric = "porcentaje"
    elif metric.upper() == metric and len(metric) <= 12:
        metric = metric.title()
    if len(dimension) > 24 and not is_executive_dimension_header(dimension):
        dimension = "categoria principal"
    if aggregation == "suma":
        return f"Base real del Excel: suma de {metric} por {dimension}"
    if aggregation == "tendencia":
        return f"Base real del Excel: tendencia de {metric} por {dimension}"
    return f"Base real del Excel: conteo de registros por {dimension}"


def detect_status_color(status):
    normalized = clean_text(status).lower()
    if normalized in {"rojo", "alta", "critico", "critica"}:
        return COLOR_RED
    if normalized in {"amarillo", "media", "moderado", "moderada"}:
        return COLOR_YELLOW
    return COLOR_GREEN


def status_label_from_value(value, warning_threshold, critical_threshold):
    if value is None:
        return "amarillo"
    if value >= critical_threshold:
        return "rojo"
    if value >= warning_threshold:
        return "amarillo"
    return "verde"


def is_identifier_header(header):
    normalized = normalize_header_key(header)
    return any(token in normalized for token in ("id", "codigo", "consecutivo", "numero", "nro", "radicado", "folio"))


def is_person_like_header(header):
    normalized = normalize_header_key(header)
    return any(token in normalized for token in ("solicitante", "nombre", "responsable", "empleado", "cliente", "proveedor", "usuario", "persona", "colaborador"))


def is_executive_dimension_header(header):
    normalized = normalize_header_key(header)
    return any(token in normalized for token in ("estado", "mes", "ciudad", "destino", "centro", "tipo", "categoria", "control", "riesgo", "hallazgo", "componente", "proceso", "area"))


def is_executive_metric_header(header):
    normalized = normalize_header_key(header)
    return any(token in normalized for token in ("valor", "monto", "total", "costo", "gasto", "ingreso", "venta", "dias", "porcentaje", "%", "cantidad", "cobertura", "avance"))


def derive_table_basis(headers, rows):
    label_index = choose_table_label_column(headers, rows)
    numeric_index = choose_table_numeric_column(headers, rows)
    if numeric_index is not None and is_identifier_header(headers[numeric_index]):
        numeric_index = None
    if label_index is not None and (is_identifier_header(headers[label_index]) or is_person_like_header(headers[label_index])):
        executive_indexes = [index for index, header in enumerate(headers) if is_executive_dimension_header(header)]
        if executive_indexes:
            label_index = executive_indexes[0]
    return {
        "label_index": label_index,
        "numeric_index": numeric_index,
        "dimension_label": clean_text(headers[label_index], 34) if label_index is not None and label_index < len(headers) else "",
        "metric_label": clean_text(headers[numeric_index], 36) if numeric_index is not None and numeric_index < len(headers) else "Registros",
        "aggregation": "suma" if numeric_index is not None else "conteo",
    }


def build_table_meta_summary(headers, rows, basis=None):
    basis = basis or derive_table_basis(headers, rows)
    parts = [f"{len(rows)} filas", f"{len(headers)} columnas"]
    numeric_index = basis.get("numeric_index")
    if numeric_index is not None:
        values = [try_number(row[numeric_index]) for row in rows if numeric_index < len(row)]
        numeric_values = [value for value in values if value is not None]
        if numeric_values:
            metric_kind = infer_metric_kind(basis.get("metric_label"), numeric_values)
            parts.append(f"Total {format_value(sum(numeric_values), metric_kind, compact=True)}")
    return " | ".join(parts)


def clamp_score(value, lower=0.0, upper=0.99):
    return max(lower, min(upper, value))


def resolve_source_sheet_name(data, explicit_source=None, table_name=None):
    explicit = clean_text(explicit_source, 42)
    if explicit:
        return explicit
    primary_sheet = clean_text(
        (data.get("metadatos", {}) or {}).get("hoja_principal")
        or (data.get("resumen_generico", {}) or {}).get("hoja_principal"),
        42,
    )
    normalized_table = clean_text(table_name).lower()
    if normalized_table and normalized_table not in {"tabla principal", "evaluacion coso", "distribucion por mes"}:
        return clean_text(table_name, 42)
    if primary_sheet:
        return primary_sheet
    return clean_text((data.get("metadatos", {}) or {}).get("archivo"), 42) or "Hoja principal"


def confidence_label(score):
    if score >= 0.8:
        return "Alta"
    if score >= 0.62:
        return "Media"
    return "Baja"


def aggregation_label(aggregation):
    mapping = {
        "suma": "Suma",
        "conteo": "Conteo",
        "tendencia": "Tendencia",
    }
    return mapping.get(clean_text(aggregation).lower(), clean_text(aggregation).title() or "N/D")


def build_traceability(source_sheet, aggregation, confidence_score, row_count=None, col_count=None, numeric_coverage=None):
    return {
        "source_sheet": clean_text(source_sheet, 42) or "Hoja principal",
        "aggregation": clean_text(aggregation) or "conteo",
        "aggregation_label": aggregation_label(aggregation),
        "confidence_score": clamp_score(confidence_score),
        "confidence_label": confidence_label(confidence_score),
        "row_count": int(row_count or 0),
        "col_count": int(col_count or 0),
        "numeric_coverage": None if numeric_coverage is None else max(0.0, min(1.0, float(numeric_coverage))),
    }


def build_traceability_caption(traceability):
    if not traceability:
        return ""
    score_pct = int(round((traceability.get("confidence_score") or 0) * 100))
    return (
        f"Fuente: {traceability.get('source_sheet')} | "
        f"Regla: {traceability.get('aggregation_label')} | "
        f"Confianza: {traceability.get('confidence_label')} {score_pct}%"
    )


def evaluate_chart_truth(chart, source_sheet=None):
    labels = [clean_text(item) for item in safe_list(chart.get("labels"))]
    values = [try_number(item) for item in safe_list(chart.get("valores"))]
    pairs = [(label, value) for label, value in zip(labels, values) if label and value is not None]
    if len(pairs) < 2:
        return {
            "score": 0.18,
            "blocked": True,
            "reason": "menos de dos pares validos",
            "numeric_coverage": 0.0,
        }
    numeric_values = [value for _, value in pairs]
    valid_pair_ratio = len(pairs) / max(1, len(labels))
    non_zero_ratio = sum(1 for value in numeric_values if abs(value) > 1e-9) / max(1, len(numeric_values))
    positive_ratio = sum(1 for value in numeric_values if value > 0) / max(1, len(numeric_values))
    unique_ratio = len({round(value, 6) for value in numeric_values}) / max(1, len(numeric_values))
    score = 0.24
    score += valid_pair_ratio * 0.18
    score += min(len(pairs), 6) / 6 * 0.16
    score += non_zero_ratio * 0.12
    score += positive_ratio * 0.12
    score += min(unique_ratio, 1.0) * 0.08
    score += 0.06 if clean_text(source_sheet) else 0.0
    score += 0.06 if clean_text(chart.get("aggregation")) in {"suma", "conteo", "tendencia"} else 0.0
    chart_type = clean_text(chart.get("tipo")).lower()
    spread = max(numeric_values) - min(numeric_values)
    if chart_type == "line" and spread <= 1e-9:
        score -= 0.22
    if len(pairs) < 3:
        score -= 0.08
    if max(abs(value) for value in numeric_values) <= 1e-9:
        score -= 0.28
    score = clamp_score(score)
    blocked = score < 0.58 or (chart_type == "line" and spread <= 1e-9)
    reason = "confianza insuficiente" if score < 0.58 else ("serie plana sin tendencia real" if chart_type == "line" and spread <= 1e-9 else "")
    return {
        "score": score,
        "blocked": blocked,
        "reason": reason,
        "numeric_coverage": valid_pair_ratio,
    }


def evaluate_table_truth(headers, rows, basis=None, source_sheet=None):
    basis = basis or derive_table_basis(headers, rows)
    profile = build_table_signal_profile(headers, rows)
    row_count = len(rows)
    informative_ratio = profile["informative_rows"] / max(1, row_count)
    score = 0.2
    score += min(row_count, 12) / 12 * 0.12
    score += informative_ratio * 0.16
    score += max(0.0, 1 - min(profile["noise_cell_ratio"], 1)) * 0.14
    score += max(0.0, 1 - min(profile["weak_row_ratio"], 1)) * 0.12
    score += max(0.0, 1 - min(profile["generic_header_ratio"], 1)) * 0.1
    score += 0.06 if clean_text(source_sheet) else 0.0

    numeric_index = basis.get("numeric_index")
    label_index = basis.get("label_index")
    numeric_coverage = None
    if numeric_index is not None:
        numeric_values = [try_number(row[numeric_index]) for row in rows if numeric_index < len(row)]
        numeric_coverage = sum(1 for value in numeric_values if value is not None) / max(1, row_count)
        score += numeric_coverage * 0.1
        if is_identifier_header(headers[numeric_index]):
            score -= 0.18
    else:
        score += 0.04

    if label_index is None:
        score -= 0.14
    elif is_identifier_header(headers[label_index]):
        score -= 0.16
    elif is_person_like_header(headers[label_index]):
        score -= 0.08
    else:
        score += 0.06

    if profile["table_kind"] == "narrative" and numeric_index is None:
        score -= 0.06
    if row_count < 3:
        score -= 0.1

    score = clamp_score(score)
    return {
        "score": score,
        "numeric_coverage": numeric_coverage,
        "row_count": row_count,
        "col_count": len(headers),
        "profile": profile,
        "basis": basis,
        "blocked_dashboard": score < 0.62,
        "blocked_detail": score < 0.56,
        "blocked_summary": score < 0.4,
    }


def choose_table_label_column(headers, rows):
    best_index = None
    best_score = -1
    for index, header in enumerate(headers):
        values = [clean_text(row[index]) for row in rows if index < len(row)]
        non_empty = [value for value in values if not is_noise_text(value)]
        if not non_empty:
            continue
        numeric_ratio = sum(1 for value in non_empty if try_number(value) is not None) / max(1, len(non_empty))
        if numeric_ratio > 0.4:
            continue
        date_like_ratio = sum(1 for value in non_empty if any(token in value for token in ("-", "/", ":"))) / max(1, len(non_empty))
        if date_like_ratio > 0.7:
            continue
        unique_ratio = len({value.casefold() for value in non_empty}) / max(1, len(non_empty))
        unique_count = len({value.casefold() for value in non_empty})
        score = unique_ratio + min(len(non_empty), 12) / 20
        if is_executive_dimension_header(header):
            score += 4
        if is_time_like_label(header):
            score += 3
        if is_person_like_header(header):
            score -= 5
        if is_identifier_header(header):
            score -= 8
        if unique_ratio >= 0.85 or unique_count > 18:
            score -= 4
        if score > best_score:
            best_score = score
            best_index = index
    return best_index if best_index is not None else 0


def choose_table_numeric_column(headers, rows):
    best_index = None
    best_score = -1
    for index, header in enumerate(headers):
        normalized_header = normalize_header_key(header)
        if any(token in normalized_header for token in ("fecha", "date", "mes", "month", "periodo", "period", "ano", "anio", "year", "semana", "week", "dia", "day")):
            continue
        values = [try_number(row[index]) for row in rows if index < len(row)]
        numeric_values = [value for value in values if value is not None]
        if len(numeric_values) < 2:
            continue
        coverage = len(numeric_values) / max(1, len(rows))
        if coverage < 0.45:
            continue
        unique_count = len({value for value in numeric_values})
        unique_ratio = unique_count / max(1, len(numeric_values))
        spread = max(numeric_values) - min(numeric_values)
        mean_abs = max(abs(sum(numeric_values) / len(numeric_values)), 1)
        non_zero_ratio = sum(1 for value in numeric_values if abs(value) > 1e-9) / max(1, len(numeric_values))
        score = (coverage * 12) + (non_zero_ratio * 4) + min(spread / mean_abs, 25)
        if is_identifier_header(header):
            score -= 12
        if is_executive_metric_header(header):
            score += 4
        metric_kind = infer_metric_kind(header, numeric_values)
        if metric_kind == "temporal":
            continue
        if metric_kind in {"currency", "percent"}:
            score += 2
        if unique_ratio >= 0.9 and not is_executive_metric_header(header):
            score -= 6
        if len(numeric_values) >= 4 and all(float(value).is_integer() for value in numeric_values[: min(10, len(numeric_values))]):
            deltas = [numeric_values[i + 1] - numeric_values[i] for i in range(min(len(numeric_values) - 1, 8))]
            if deltas and len(set(round(delta, 6) for delta in deltas)) <= 2:
                score -= 8
        if len(numeric_values) >= 5 and all(float(value).is_integer() for value in numeric_values[: min(12, len(numeric_values))]):
            ordered = numeric_values[: min(12, len(numeric_values))]
            deltas = [ordered[i + 1] - ordered[i] for i in range(len(ordered) - 1)]
            if deltas and len(set(round(delta, 6) for delta in deltas)) <= 2 and max(abs(delta) for delta in deltas) <= 10:
                score -= 10
        if score > best_score:
            best_score = score
            best_index = index
    return best_index


def collect_table_candidates(data):
    candidates = []
    if data.get("muestra_tabla"):
        primary_source = resolve_source_sheet_name(data, data.get("muestra_tabla", {}).get("hoja_origen"), "Tabla Principal")
        candidates.append((present_sheet_title(primary_source) or "Tabla Principal", data.get("muestra_tabla"), primary_source))
    for name, table in (data.get("otras_tablas") or {}).items():
        candidates.append((clean_text(name, 70) or "Tabla", table, resolve_source_sheet_name(data, table.get("hoja_origen"), name)))
    for name, table in (data.get("genericas") or {}).items():
        candidates.append((clean_text(name, 70) or "Tabla", table, resolve_source_sheet_name(data, table.get("hoja_origen"), name)))
    return candidates


def score_table_candidate(table_name, table, source_sheet=None):
    headers, rows = extract_table_payload(table)
    if len(headers) < 2 or not rows:
        return -999

    profile = build_table_signal_profile(headers, rows)
    basis = derive_table_basis(headers, rows)
    truth = evaluate_table_truth(headers, rows, basis=basis, source_sheet=source_sheet)
    score = 0.0
    score += (truth.get("score") or 0) * 100
    score += min(profile.get("informative_rows", 0), 12) * 2.5
    score += min(len(rows), 20) * 0.8
    score += min(len(headers), 6) * 1.4
    score += (profile.get("numeric_cell_ratio") or 0) * 18
    score += (1 - min(profile.get("noise_cell_ratio") or 0, 1)) * 10

    if basis.get("numeric_index") is not None:
        score += 10
    if profile.get("detail_allowed"):
        score += 8

    if profile.get("table_kind") == "narrative":
        score -= 8
    if profile.get("avg_text_len", 0) >= 55 and profile.get("sentence_cell_ratio", 0) >= 0.30:
        score -= 14
    if is_low_signal_sheet_name(table_name) or is_low_signal_sheet_name(source_sheet):
        score -= 35

    return score


def build_derived_chart_blocks(data):
    blocks = []
    seen_titles = set()
    ranked_candidates = sorted(
        collect_table_candidates(data),
        key=lambda item: score_table_candidate(item[0], item[1], item[2]),
        reverse=True,
    )
    for table_name, table, source_sheet in ranked_candidates:
        if is_low_signal_sheet_name(table_name) or is_low_signal_sheet_name(source_sheet):
            continue
        headers, rows = extract_table_payload(table)
        if len(headers) < 2 or len(rows) < 3:
            continue

        profile = build_table_signal_profile(headers, rows)
        if profile["informative_rows"] < 3:
            continue
        # Only skip chart generation if the ENTIRE table is highly narrative (avg >50 chars)
        # Short-column tables (checklist, risk matrix) can still produce charts
        if profile["sentence_cell_ratio"] > 0.30 or profile["avg_text_len"] > 50:
            continue

        numeric_index = choose_table_numeric_column(headers, rows)
        if numeric_index is not None:
            normalized_numeric = normalize_header_key(headers[numeric_index])
            if any(token in normalized_numeric for token in ("id", "codigo", "consecutivo", "numero", "nro", "fecha", "date", "month", "mes", "periodo", "semana", "dia", "year", "anio", "ano")):
                numeric_index = None
        dimension_candidates = []
        for index, header in enumerate(headers):
            if index == numeric_index:
                continue
            score = score_chart_dimension(header, rows, index)
            if score <= 0:
                continue
            dimension_candidates.append((score, index))
        dimension_candidates.sort(reverse=True)
        if not dimension_candidates:
            continue

        for _, label_index in dimension_candidates[:2]:
            if numeric_index is not None:
                grouped = build_group_totals(rows, label_index, numeric_index)
                metric_label = clean_text(headers[numeric_index]) or "valor"
            else:
                grouped = build_group_counts(rows, label_index)
                metric_label = "registros"

            grouped = [(label, value) for label, value in grouped if not is_noise_text(label)][:6]
            if len(grouped) < 3:
                continue

            labels = [clean_text(label) for label, _ in grouped]
            values = [float(value) for _, value in grouped]
            chart_type = resolve_derived_chart_type(table_name, headers[label_index], labels, len(grouped))
            limit = 5 if chart_type in {"pie", "doughnut"} else 6
            labels = labels[:limit]
            values = values[:limit]
            total_value = sum(values) or 1
            lead_share = (values[0] / total_value) * 100 if total_value else 0
            metric_kind = infer_metric_kind(metric_label, values)
            exec_table_name = present_sheet_title(table_name)
            title = clean_text(f"{exec_table_name} \u00b7 {headers[label_index]}", 78) or "Grafica derivada"
            if title.casefold() in seen_titles:
                continue
            chart_payload = {
                "tipo": chart_type,
                "titulo": title,
                "labels": labels,
                "valores": values,
                "dimension_label": clean_text(headers[label_index]),
                "metric_label": metric_label,
                "aggregation": "suma" if numeric_index is not None else "conteo",
                "hoja_origen": source_sheet,
            }
            truth = evaluate_chart_truth(chart_payload, source_sheet=source_sheet)
            if truth["blocked"]:
                continue
            seen_titles.add(title.casefold())

            blocks.append({
                "type": "chart",
                "title": title,
                "chart": chart_payload,
                "insight": clean_text(
                    f"{labels[0]} lidera {headers[label_index]} con {format_metric_full(values[0], metric_kind)} en {metric_label} y concentra {lead_share:.0f}% del grupo visualizado."
                ),
                "source_basis": build_source_basis_text(headers[label_index], metric_label, "suma" if numeric_index is not None else "conteo"),
                "traceability": build_traceability(source_sheet, "suma" if numeric_index is not None else "conteo", truth["score"], row_count=len(rows), col_count=len(headers), numeric_coverage=truth["numeric_coverage"]),
            })
            if len(blocks) >= MAX_DERIVED_CHARTS:
                return blocks

    return blocks


def infer_excel_topic(data, excel_path):
    generic = data.get("resumen_generico") or {}
    primary_sheet = clean_text(generic.get("hoja_principal"), 52)
    normalized_primary = normalize_header_key(primary_sheet)
    normalized_columns = " ".join(normalize_header_key(item) for item in safe_list(generic.get("columnas")))
    if any(
        token in normalized_columns or token in normalized_primary
        for token in ("comision", "tiquete", "hospedaje", "transporte", "viatico", "centro de costos", "ciudad destino")
    ):
        return "comisiones, viaticos y gastos de viaje"
    if any(token in normalized_columns for token in ("inventario", "stock", "producto", "bodega")):
        return "inventario, referencias y niveles de stock"
    if any(token in normalized_columns for token in ("venta", "ingreso", "factura", "cliente")):
        return "ventas, ingresos y comportamiento comercial"

    sheet_tokens = " ".join(
        clean_text(name)
        for name in safe_list(data.get("metadatos", {}).get("hojas_encontradas"))
        if not is_low_signal_sheet_name(name)
    ).lower()
    if any(token in sheet_tokens for token in ("auditoria", "hallazgo", "oportunidad", "coso", "control")):
        return "auditoria, controles y oportunidades de mejora"
    if primary_sheet:
        return f"datos operativos concentrados en {primary_sheet}"
    return f"datos ejecutivos extraidos de {basename_label(excel_path)}"


def build_management_highlights(data, excel_path, limit=4):
    highlights = []
    topic = infer_excel_topic(data, excel_path)
    highlights.append(f"El analisis se enfoca en {topic}.")

    summary = data.get("resumen_ejecutivo") or {}
    if summary:
        total = summary.get("total_registros")
        value = try_number(summary.get("valor_total"))
        if total:
            highlights.append(f"Se analizaron {format_metric_full(total)} registros base.")
        if value is not None and value > 0:
            highlights.append(f"El valor total observado asciende a {format_metric_full(value, 'currency')}.")

    financial = data.get("contexto_financiero") or {}
    currencies = safe_list(financial.get("currencies_detected"))
    if currencies:
        if financial.get("has_mixed_currency") and not financial.get("has_unresolved_conversion"):
            highlights.append(f"Se detectan monedas mixtas ({', '.join(currencies[:3])}) y el motor normaliza los montos con tasa explicita cuando esta disponible.")
        elif financial.get("has_unresolved_conversion"):
            highlights.append(f"Se detectan monedas {', '.join(currencies[:3])} con conversion parcial; los acumulados financieros se calculan solo con filas resueltas.")
        elif currencies != ["COP"]:
            highlights.append(f"La lectura financiera identifica montos en {', '.join(currencies[:3])} y mantiene trazabilidad de la moneda fuente.")

    quality = data.get("calidad_datos") or {}
    quality_score = quality.get("quality_score")
    if quality_score is not None and float(quality_score) < 0.9:
        highlights.append(f"La validacion de calidad estima una integridad del {int(round(float(quality_score) * 100))}%.")
    duplicate_rows = try_number(quality.get("duplicate_rows")) or 0
    total_rows = try_number((data.get("resumen_generico") or {}).get("total_filas")) or 0
    if duplicate_rows and total_rows and (duplicate_rows / max(total_rows, 1)) >= 0.12:
        highlights.append(f"El archivo exige depuracion previa porque el volumen de duplicados puede distorsionar la lectura gerencial.")
    sparse_columns = safe_list(quality.get("sparse_numeric_columns"))
    if sparse_columns and total_rows >= 20:
        column_name = clean_text((sparse_columns[0] or {}).get("columna"), 28)
        missing_pct = int(round(((sparse_columns[0] or {}).get("missing_ratio") or 0) * 100))
        if column_name and missing_pct >= 35:
            highlights.append(f"La columna {column_name} presenta {missing_pct}% de vacios y se interpreta de forma conservadora.")

    generic = data.get("resumen_generico") or {}
    if generic:
        total_rows = generic.get("total_filas")
        total_cols = generic.get("total_columnas")
        if total_rows and total_cols:
            highlights.append(
                f"La hoja principal consolida {format_metric_full(total_rows)} filas y {format_metric_full(total_cols)} columnas utiles."
            )

    advanced = data.get("analisis_avanzado") or {}
    for insight in safe_list(advanced.get("insights")):
        if isinstance(insight, dict):
            text = compact_line(insight.get("texto"), 96)
            if text:
                highlights.append(text)

    for conclusion in unique_texts(data.get("conclusiones"), limit=3, min_len=18):
        highlights.append(compact_line(conclusion, 96))

    for chart in safe_list(data.get("graficas_automaticas"))[:2]:
        insight = compact_line(chart.get("insight_auto"), 96)
        if insight:
            highlights.append(insight)

    curated = []
    low_value_tokens = (
        "filas potencialmente duplicadas",
        "hojas detectadas",
        "hojas funcionales",
        "archivo corresponde",
        "base principal",
        "columnas utiles",
        "integridad del",
        "vacios y se interpreta",
        "motor",
        "sistema",
        "placeholders",
        "celdas fantasma",
    )
    for line in unique_texts(highlights, limit=limit * 3, min_len=14):
        normalized = clean_text(line).lower()
        if any(token in normalized for token in low_value_tokens):
            continue
        curated.append(line)
    return curated[:limit]


def get_primary_table_profile(data):
    if not data.get("muestra_tabla"):
        return None, []
    headers, rows = extract_table_payload(data.get("muestra_tabla"))
    if len(headers) < 2 or not rows:
        return None, []
    return build_table_signal_profile(headers, rows), headers


SEMANTIC_HEADER_GROUPS = {
    "objective": ("objetivo", "objetivos", "objective", "alcance", "proposito", "purpose", "meta"),
    "activity": ("actividad", "actividades", "activity", "accion", "acciones", "tarea", "tareas", "plan de trabajo", "plan"),
    "question": ("pregunta", "preguntas", "question", "consulta", "criterio", "prueba", "chequeo", "validacion"),
    "risk": ("riesgo", "riesgos", "hallazgo", "hallazgos", "issue", "problema", "observacion", "alerta"),
    "owner": ("responsable", "encargado", "owner", "lider", "líder"),
}


def normalize_header_key(value):
    text = clean_text(value).casefold()
    replacements = {
        "á": "a",
        "é": "e",
        "í": "i",
        "ó": "o",
        "ú": "u",
        "ñ": "n",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return " ".join(text.split()).strip()


def match_semantic_role(header):
    normalized = normalize_header_key(header)
    if not normalized:
        return None
    for role, aliases in SEMANTIC_HEADER_GROUPS.items():
        if any(alias in normalized for alias in aliases):
            return role
    return None


def build_semantic_samples(data, per_role_limit=3):
    samples = {key: [] for key in SEMANTIC_HEADER_GROUPS}
    role_headers = {key: [] for key in SEMANTIC_HEADER_GROUPS}

    for table_name, table, _source_sheet in collect_table_candidates(data):
        headers, rows = extract_table_payload(table)
        if len(headers) < 2 or not rows:
            continue

        for index, header in enumerate(headers):
            role = match_semantic_role(header)
            if not role:
                continue

            role_headers[role].append(clean_text(header, 28))
            for row in rows:
                if index >= len(row):
                    continue
                value = clean_text(row[index], 180)
                if is_noise_text(value):
                    continue
                if role in {"objective", "activity", "question", "risk"} and len(value) < 10:
                    continue
                if role == "owner":
                    if looks_like_sentence(value) or len(value) > 42:
                        continue
                    candidate = clean_text(value, 42)
                else:
                    header_label = clean_text(header, 24)
                    candidate = compact_line(f"{header_label}: {value}" if header_label else value, 96)
                if candidate:
                    samples[role].append(candidate)

    return {
        "samples": {role: unique_texts(values, limit=per_role_limit, min_len=12) for role, values in samples.items()},
        "headers": {role: unique_texts(values, limit=3, min_len=4) for role, values in role_headers.items()},
    }


def build_real_objectives_from_semantics(data, limit=3):
    semantic = build_semantic_samples(data, per_role_limit=3)
    samples = semantic["samples"]
    objectives = []

    if samples["objective"]:
        objectives.append(f"Abordar los objetivos declarados en el Excel, empezando por {compact_line(samples['objective'][0], 70)}.")
    if samples["risk"]:
        objectives.append(f"Escalar riesgos o hallazgos reales como {compact_line(samples['risk'][0], 68)}.")
    if samples["question"]:
        objectives.append(f"Responder preguntas clave del archivo, incluyendo {compact_line(samples['question'][0], 68)}.")
    if samples["activity"]:
        objectives.append(f"Traducir actividades del Excel en un plan de seguimiento gerencial accionable.")
    if samples["owner"]:
        objectives.append("Vincular frentes de trabajo con responsables o areas cuando el Excel los identifica.")

    return unique_texts(objectives, limit=limit, min_len=18)


def build_workplan_lines(data, limit=4):
    semantic = build_semantic_samples(data, per_role_limit=3)
    samples = semantic["samples"]
    lines = []

    if samples["activity"]:
        activity_line = compact_line(samples["activity"][0], 78)
        normalized_activity = clean_text(activity_line).lower()
        if not any(token in normalized_activity for token in ("accion correctiva", "oportunidad de mejora", "hallazgo", "procedimiento")):
            lines.append(f"Frente de trabajo: {activity_line}.")
    if samples["question"]:
        lines.append(f"Pregunta de control: {compact_line(samples['question'][0], 78)}.")
    if samples["risk"]:
        lines.append(f"Riesgo a monitorear: {compact_line(samples['risk'][0], 78)}.")
    if samples["owner"]:
        lines.append(f"Responsable o area visible: {compact_line(samples['owner'][0], 78)}.")
    if samples["objective"]:
        lines.append(f"Objetivo operativo: {compact_line(samples['objective'][0], 78)}.")

    return unique_texts(lines, limit=limit, min_len=18)


SHEET_TYPE_RULES = {
    "workplan": ("plan de trabajo", "cronograma", "actividad", "actividades", "seguimiento", "plan"),
    "findings": ("hallazgo", "hallazgos", "observacion", "observaciones", "no conform", "mejora"),
    "evidence": ("evidencia", "evidencias", "soporte", "documento", "documental", "prueba"),
    "risk": ("riesgo", "riesgos", "control", "alerta", "criticidad"),
    "owner": ("responsable", "responsables", "owner", "encargado", "lider", "líder"),
    "timeline": ("cronograma", "fecha", "fechas", "periodo", "mes", "semana", "trimestre", "timeline"),
}


def classify_sheet_type(table_name, headers, rows):
    combined = " ".join([table_name] + list(headers))
    normalized = normalize_header_key(combined)
    scores = {sheet_type: 0 for sheet_type in SHEET_TYPE_RULES}

    for sheet_type, keywords in SHEET_TYPE_RULES.items():
        for keyword in keywords:
            normalized_keyword = normalize_header_key(keyword)
            if normalized_keyword and normalized_keyword in normalized:
                scores[sheet_type] += 2

    semantic = build_table_signal_profile(headers, rows)
    semantic_headers = [match_semantic_role(header) for header in headers]
    for role in semantic_headers:
        if role == "objective" or role == "activity":
            scores["workplan"] += 1
        elif role == "question":
            scores["findings"] += 1
        elif role == "risk":
            scores["risk"] += 2
        elif role == "owner":
            scores["owner"] += 2

    if semantic["table_kind"] == "narrative":
        scores["evidence"] += 1
        scores["findings"] += 1

    winner = max(scores.items(), key=lambda item: item[1])
    if winner[1] <= 1:
        return "generic"
    return winner[0]


def build_native_sheet_highlights(sheet_type, table_name, headers, rows):
    lines = []
    semantic = build_semantic_samples({"muestra_tabla": {"encabezados": headers, "filas": rows}}, per_role_limit=2)["samples"]

    if sheet_type == "workplan":
        if semantic["activity"]:
            lines.append(f"Actividad principal: {semantic['activity'][0]}.")
        if semantic["objective"]:
            lines.append(f"Objetivo asociado: {semantic['objective'][0]}.")
        if semantic["owner"]:
            lines.append(f"Responsable identificado: {semantic['owner'][0]}.")
        lines.append("Esta hoja funciona como plan de trabajo y conviene presentarla por frentes, responsables y seguimiento.")
    elif sheet_type == "findings":
        if semantic["risk"]:
            lines.append(f"Hallazgo o riesgo clave: {semantic['risk'][0]}.")
        if semantic["question"]:
            lines.append(f"Pregunta de control asociada: {semantic['question'][0]}.")
        lines.append("La hoja debe escalar solo hallazgos repetitivos, criticos o con impacto gerencial.")
    elif sheet_type == "evidence":
        if semantic["question"]:
            lines.append(f"Prueba o evidencia central: {semantic['question'][0]}.")
        lines.append("La hoja contiene soportes o evidencias y debe sintetizarse, evitando saturar con tabla documental completa.")
    elif sheet_type == "risk":
        if semantic["risk"]:
            lines.append(f"Riesgo principal: {semantic['risk'][0]}.")
        lines.append("La lectura debe priorizar criticidad, exposicion y acciones de mitigacion.")
    elif sheet_type == "owner":
        if semantic["owner"]:
            lines.append(f"Responsable visible: {semantic['owner'][0]}.")
        lines.append("La hoja facilita asignacion de frentes y seguimiento por responsable o area.")
    elif sheet_type == "timeline":
        lines.append("La hoja tiene un comportamiento temporal y conviene llevarla a comparativos o seguimiento por hitos.")

    return unique_texts([clean_text(item) for item in lines], limit=4, min_len=18)


def build_native_sheet_briefing(data, limit=4):
    lines = []
    for table_name, table, _source_sheet in collect_table_candidates(data):
        headers, rows = extract_table_payload(table)
        if len(headers) < 2 or not rows:
            continue
        sheet_type = classify_sheet_type(table_name, headers, rows)
        if sheet_type == "generic":
            continue
        label_map = {
            "workplan": "plan de trabajo",
            "findings": "hallazgos",
            "evidence": "evidencias",
            "risk": "riesgos",
            "owner": "responsables",
            "timeline": "cronograma",
        }
        lines.append(f"La hoja {table_name} se interpreta como {label_map.get(sheet_type, sheet_type)}.")
    return unique_texts(lines, limit=limit, min_len=18)


def build_technical_highlights(data, excel_path, limit=4):
    generic = data.get("resumen_generico") or {}
    summary = data.get("resumen_ejecutivo") or {}
    metadata = data.get("metadatos") or {}
    budget = data.get("presupuesto_slides") or {}
    table_profile, headers = get_primary_table_profile(data)

    lines = []
    sheet_count = len(safe_list(metadata.get("hojas_encontradas")))
    if sheet_count:
        lines.append(f"El archivo contiene {format_metric_full(sheet_count)} hojas detectadas.")

    primary_sheet = clean_text(generic.get("hoja_principal"), 38)
    if primary_sheet:
        lines.append(f"La hoja principal priorizada es {primary_sheet}.")

    total_rows = generic.get("total_filas")
    total_cols = generic.get("total_columnas")
    if total_rows and total_cols:
        lines.append(
            f"La base principal aporta {format_metric_full(total_rows)} filas y {format_metric_full(total_cols)} columnas utiles."
        )
    elif total_rows:
        lines.append(f"La base principal aporta {format_metric_full(total_rows)} registros utiles.")

    kpi_count = len(safe_list(data.get("kpis_automaticos")))
    chart_count = len(safe_list(data.get("graficas_automaticas")))
    table_count = len(collect_table_candidates(data))
    lines.append(
        f"El motor identifico {format_metric_full(kpi_count)} KPIs, {format_metric_full(chart_count)} graficas y {format_metric_full(table_count)} tablas utilizables."
    )

    if table_profile:
        lines.append(
            f"La tabla principal se clasifica como {('narrativa' if table_profile['table_kind'] == 'narrative' else 'operativa')} con {format_metric_full(table_profile['informative_rows'])} filas informativas."
        )
        if headers:
            sample_headers = ", ".join([clean_text(item, 18) for item in headers[:3] if clean_text(item, 18)])
            if sample_headers:
                lines.append(f"Los campos dominantes del analisis son {sample_headers}.")

    if budget:
        slide_budget = budget.get("total_estimado")
        if slide_budget:
            lines.append(f"El planeamiento proyecta {format_metric_full(slide_budget)} slides de contenido util.")

    semantic = build_semantic_samples(data, per_role_limit=2)
    detected_roles = [role for role, values in semantic["samples"].items() if values]
    if detected_roles:
        lines.append(
            f"Se detectaron columnas semanticas para {', '.join(detected_roles[:4])}, utiles para construir objetivos y plan de trabajo."
        )

    financial = data.get("contexto_financiero") or {}
    if safe_list(financial.get("currency_columns")):
        monitored = safe_list(financial.get("currency_columns"))[:2]
        summary = []
        for item in monitored:
            column_name = clean_text(item.get("columna"), 24)
            currencies = "/".join(safe_list(item.get("currencies_detected"))[:3]) or "COP"
            resolved_pct = int(round((item.get("resolved_ratio") or 0) * 100))
            summary.append(f"{column_name}: {currencies} ({resolved_pct}% resuelto)")
        if summary:
            lines.append(f"Contexto financiero monitoreado: {'; '.join(summary)}.")
    quality = data.get("calidad_datos") or {}
    if quality.get("duplicate_rows"):
        lines.append(f"Control de consistencia: {format_metric_full(quality.get('duplicate_rows'))} filas duplicadas detectadas para depuracion.")
    placeholder_pct = int(round(float(quality.get("placeholder_ratio") or 0) * 100))
    if placeholder_pct > 0:
        lines.append(f"Los placeholders o celdas fantasma representan {placeholder_pct}% del contenido textual revisado.")
    if safe_list(quality.get("currency_alerts")):
        alert = quality["currency_alerts"][0]
        lines.append(
            f"Se valida conversion financiera en {clean_text(alert.get('columna'), 24)} con {int(round((alert.get('resolved_ratio') or 0) * 100))}% de cobertura monetaria."
        )
    lines.extend(build_native_sheet_briefing(data, limit=2))

    return unique_texts([clean_text(item) for item in lines], limit=limit, min_len=18)


def build_dynamic_objectives(data, excel_path, limit=3):
    topic = infer_excel_topic(data, excel_path)
    objectives = [
        f"Explicar con claridad de gerencia el alcance del Excel sobre {topic}.",
        "Ordenar la conversacion en indicadores, hallazgos, prioridades y decisiones posibles.",
    ]

    if len(safe_list(data.get("graficas_automaticas"))) > 0:
        objectives.append("Mostrar tendencias y concentraciones visibles en las graficas generadas desde el Excel.")
    if data.get("analisis_avanzado"):
        objectives.append("Escalar riesgos, anomalias o focos de control detectados en el analisis tecnico.")
    if data.get("muestra_tabla"):
        objectives.append("Usar la tabla principal como soporte tecnico solo cuando aporte lectura real y no ruido.")
    financial = data.get("contexto_financiero") or {}
    if financial.get("has_mixed_currency"):
        objectives.append("Aclarar la lectura financiera cuando el Excel mezcle monedas y priorizar montos convertidos de forma trazable.")
    if IS_BOARDROOM_MODE:
        objectives.append("Traducir la evidencia del Excel a un lenguaje boardroom con alertas y prioridades.")
    objectives.extend(build_real_objectives_from_semantics(data, limit=3))
    objectives.extend(build_native_sheet_briefing(data, limit=2))

    return unique_texts([clean_text(item) for item in objectives], limit=limit, min_len=18)


def build_presentation_topics(data, limit=4):
    topics = []
    visual_curation_ready = is_visual_curation_ready(data)

    for kpi in safe_list(data.get("kpis_automaticos"))[:3]:
        label = clean_text((kpi or {}).get("label"), 36).lower()
        if not label:
            continue
        if "total registros" in label:
            topics.append("Dimension y volumen del dataset.")
        elif "promedio" in label:
            topics.append("Promedios y dispersion del indicador principal.")
        elif "total" in label:
            topics.append("Acumulados principales con impacto para la decision.")
        else:
            topics.append("Indicadores clave para priorizar seguimiento.")

    if visual_curation_ready:
        for chart in safe_list(data.get("graficas_automaticas"))[:2]:
            title = clean_text((chart or {}).get("titulo"), 42)
            if title and not is_low_signal_sheet_name(title):
                topics.append(f"Tendencia visual: {title}.")

    if data.get("analisis_avanzado"):
        advanced = data.get("analisis_avanzado") or {}
        if safe_list(advanced.get("pareto")):
            topics.append("Concentraciones y distribuciones relevantes.")
        if safe_list(advanced.get("outliers")):
            topics.append("Anomalias y valores atipicos detectados.")
        if advanced.get("tendencia"):
            topics.append("Cambios temporales y comparativos del periodo.")

    table_profile, _ = get_primary_table_profile(data)
    if table_profile:
        if table_profile["table_kind"] == "narrative":
            topics.append("Focos documentales, responsables y mensajes accionables.")
        else:
            topics.append("Tabla principal con soporte operativo para profundizar.")

    financial = data.get("contexto_financiero") or {}
    if financial.get("has_mixed_currency"):
        topics.append("Normalizacion financiera y conversion de moneda con evidencia explicita.")
    elif safe_list(financial.get("currencies_detected")) not in ([], ["COP"]):
        topics.append("Lectura financiera multimoneda con trazabilidad.")

    topics.extend(build_workplan_lines(data, limit=3))
    topics.extend(build_native_sheet_briefing(data, limit=2))

    filtered = []
    for item in unique_texts([clean_text(value) for value in topics], limit=limit * 2, min_len=16):
        normalized = clean_text(item).lower()
        if any(token in normalized for token in ("base principal", "hoja ", "archivo corresponde")):
            continue
        filtered.append(item)
    return filtered[:limit]


def build_intro_briefing(data, excel_path):
    summary = build_management_highlights(data, excel_path, limit=4)
    technical = build_technical_highlights(data, excel_path, limit=4)
    objectives = build_dynamic_objectives(data, excel_path, limit=3)
    topics = build_presentation_topics(data, limit=4)
    workplan = build_workplan_lines(data, limit=4)
    executive_ai = data.get("briefing_ejecutivo_ia") if isinstance(data.get("briefing_ejecutivo_ia"), dict) else None

    return {
        "summary": summary,
        "technical": technical,
        "objectives": objectives,
        "topics": topics,
        "workplan": workplan,
        "executive_ai": executive_ai,
    }


def build_showcase_items(data):
    items = []
    if IS_BOARDROOM_MODE and is_visual_curation_ready(data):
        items.append("semaforos ejecutivos")
        items.append("riesgos y oportunidades")
        items.append("comparativos antes/despues")
        items.append("top 5 con alertas")
    if is_visual_curation_ready(data) and len(build_chart_blocks(data)) > 0:
        items.append("graficas priorizadas")
    if is_visual_curation_ready(data) and (data.get("muestra_tabla") or data.get("otras_tablas") or data.get("genericas")):
        items.append("tablas resumidas")
    return unique_texts(items, limit=4, min_len=6)


def build_boardroom_traffic_payload(data):
    advanced = data.get("analisis_avanzado") or {}
    items = []

    tendencia = advanced.get("tendencia") or {}
    if tendencia:
        change = try_number(tendencia.get("cambio_pct"))
        status = "verde"
        if tendencia.get("tendencia") == "decreciente":
            status = "rojo" if (change is not None and abs(change) >= 15) else "amarillo"
        elif tendencia.get("tendencia") == "estable":
            status = "amarillo"
        items.append({
            "label": "Tendencia",
            "status": status,
            "metric": f"{change:+.1f}%" if change is not None else clean_text(tendencia.get("tendencia"), 18),
            "detail": compact_line(
                f"{clean_text(tendencia.get('columna_valor'), 32)} con tendencia {clean_text(tendencia.get('tendencia'), 18)}.",
                66,
            ),
        })

    paretos = safe_list(advanced.get("pareto"))
    if paretos:
        top = paretos[0]
        pct = try_number(top.get("lider_pct"))
        items.append({
            "label": "Concentracion",
            "status": status_label_from_value(pct, 30, 45),
            "metric": f"{pct:.0f}%" if pct is not None else "N/D",
            "detail": compact_line(
                f"{clean_text(top.get('lider'), 28)} lidera {clean_text(top.get('columna_categoria'), 28)}.",
                66,
            ),
        })

    outliers = safe_list(advanced.get("outliers"))
    if outliers:
        top = outliers[0]
        total = try_number(top.get("total_outliers")) or 0
        items.append({
            "label": "Anomalias",
            "status": "rojo" if total >= 3 else "amarillo",
            "metric": format_metric_full(total),
            "detail": compact_line(
                f"Outliers detectados en {clean_text(top.get('columna'), 30)}.",
                66,
            ),
        })

    quality_issue = None
    for insight in safe_list(advanced.get("insights")):
        if isinstance(insight, dict) and clean_text(insight.get("tipo")).lower() == "calidad_datos":
            quality_issue = insight
            break
    if quality_issue:
        items.append({
            "label": "Calidad de datos",
            "status": "amarillo",
            "metric": "Revision",
            "detail": compact_line(quality_issue.get("texto"), 66),
        })

    if len(items) < 2:
        return None

    return {
        "type": "boardroom_traffic",
        "title": "Semaforo de Control",
        "subtitle": "Estado sintetico de los principales frentes de auditoria",
        "items": items[:4],
    }


def classify_insight_direction(insight):
    insight_type = clean_text((insight or {}).get("tipo")).lower()
    text = clean_text((insight or {}).get("texto")).lower()
    if insight_type in {"anomalia", "calidad_datos"}:
        return "risk"
    if insight_type == "tendencia" and "decreciente" in text:
        return "risk"
    if insight_type == "concentracion" and any(token in text for token in ("critica", "alta", "solo")):
        return "risk"
    if insight_type in {"correlacion"}:
        return "opportunity"
    if insight_type == "tendencia" and "creciente" in text:
        return "opportunity"
    if insight_type == "desbalance":
        return "risk"
    return "opportunity"


def build_risk_opportunity_payload(data):
    advanced = data.get("analisis_avanzado") or {}
    risks = []
    opportunities = []

    for insight in safe_list(advanced.get("insights")):
        if not isinstance(insight, dict):
            continue
        card = {
            "title": clean_text(insight.get("tipo"), 24).replace("_", " ").title() or "Insight",
            "body": compact_line(insight.get("texto"), 98),
            "action": compact_line(insight.get("accion"), 92),
        }
        if not card["body"]:
            continue
        if classify_insight_direction(insight) == "risk":
            risks.append(card)
        else:
            opportunities.append(card)

    if not risks:
        for line in build_management_highlights(data, "archivo", limit=2):
            risks.append({
                "title": "Seguimiento",
                "body": compact_line(line, 98),
                "action": "Mantener monitoreo de este frente en la revision gerencial.",
            })
            break

    if not opportunities:
        fallback_lines = []
        for line in build_management_highlights(data, "archivo", limit=4):
            normalized = clean_text(line).lower()
            if any(token in normalized for token in ("atipico", "faltantes", "outlier", "anomalia")):
                continue
            if line.startswith("'"):
                continue
            fallback_lines.append(line)
        for line in fallback_lines[:1]:
            opportunities.append({
                "title": "Oportunidad",
                "body": line,
                "action": "Usar esta concentracion para priorizar decisiones y recursos.",
            })

    if not risks and not opportunities:
        return None

    return {
        "type": "risk_opportunity_cards",
        "title": "Riesgos y Oportunidades",
        "subtitle": "Analisis de brechas construido con datos reales del Excel",
        "risks": risks[:2],
        "opportunities": opportunities[:2],
    }


def build_before_after_payload(data):
    advanced = data.get("analisis_avanzado") or {}
    tendencia = advanced.get("tendencia") or {}
    if not tendencia or not tendencia.get("columna_valor"):
        return None

    start = try_number(tendencia.get("promedio_inicio"))
    end = try_number(tendencia.get("promedio_fin"))
    if start is None and end is None:
        return None
    change = try_number(tendencia.get("cambio_pct"))
    direction = "mejora" if (change is not None and change > 0) else "ajuste"
    metric_kind = infer_metric_kind(tendencia.get("columna_valor"), [value for value in (start, end) if value is not None])

    return {
        "type": "before_after_compare",
        "title": f"Comparativo antes/despues de {clean_text(tendencia.get('columna_valor'), 42)}",
        "subtitle": "Evolucion temporal inferida desde datos del propio Excel",
        "before_label": "Inicio",
        "before_value": format_metric_full(start, metric_kind),
        "after_label": "Cierre",
        "after_value": format_metric_full(end, metric_kind),
        "delta": f"{change:+.1f}%" if change is not None else "N/D",
        "commentary": compact_line(
            f"El promedio pasa de {format_metric_full(start, metric_kind)} a {format_metric_full(end, metric_kind)}, lo que sugiere {direction} en el periodo.",
            106,
        ),
    }


def build_top5_alerts_payload(data):
    for table_name, table, _source_sheet in collect_table_candidates(data):
        headers, rows = extract_table_payload(table)
        if len(headers) < 2 or len(rows) < 5:
            continue

        basis = derive_table_basis(headers, rows)
        label_index = basis["label_index"]
        numeric_index = basis["numeric_index"]
        if label_index is None:
            continue

        grouped = build_group_totals(rows, label_index, numeric_index) if numeric_index is not None else build_group_counts(rows, label_index)
        if len(grouped) < 3:
            continue

        values = [try_number(item[1]) or 0 for item in grouped[:MAX_BOARDROOM_ITEMS]]
        avg_value = sum(values) / len(values)
        critical = avg_value * 1.45
        warning = avg_value * 1.15
        items = []
        for label, value in grouped[:MAX_BOARDROOM_ITEMS]:
            raw_value = try_number(value) or 0
            status = status_label_from_value(raw_value, warning, critical)
            items.append({
                "label": clean_text(label, 34),
                "raw_value": raw_value,
                "value": format_metric_full(raw_value),
                "status": status,
                "detail": "Supera el umbral critico." if status == "rojo" else ("Requiere seguimiento." if status == "amarillo" else "Dentro del rango esperado."),
            })

        threshold_text = f"Promedio top 5: {format_metric_full(avg_value)}"
        return {
            "type": "top5_alerts",
            "title": f"Top 5 con alertas de {table_name}",
            "subtitle": "Umbrales calculados automaticamente sobre datos reales",
            "threshold": threshold_text,
            "items": items,
        }
    return None


def build_boardroom_blocks(data, excel_path):
    blocks = []
    topic = infer_excel_topic(data, excel_path)
    focus_phrase = build_prompt_focus_phrase(data, excel_path, max_len=66)
    summary_lines = build_management_highlights(data, excel_path, limit=4)
    if summary_lines:
        blocks.append({
            "type": "text",
            "title": compact_line(f"Panorama para decidir: {focus_phrase}", 78),
            "subtitle": compact_line(f"Lectura gerencial sobre {topic}", 84) if topic else "Lectura gerencial basada en datos reales",
            "lines": summary_lines,
        })

    for builder in (
        build_boardroom_traffic_payload,
        build_risk_opportunity_payload,
        build_before_after_payload,
        build_top5_alerts_payload,
    ):
        payload = builder(data)
        if payload:
            blocks.append(payload)
    return blocks


def build_table_summary_lines(table_name, headers, rows):
    profile = build_table_signal_profile(headers, rows)
    row_count = len(rows)
    col_count = len(headers)
    basis = derive_table_basis(headers, rows)
    label_index = basis["label_index"]
    numeric_index = basis["numeric_index"]
    lines = [
        f"El analisis cubre {row_count} registros y {col_count} variables relevantes.",
    ]

    if profile["table_kind"] == "narrative":
        lines.append("La estructura es principalmente documental y conviene leerla como resumen ejecutivo, no como detalle operativo.")
        if profile["sentence_cell_ratio"] >= 0.2:
            lines.append("Predominan descripciones extensas, observaciones o actividades que requieren sintesis por focos.")

    if numeric_index is not None:
        numeric_header = clean_text(headers[numeric_index], 40) or "valor"
        numeric_values = [try_number(row[numeric_index]) for row in rows if numeric_index < len(row)]
        numeric_values = [value for value in numeric_values if value is not None]
        if numeric_values:
            total = sum(numeric_values)
            average = total / len(numeric_values)
            metric_kind = infer_metric_kind(numeric_header, numeric_values)
            lines.append(
                f"El indicador {numeric_header} suma {format_metric_full(total, metric_kind)} y promedia {format_metric_full(average, metric_kind)}."
            )

            if label_index is not None and label_index < len(headers):
                ranked_rows = []
                for row in rows:
                    if label_index >= len(row) or numeric_index >= len(row):
                        continue
                    label = clean_text(row[label_index], 42)
                    numeric_value = try_number(row[numeric_index])
                    if label and not is_noise_text(label) and numeric_value is not None:
                        ranked_rows.append((label, numeric_value))
                ranked_rows.sort(key=lambda item: item[1], reverse=True)
                if ranked_rows:
                    top_label, top_value = ranked_rows[0]
                    lines.append(
                        f"El mayor peso se concentra en {top_label} con {format_metric(top_value, metric_kind)}."
                    )

    if label_index is not None and label_index < len(headers):
        counts = {}
        for row in rows:
            if label_index >= len(row):
                continue
            label = clean_text(row[label_index], 36)
            if label and not is_noise_text(label):
                counts[label] = counts.get(label, 0) + 1
        if counts:
            top_label, top_count = sorted(counts.items(), key=lambda item: item[1], reverse=True)[0]
            if top_count > 1 or len(counts) <= 3:
                lines.append(f"La categoria dominante es {top_label} con {top_count} registros.")

    if "hallazgo" in clean_text(table_name).lower():
        lines.append("Se recomienda escalar solo los casos repetitivos o de mayor impacto para decision.")
    elif "coso" in clean_text(table_name).lower():
        lines.append("La lectura prioriza los componentes con mayor concentracion de observaciones.")
    elif profile["table_kind"] == "narrative":
        lines.append("Se recomienda convertir esta hoja en mensajes accionables y responsables, evitando sobrecargar la gerencia con tabla cruda.")

    return unique_texts([compact_line(line, 92) for line in lines], limit=4, min_len=16)


def build_narrative_focus_lines(table_name, headers, rows):
    profile = build_table_signal_profile(headers, rows)
    if profile["table_kind"] != "narrative":
        return []

    candidates = []
    for header_index, header in enumerate(headers):
        if try_number(header) is not None:
            continue
        header_label = clean_text(header, 28) or "Campo"
        values = []
        for row in rows:
            if header_index >= len(row):
                continue
            value = clean_text(row[header_index], 180)
            if is_noise_text(value) or not looks_like_sentence(value):
                continue
            values.append(f"{header_label}: {compact_line(value, 88)}")
        candidates.extend(values[:2])

    if not candidates:
        first_row = rows[0] if rows else []
        for index, value in enumerate(first_row[: min(len(headers), 4)]):
            if is_noise_text(value):
                continue
            header_label = clean_text(headers[index], 28) or f"Campo {index + 1}"
            candidates.append(f"{header_label}: {compact_line(value, 88)}")

    prefix = compact_line(
        f"La hoja {table_name} es mayoritariamente narrativa y se resume por frentes accionables.",
        92,
    )
    result = [prefix] if prefix else []
    result.extend(unique_texts(candidates, limit=3, min_len=18))
    return result[:4]


def split_long_text(text, max_len=110):
    cleaned = clean_text(text)
    if not cleaned:
        return []
    if len(cleaned) <= max_len:
        return [cleaned]

    segments = []
    current = ""
    for word in cleaned.split():
        candidate = f"{current} {word}".strip()
        if len(candidate) <= max_len:
            current = candidate
            continue
        if current:
            segments.append(current)
        current = word
    if current:
        segments.append(current)
    return segments


def wrap_chart_label(text, max_len=16, max_lines=3):
    parts = split_long_text(text, max_len=max_len)
    if not parts:
        return ""
    return "\n".join(parts[:max_lines])


def build_chart_layout_profile(labels, values):
    count = len(labels)
    longest_label = max((len(clean_text(label)) for label in labels), default=8)
    max_value = max((abs(value) for value in values), default=0)
    if count >= 5:
        label_max_len = 12
        label_font = 9.0
    elif count == 4:
        label_max_len = 14
        label_font = 9.4
    else:
        label_max_len = 16
        label_font = 9.8
    if longest_label >= 26:
        label_font = max(8.2, label_font - 0.8)
    value_font = 9.2 if max_value >= 1_000_000 else 9.6
    return {
        "count": count,
        "longest_label": longest_label,
        "label_max_len": label_max_len,
        "label_font": label_font,
        "value_font": value_font,
        "axis_font": 8.1 if count >= 5 else 8.6,
        "prefer_compact_callout": count >= 5 or max_value >= 1_000_000,
        "left_margin": min(0.34, max(0.22, longest_label / 105)),
        "bottom_margin": 0.22 if count >= 5 else 0.18,
    }


def draw_line_annotation(ax, x_pos, y_pos, text, above=True, offset_x=0):
    offset_y = 12 if above else -16
    va = "bottom" if above else "top"
    ax.annotate(
        text,
        (x_pos, y_pos),
        textcoords="offset points",
        xytext=(offset_x, offset_y),
        ha="center",
        va=va,
        fontsize=8.6,
        color="#243041",
        fontweight="bold",
        bbox={"boxstyle": "round,pad=0.18", "fc": "white", "ec": "none", "alpha": 0.92},
        zorder=5,
        clip_on=False,
    )


def draw_bar_value_label(ax, bar, value, max_value, text):
    width = bar.get_width()
    y_center = bar.get_y() + (bar.get_height() / 2)
    inside = max_value > 0 and (width / max_value) >= 0.34 and len(text) <= 18
    if inside:
        ax.text(
            max(width - max(max_value * 0.03, 0.06), width * 0.55),
            y_center,
            text,
            va="center",
            ha="right",
            fontsize=8.8,
            color="white",
            fontweight="bold",
            clip_on=True,
        )
        return
    ax.text(
        width + max(max_value * 0.025, 0.06),
        y_center,
        text,
        va="center",
        ha="left",
        fontsize=8.8,
        color="#243041",
        fontweight="bold",
        clip_on=False,
    )


def summarize_chart_payload(chart):
    metric_kind = infer_metric_kind(chart.get("metric_label") or chart.get("titulo"), safe_list(chart.get("valores")))
    _, pairs = normalize_chart_pairs(chart)
    if not pairs:
        return {
            "top_label": "N/D",
            "top_value": "N/D",
            "share": "N/D",
            "count": 0,
            "total_value": "N/D",
            "metric_kind": metric_kind,
        }
    total = sum(value for _, value in pairs) or 1
    top_label, top_value = max(pairs, key=lambda item: item[1])
    share = (top_value / total) * 100 if total else 0
    return {
        "top_label": top_label,
        "top_value": format_value(top_value, metric_kind, compact=False),
        "share": f"{share:.0f}%",
        "count": len(pairs),
        "total_value": format_value(total, metric_kind, compact=False),
        "metric_kind": metric_kind,
    }


def normalize_chart_pairs(chart):
    chart_type = clean_text(chart.get("tipo")).lower()
    labels = [clean_text(item) for item in safe_list(chart.get("labels"))]
    values = []
    for value in safe_list(chart.get("valores")):
        values.append(try_number(value) or 0)

    pairs = [(label, value) for label, value in zip(labels, values) if label and value > 0]
    if not pairs:
        return chart_type, []

    if chart_type in {"pie", "doughnut"}:
        max_items = 4
    elif chart_type == "bar":
        max_items = 5
    else:
        max_items = 6

    if chart_type in {"pie", "doughnut", "bar"} and len(pairs) > max_items:
        visible = pairs[: max_items - 1]
        remaining_total = sum(value for _, value in pairs[max_items - 1 :])
        if remaining_total > 0:
            visible.append(("Otros", remaining_total))
        pairs = visible
    elif chart_type == "line" and len(pairs) > max_items:
        pairs = pairs[-max_items:]

    return chart_type, pairs


def build_chart_detail_lines(chart, limit=4):
    _, pairs = normalize_chart_pairs(chart)
    if not pairs:
        return []
    metric_kind = infer_metric_kind(chart.get("metric_label") or chart.get("titulo"), [value for _, value in pairs])
    total = sum(value for _, value in pairs) or 1
    lines = []
    for label, value in pairs[:limit]:
        share = (value / total) * 100 if total else 0
        lines.append(f"{label}: {share:.0f}% | {format_value(value, metric_kind, compact=False)}")
    return lines


def expand_lines_for_slide(lines, max_len=104, max_items=None):
    expanded = []
    for line in lines or []:
        expanded.extend(split_long_text(line, max_len=max_len))
    expanded = unique_texts(expanded, limit=max_items, min_len=10) if max_items else unique_texts(expanded, min_len=10)
    return expanded


def build_group_totals(rows, label_index, numeric_index):
    grouped = {}
    for row in rows:
        if label_index >= len(row) or numeric_index >= len(row):
            continue
        label = clean_text(row[label_index], 38)
        numeric_value = try_number(row[numeric_index])
        if is_noise_text(label) or numeric_value is None:
            continue
        grouped[label] = grouped.get(label, 0) + numeric_value
    return sorted(grouped.items(), key=lambda item: item[1], reverse=True)


def build_group_counts(rows, label_index):
    grouped = {}
    for row in rows:
        if label_index >= len(row):
            continue
        label = clean_text(row[label_index], 38)
        if is_noise_text(label):
            continue
        grouped[label] = grouped.get(label, 0) + 1
    return sorted(grouped.items(), key=lambda item: item[1], reverse=True)


def build_table_signal_profile(headers, rows):
    if not headers or not rows:
        return {
            "row_count": 0,
            "col_count": len(headers),
            "informative_rows": 0,
            "weak_row_ratio": 1,
            "noise_cell_ratio": 1,
            "generic_header_ratio": 1,
            "dominant_label_ratio": 0,
            "detail_allowed": False,
        }

    row_count = len(rows)
    col_count = len(headers)
    generic_headers = sum(1 for header in headers if is_generic_header(header))
    weak_rows = 0
    noise_cells = 0
    informative_rows = 0
    textual_cells = 0
    numeric_cells = 0
    sentence_cells = 0
    total_chars = 0
    meaningful_cells = 0
    min_meaningful = 1 if col_count <= 2 else 2

    for row in rows:
        meaningful = meaningful_cell_count(row)
        informative_rows += 1 if meaningful >= min_meaningful else 0
        weak_rows += 1 if meaningful <= max(1, min_meaningful - 1) else 0
        for value in row[:col_count]:
            if is_noise_text(value):
                noise_cells += 1
                continue
            meaningful_cells += 1
            total_chars += len(clean_text(value))
            if try_number(value) is not None:
                numeric_cells += 1
            else:
                textual_cells += 1
                if looks_like_sentence(value):
                    sentence_cells += 1

    label_index = choose_table_label_column(headers, rows) if col_count >= 1 else None
    dominant_label_ratio = 0
    if label_index is not None:
        counts = build_group_counts(rows, label_index)
        total_count = sum(count for _, count in counts)
        if counts and total_count:
            dominant_label_ratio = counts[0][1] / total_count

    weak_row_ratio = weak_rows / max(1, row_count)
    noise_cell_ratio = noise_cells / max(1, row_count * col_count)
    generic_header_ratio = generic_headers / max(1, col_count)
    avg_text_len = total_chars / max(1, meaningful_cells)
    numeric_cell_ratio = numeric_cells / max(1, meaningful_cells)
    sentence_cell_ratio = sentence_cells / max(1, meaningful_cells)
    textual_density = textual_cells / max(1, meaningful_cells)
    narrative_score = (
        (1 if numeric_cell_ratio < 0.22 else 0)
        + (1 if sentence_cell_ratio >= 0.18 else 0)
        + (1 if avg_text_len >= 28 else 0)
        + (1 if textual_density >= 0.72 else 0)
    )
    narrative_like = narrative_score >= 3
    operational_like = numeric_cell_ratio >= 0.28 or dominant_label_ratio >= 0.22 or col_count >= 5
    table_kind = "narrative" if narrative_like and not operational_like else "operational"
    # Allow detail even for narrative-leaning tables if they have structured columns
    # Only block when BOTH sentence ratio is very high AND avg text is very long
    hard_narrative = table_kind == "narrative" and sentence_cell_ratio >= 0.28 and avg_text_len >= 45
    detail_allowed = (
        informative_rows >= max(3, math.ceil(row_count * 0.45))
        and weak_row_ratio <= 0.42
        and noise_cell_ratio <= 0.58
        and generic_header_ratio <= 0.6
        and not hard_narrative
    )

    return {
        "row_count": row_count,
        "col_count": col_count,
        "informative_rows": informative_rows,
        "weak_row_ratio": weak_row_ratio,
        "noise_cell_ratio": noise_cell_ratio,
        "generic_header_ratio": generic_header_ratio,
        "dominant_label_ratio": dominant_label_ratio,
        "avg_text_len": avg_text_len,
        "numeric_cell_ratio": numeric_cell_ratio,
        "sentence_cell_ratio": sentence_cell_ratio,
        "textual_density": textual_density,
        "table_kind": table_kind,
        "detail_allowed": detail_allowed,
    }


def has_oversized_table_cells(rows, threshold=120, sample_rows=6):
    for row in (rows or [])[:sample_rows]:
        for value in row:
            if len(clean_text(value)) >= threshold:
                return True
    return False


def resolve_rows_per_table_page(headers, rows):
    profile = build_table_signal_profile(headers, rows)
    col_count = len(headers)
    
    # Very heavy narrative or extremely long cells
    if profile["sentence_cell_ratio"] >= 0.15 or profile["avg_text_len"] >= 45 or has_oversized_table_cells(rows, threshold=130):
        return 3
    # Heavy narrative or long cells
    if profile["sentence_cell_ratio"] >= 0.08 or profile["avg_text_len"] >= 30 or has_oversized_table_cells(rows, threshold=85):
        return 4
    # Moderate text or many columns
    if profile["avg_text_len"] >= 22 or profile["sentence_cell_ratio"] >= 0.04 or col_count >= 7:
        return 6
    # Dense columns
    if col_count >= 6:
        return 8
        
    return ROWS_PER_TABLE_PAGE


def choose_table_focus_column(headers, rows, excluded_indexes=None):
    excluded = set(excluded_indexes or [])
    best_index = None
    best_score = -1
    for index, header in enumerate(headers):
        if index in excluded:
            continue
        values = [clean_text(row[index]) for row in rows if index < len(row)]
        non_empty = [value for value in values if not is_noise_text(value)]
        if len(non_empty) < 2:
            continue
        numeric_ratio = sum(1 for value in non_empty if try_number(value) is not None) / max(1, len(non_empty))
        if numeric_ratio > 0.4:
            continue
        unique_count = len({value.casefold() for value in non_empty})
        if unique_count < 2:
            continue
        score = min(unique_count, 12) + len(non_empty) / 10
        if is_executive_dimension_header(header):
            score += 4
        if is_person_like_header(header) or is_identifier_header(header):
            score -= 6
        if score > best_score:
            best_score = score
            best_index = index
    return best_index


def build_table_dashboard_payload(table_name, headers, rows, source_sheet=None):
    profile = build_table_signal_profile(headers, rows)
    display_title = present_sheet_title(table_name)
    row_count = len(rows)
    col_count = len(headers)
    basis = derive_table_basis(headers, rows)
    truth = evaluate_table_truth(headers, rows, basis=basis, source_sheet=source_sheet)
    label_index = basis["label_index"]
    numeric_index = basis["numeric_index"]
    focus_index = choose_table_focus_column(headers, rows, excluded_indexes={label_index, numeric_index} if numeric_index is not None else {label_index})

    kpis = [
        {"label": "Registros analizados", "value": format_metric_full(row_count)},
        {"label": "Variables", "value": format_metric_full(col_count)},
    ]

    ranking = []
    insights = []

    if numeric_index is not None:
        numeric_header = clean_text(headers[numeric_index], 36) or "Indicador"
        numeric_values = [try_number(row[numeric_index]) for row in rows if numeric_index < len(row)]
        numeric_values = [value for value in numeric_values if value is not None]
        total = sum(numeric_values) if numeric_values else None
        average = (total / len(numeric_values)) if numeric_values else None
        metric_kind = infer_metric_kind(numeric_header, numeric_values)
        kpis.append({"label": numeric_header, "value": format_metric(total, metric_kind)})

        if label_index is not None:
            grouped = build_group_totals(rows, label_index, numeric_index)
            ranking = [{"label": label, "value": value, "display": format_metric(value, metric_kind)} for label, value in grouped[:MAX_DASHBOARD_BARS]]
            if grouped:
                insights.append(f"Mayor concentracion en {grouped[0][0]} con {format_metric(grouped[0][1], metric_kind)}.")
            if average is not None:
                insights.append(f"Promedio del indicador {numeric_header}: {format_metric_full(average, metric_kind)}.")
    elif label_index is not None:
        grouped = build_group_counts(rows, label_index)
        ranking = [{"label": label, "value": value, "display": format_metric_full(value)} for label, value in grouped[:MAX_DASHBOARD_BARS]]
        if grouped:
            if grouped[0][1] > 1 or len(grouped) <= 3:
                insights.append(f"La categoria dominante es {grouped[0][0]} con {grouped[0][1]} registros.")
        kpis.append({"label": "Cobertura", "value": format_metric_full(len(grouped))})

    if focus_index is not None:
        focus_header = clean_text(headers[focus_index], 34)
        focus_counts = build_group_counts(rows, focus_index)
        if focus_counts:
            if focus_counts[0][1] > 1 or len(focus_counts) <= 3:
                insights.append(f"El foco {focus_header} se concentra en {focus_counts[0][0]}.")

    if not ranking and not insights and not profile["detail_allowed"]:
        return None

    return {
        "type": "kpi_dashboard",
        "title": display_title,
        "subtitle": "Indicadores clave",
        "source_basis": build_source_basis_text(basis["dimension_label"], basis["metric_label"], basis["aggregation"]),
        "traceability": build_traceability(source_sheet, basis["aggregation"], truth["score"], row_count=row_count, col_count=col_count, numeric_coverage=truth["numeric_coverage"]),
        "dimension_label": basis["dimension_label"],
        "metric_label": basis["metric_label"],
        "kpis": kpis[:3],
        "ranking": ranking[:MAX_DASHBOARD_BARS],
        "insights": unique_texts([compact_line(item, 88) for item in insights], limit=3, min_len=14),
    }


def should_render_kpi_dashboard(table_name, headers, rows, budget):
    profile = build_table_signal_profile(headers, rows)
    if profile["table_kind"] == "narrative" and profile["numeric_cell_ratio"] < 0.35:
        return False
    basis = derive_table_basis(headers, rows)
    numeric_index = basis["numeric_index"]
    label_index = basis["label_index"]
    numeric_is_good = numeric_index is not None and not is_identifier_header(headers[numeric_index])
    label_is_good = label_index is not None and not is_identifier_header(headers[label_index]) and not is_person_like_header(headers[label_index])
    if numeric_index is not None and label_index is not None and len(rows) >= 8 and profile["informative_rows"] >= 4:
        if numeric_is_good and label_is_good:
            return True
    return (
        len(headers) >= 5
        and len(rows) >= max(6, int(budget.get("genericas", 2) or 2) * 2)
        and profile["informative_rows"] >= 4
        and (numeric_is_good or label_is_good)
    )


def should_render_table_summary(table_name, headers, rows):
    profile = build_table_signal_profile(headers, rows)
    if profile["informative_rows"] < 2:
        return False
    if len(rows) < 4:
        return False
    if profile["table_kind"] == "narrative":
        return True
    if len(rows) >= 5:
        return True
    if len(headers) >= 5:
        return True
    return False


def should_render_table_detail(table_name, headers, rows, budget, has_chart_blocks):
    profile = build_table_signal_profile(headers, rows)
    if not profile["detail_allowed"]:
        return False
    # Only block tables that are overwhelmingly narrative long text
    if profile["avg_text_len"] >= 48 and profile["sentence_cell_ratio"] >= 0.25:
        return False
    # Long cells are truncated at render time — no blocking needed here
    if len(headers) > 6:
        return False
    # Allow tables even when charts exist — charts and tables complement each other
    if len(rows) <= 6:
        return True
    return int(budget.get("genericas", 0) or 0) >= 1 and len(rows) <= 25


def build_slide_2_cards(data, excel_path):
    briefing = build_intro_briefing(data, excel_path)
    objectives = briefing["objectives"]
    topics = briefing["topics"]
    summary = briefing["summary"]
    workplan = briefing["workplan"]
    executive_ai = briefing.get("executive_ai") or {}
    visual_plan = data.get("visual_plan_ia") or {}
    visual_curation_ready = is_visual_curation_ready(data)
    topic = compact_line(infer_excel_topic(data, excel_path), 44)
    focus_phrase = build_prompt_focus_phrase(data, excel_path, max_len=44)
    mode_label = "Boardroom" if IS_BOARDROOM_MODE else "Ejecutivo"
    fallback_topic = f"{focus_phrase}." if focus_phrase else (f"Analisis ejecutivo de {topic}." if topic else "Analisis ejecutivo con datos reales del Excel.")
    fallback_importance = [
        (data.get("resumen_ejecutivo_ia") or {}).get("alerta_principal") or (summary[1] if len(summary) > 1 else ""),
        (data.get("resumen_ejecutivo_ia") or {}).get("recomendacion") or (topics[0] if topics else ""),
    ]
    fallback_chart_topics = [
        compact_line((chart or {}).get("insight_auto"), 42)
        for chart in safe_list(data.get("graficas_automaticas"))[:2]
        if compact_line((chart or {}).get("insight_auto"), 42)
    ] if visual_curation_ready else []
    fallback_planning = [
        objectives[0] if objectives else f"{mode_label} orientado a decisiones y seguimiento.",
        "Se prioriza evidencia util para decidir y actuar, evitando anexos o soportes secundarios.",
    ]
    fallback_topics = safe_list(visual_plan.get("storyline"))[:2] or fallback_chart_topics[:2] or [
        topics[0] if topics else "Se priorizan mensajes con impacto real para la gerencia.",
        topics[1] if len(topics) > 1 else "Se evita informacion tecnica que no ayude a decidir.",
    ]

    ai_topic = clean_text(executive_ai.get("de_que_trata"))
    ai_importance = safe_list(executive_ai.get("breve_resumen")) or safe_list(executive_ai.get("datos_tecnicos"))
    ai_planning = safe_list(executive_ai.get("planeamiento")) or safe_list(executive_ai.get("objetivos"))
    ai_topics = safe_list(executive_ai.get("puntos_a_tratar")) or safe_list(executive_ai.get("elementos_prioritarios"))

    card_1_lines = expand_lines_for_slide([ai_topic or fallback_topic], max_len=42, max_items=3)
    card_2_lines = expand_lines_for_slide(
        ai_importance[:2] or fallback_importance,
        max_len=42,
        max_items=4,
    )
    card_3_lines = expand_lines_for_slide(ai_planning[:2] or fallback_planning, max_len=42, max_items=4)
    card_4_lines = expand_lines_for_slide(ai_topics[:2] or fallback_topics, max_len=42, max_items=4)

    cards = [
        {
            "badge": "01",
            "title": "Contexto",
            "body": "\n".join(card_1_lines),
            "x": Inches(1.2),
            "y": Inches(1.85),
            "w": Inches(3.55),
            "h": Inches(1.46),
        },
        {
            "badge": "02",
            "title": "Por que importa",
            "body": "\n".join(card_2_lines),
            "x": Inches(6.15),
            "y": Inches(1.85),
            "w": Inches(3.55),
            "h": Inches(1.46),
        },
        {
            "badge": "03",
            "title": "Enfoque",
            "body": "\n".join(card_3_lines),
            "x": Inches(1.2),
            "y": Inches(3.92),
            "w": Inches(3.55),
            "h": Inches(1.46),
        },
        {
            "badge": "04",
            "title": "Decisiones",
            "body": "\n".join(card_4_lines),
            "x": Inches(6.15),
            "y": Inches(3.92),
            "w": Inches(3.55),
            "h": Inches(1.46),
        },
    ]
    return cards


def build_slide_3_sections(data, excel_path):
    briefing = build_intro_briefing(data, excel_path)
    summary = briefing["summary"]
    objectives = briefing["objectives"]
    priorities = briefing["topics"][:]
    workplan = briefing["workplan"]
    executive_ai = briefing.get("executive_ai") or {}
    visual_plan = data.get("visual_plan_ia") or {}
    visual_curation_ready = is_visual_curation_ready(data)

    if IS_BOARDROOM_MODE and visual_curation_ready:
        boardroom_blocks = build_boardroom_blocks(data, excel_path)
        for block in boardroom_blocks:
            if block["type"] == "boardroom_traffic":
                priorities.append("Se prioriza un semaforo ejecutivo para enfocar estados criticos, medios y controlados.")
            elif block["type"] == "risk_opportunity_cards":
                priorities.append("Se incluyen tarjetas de riesgo y oportunidad para orientar decisiones de seguimiento.")
            elif block["type"] == "before_after_compare":
                priorities.append("Se incorpora un comparativo antes/despues cuando el Excel evidencia secuencia temporal.")
            elif block["type"] == "top5_alerts":
                priorities.append("Se construye un top 5 con umbrales y alertas automaticas sobre la tabla mas relevante.")
    if visual_curation_ready:
        for chart in safe_list(data.get("graficas_automaticas")):
            insight = compact_line(chart.get("insight_auto"), 84)
            if insight:
                priorities.append(insight)
    priorities.extend(workplan[:2])
    priorities.extend(safe_list(visual_plan.get("storyline"))[:2])
    priorities = unique_texts(priorities, limit=SLIDE3_MAX_BULLETS, min_len=18)
    summary = expand_lines_for_slide(summary[:SLIDE3_MAX_BULLETS], max_len=76, max_items=4)
    objectives = expand_lines_for_slide(unique_texts(objectives, limit=SLIDE3_MAX_BULLETS, min_len=18), max_len=76, max_items=4)
    priorities = expand_lines_for_slide(priorities[:SLIDE3_MAX_BULLETS], max_len=76, max_items=4)

    if executive_ai:
        return [
            ("BREVE RESUMEN", safe_list(executive_ai.get("breve_resumen"))[:SLIDE3_MAX_BULLETS] or summary[:SLIDE3_MAX_BULLETS] or ["Lectura ejecutiva generada con datos reales del archivo."]),
            ("OBJETIVOS", safe_list(executive_ai.get("objetivos"))[:SLIDE3_MAX_BULLETS] or objectives[:SLIDE3_MAX_BULLETS] or ["Sintetizar la informacion critica del Excel para decision gerencial."]),
            ("ELEMENTOS PRIORITARIOS", safe_list(executive_ai.get("elementos_prioritarios"))[:SLIDE3_MAX_BULLETS] or priorities[:SLIDE3_MAX_BULLETS] or ["Se priorizan los frentes con mayor impacto sobre seguimiento y control."]),
        ]

    resumen_ia = data.get("resumen_ejecutivo_ia")
    if resumen_ia and isinstance(resumen_ia, dict):
        return [
            ("VISIÓN GENERAL (AI)", [resumen_ia.get("vision_general", "N/D")]),
            ("ALERTA PRINCIPAL", [resumen_ia.get("alerta_principal", "N/D")]),
            ("RECOMENDACIÓN ESTRATÉGICA", [resumen_ia.get("recomendacion", "N/D")]),
        ]

    return [
        ("BREVE RESUMEN", summary[:SLIDE3_MAX_BULLETS] or ["No se identifico un resumen confiable del Excel analizado."]),
        ("OBJETIVOS", objectives[:SLIDE3_MAX_BULLETS] or ["Sintetizar el Excel con criterio ejecutivo y soporte tecnico real."]),
        ("ELEMENTOS PRIORITARIOS", priorities[:SLIDE3_MAX_BULLETS] or ["Se prioriza la informacion principal del archivo sin agregar contenido no evidenciado."]),
    ]


def build_text_blocks(data):
    blocks = []
    for item in safe_list(data.get("bloques_textuales"))[:4]:
        lines = unique_texts(item.get("lines"), limit=6, min_len=16)
        raw_title = clean_text(item.get("title")) or ""
        normalized_title = raw_title.lower()
        if any(token in normalized_title for token in ("coso", "procedimiento", "oportunidad", "hallazgo", "evidencia", "cuestionario")):
            continue
        if not lines:
            continue
        display_title = present_sheet_title(raw_title) if raw_title else "Analisis complementario"
        raw_subtitle = clean_text(item.get("subtitle"), 80) or ""
        # Replace internal subtitle labels with consultive equivalents
        _subtitle_map = {
            "hallazgos relevantes": "Puntos de atencion criticos",
            "lectura nativa": "Contenido analitico",
            "focos documentales": "Puntos de atencion prioritarios",
            "lectura ejecutiva": "Vision consolidada",
            "lectura documental": "Analisis documental",
            "contenido textual del excel": "Contenido analitico del archivo",
        }
        display_subtitle = _subtitle_map.get(raw_subtitle.lower(), raw_subtitle) or "Contenido analitico del archivo"
        blocks.append({
            "type": "text",
            "origin": "sheet_text",
            "title": display_title,
            "subtitle": display_subtitle,
            "lines": expand_lines_for_slide(lines, max_len=108, max_items=6),
        })

    filtered_conclusions = []
    for line in unique_texts(data.get("conclusiones"), limit=12, min_len=18):
        normalized = clean_text(line).lower()
        if any(
            token in normalized
            for token in (
                "hojas funcionales",
                "auditoria y control",
                "frentes de distribucion",
                "coso",
                "procedimiento",
                "base de datos",
                "columnas con informacion",
                "archivo corresponde",
                "ids",
                "identificador",
            )
        ):
            continue
        filtered_conclusions.append(line)
    detailed = filtered_conclusions[:6]
    if not detailed:
        return blocks

    expanded = expand_lines_for_slide(detailed, max_len=114, max_items=10)
    if blocks:
        expanded = expanded[:5]
    chunk_size = 5
    focus_phrase = build_prompt_focus_phrase(data, data.get("metadatos", {}).get("archivo") or "archivo", max_len=58)
    for index in range(0, len(expanded), chunk_size):
        chunk_num = index // chunk_size + 1
        blocks.append({
            "type": "text",
            "origin": "conclusions",
            "title": compact_line(f"Conclusiones sobre {focus_phrase}", 72) if focus_phrase else "Conclusiones del Analisis",
            "subtitle": f"Mensajes clave · Seccion {chunk_num}" if chunk_num > 1 else "Hallazgos y recomendaciones alineadas al prompt",
            "lines": expanded[index:index + chunk_size],
        })
    return blocks


def build_chart_blocks(data):
    blocks = []
    for chart in safe_list(data.get("graficas_automaticas")):
        if not is_valid_chart(chart):
            continue
        source_sheet = resolve_source_sheet_name(data, chart.get("hoja_origen"), chart.get("titulo"))
        if is_low_signal_sheet_name(source_sheet):
            continue
        truth = evaluate_chart_truth(chart, source_sheet=source_sheet)
        if truth["blocked"]:
            continue
        source_basis = build_source_basis_text(chart.get("dimension_label"), chart.get("metric_label"), chart.get("aggregation"))
        raw_chart_title = clean_text(chart.get("titulo"), 80) or "Grafica"
        # Translate sheet name segment if title has pattern "SheetName · ColumnName"
        if " · " in raw_chart_title:
            sheet_part, _, col_part = raw_chart_title.partition(" · ")
            exec_sheet = present_sheet_title(sheet_part.strip())
            chart_display_title = f"{exec_sheet} · {col_part.strip()}"
        else:
            chart_display_title = present_sheet_title(raw_chart_title) if raw_chart_title else "Grafica"
        blocks.append({
            "type": "chart",
            "title": chart_display_title,
            "chart": chart,
            "insight": clean_text(chart.get("insight_auto"), 160),
            "visual_id": clean_text(chart.get("_visual_ai_id"), 80),
            "ai_rationale": clean_text(chart.get("_visual_ai_rationale"), 160),
            "ai_selected": bool(chart.get("_visual_ai_selected")),
            "source_basis": source_basis,
            "traceability": build_traceability(source_sheet, chart.get("aggregation"), truth["score"], row_count=len(safe_list(chart.get("labels"))), col_count=2, numeric_coverage=truth["numeric_coverage"]),
        })
    existing_titles = {clean_text(block.get("title")).casefold() for block in blocks}
    for chart_block in build_derived_chart_blocks(data):
        if clean_text(chart_block.get("title")).casefold() in existing_titles:
            continue
        blocks.append(chart_block)
        existing_titles.add(clean_text(chart_block.get("title")).casefold())
    return blocks[:MAX_CHART_BLOCKS]


def classify_table_bucket(table_name):
    return "genericas"


def normalize_sheet_family_hint(table):
    return clean_text((table or {}).get("sheet_family")).lower()


def is_low_signal_sheet_name(name):
    normalized = normalize_header_key(name)
    if not normalized:
        return True
    weak_tokens = (
        "hoja",
        "sheet",
        "soporte",
        "anexo",
        "aux",
        "auxiliar",
        "backup",
        "copia",
        "borrador",
        "basetemp",
        "tmp",
    )
    if normalized in {"td", "hoja1", "hoja 1"}:
        return True
    if normalized.startswith("muestra total"):
        return True
    if normalized.startswith("modelo"):
        return True
    if re.fullmatch(r"\d+(?:[ ._-]\d+)*", normalized):
        return True
    if any(normalized.startswith(token) for token in weak_tokens):
        return True
    return False


def present_sheet_title(name):
    if is_low_signal_sheet_name(name):
        return "Hoja Analizada"
    return clean_text(name)


def resolve_sheet_type_from_family(sheet_family, fallback_type):
    mapping = {
        "auditoria": "findings",
        "checklist": "evidence",
        "procedimiento": "evidence",
        "cuestionario": "findings",
        "evidencias": "evidence",
        "arqueo": "evidence",
        "hallazgos": "findings",
        "oportunidades": "findings",
        "matriz_riesgos": "risk",
        "coso": "risk",
    }
    return mapping.get(sheet_family, fallback_type)


def is_documentary_sheet_family(sheet_family):
    # auditoria and checklist have structured tabular data → allow table detail rendering
    return sheet_family in {"procedimiento", "cuestionario", "evidencias", "arqueo"}


def resolve_table_page_limit(table_name, budget, has_chart_blocks):
    bucket = classify_table_bucket(table_name)
    raw_limit = int(budget.get(bucket, 1) or 1)
    # Charts and tables coexist — don't penalize table pages when charts are present
    return max(1, min(raw_limit, 4))


def split_table_columns(headers, rows, max_cols=MAX_TABLE_COLS):
    if len(headers) <= max_cols:
        return [(headers, rows, "")]

    groups = []
    anchor = headers[0]
    other_headers = headers[1:]
    chunk_width = max_cols - 1
    for index in range(0, len(other_headers), chunk_width):
        subset_headers = [anchor] + other_headers[index:index + chunk_width]
        indexes = [headers.index(item) for item in subset_headers]
        subset_rows = [[row[col_index] if col_index < len(row) else "" for col_index in indexes] for row in rows]
        label = f"Columnas {index + 2}-{min(index + chunk_width + 1, len(headers))}"
        groups.append((subset_headers, subset_rows, label))
    return groups


def build_table_blocks(data):
    budget = data.get("presupuesto_slides") or {}
    has_chart_blocks = len(build_chart_blocks(data)) > 0
    # Use available slots budget — allow more blocks when there are many sheets
    slots_available = int(budget.get("slots_restantes", 8) or 8)
    base_budget = int(budget.get("tabla_principal", 2) or 2) + int(budget.get("genericas", 3) or 3) + int(budget.get("hallazgos", 2) or 2) + int(budget.get("coso", 1) or 1)
    max_table_blocks = max(6, min(16, base_budget + max(0, (slots_available - 4) // 2)))
    dashboard_count = 0
    raw_candidates = collect_table_candidates(data)
    # Prioritize tables with real data (detail candidates) before narrative/text-only tables
    def _candidate_priority(candidate):
        name, table, source_sheet = candidate
        h, r = extract_table_payload(table)
        if not h or not r:
            return 2
        sf = table.get("sheet_family", "")
        if is_documentary_sheet_family(sf):
            return 2
        profile = build_table_signal_profile(h, r)
        if should_render_table_detail(name, h, r, budget, has_chart_blocks):
            return 0  # real table data first
        if profile["table_kind"] != "narrative":
            return 1  # operational but no detail (KPI dashboards)
        return 2  # narrative/text-only last
    ordered_tables = sorted(
        raw_candidates,
        key=lambda candidate: (
            _candidate_priority(candidate),
            -score_table_candidate(candidate[0], candidate[1], candidate[2]),
        ),
    )

    blocks = []
    for table_name, table, source_sheet in ordered_tables:
        display_title = present_sheet_title(table_name)
        headers, rows = extract_table_payload(table)
        if len(headers) < 2 or not rows:
            continue
        if is_low_signal_sheet_name(table_name) or is_low_signal_sheet_name(source_sheet):
            continue
        sheet_family = normalize_sheet_family_hint(table)
        profile = build_table_signal_profile(headers, rows)
        basis = derive_table_basis(headers, rows)
        truth = evaluate_table_truth(headers, rows, basis=basis, source_sheet=source_sheet)
        sheet_type = resolve_sheet_type_from_family(sheet_family, classify_sheet_type(table_name, headers, rows))
        native_lines = build_native_sheet_highlights(sheet_type, table_name, headers, rows)
        detail_candidate = should_render_table_detail(table_name, headers, rows, budget, has_chart_blocks) and not truth["blocked_detail"]
        has_sensitive_headers = any(is_identifier_header(header) or is_person_like_header(header) for header in headers)
        if has_sensitive_headers:
            detail_candidate = False
        if is_documentary_sheet_family(sheet_family):
            detail_candidate = False
        allow_dashboard = dashboard_count < MAX_KPI_DASHBOARDS and not is_documentary_sheet_family(sheet_family)
        dashboard_candidate = None
        if allow_dashboard and should_render_kpi_dashboard(table_name, headers, rows, budget) and not truth["blocked_dashboard"]:
            dashboard_candidate = build_table_dashboard_payload(table_name, headers, rows, source_sheet=source_sheet)

        # For purely narrative/documentary tables (no detail), limit to 1 block to save slide budget
        is_text_only = not dashboard_candidate and not detail_candidate
        if native_lines and (profile["table_kind"] == "narrative" or is_text_only):
            _native_subtitles = {
                "findings": "Hallazgos identificados",
                "evidence": "Soporte documental",
                "risk": "Gestion de riesgos",
                "operational": "Analisis operacional",
            }
            native_subtitle = _native_subtitles.get(sheet_type, "Contenido analitico")
            blocks.append({
                "type": "text",
                "title": display_title,
                "subtitle": native_subtitle,
                "lines": expand_lines_for_slide(native_lines, max_len=92, max_items=5),
                "visual_id": clean_text(table.get("_visual_ai_id"), 80),
                "ai_rationale": clean_text(table.get("_visual_ai_rationale"), 160),
                "ai_selected": bool(table.get("_visual_ai_selected")),
            })
            if len(blocks) >= max_table_blocks:
                return blocks
            # For documentary-only sheets, stop here — don't add summary + focos too
            if is_text_only and is_documentary_sheet_family(sheet_family):
                continue

        kpi_rendered = False
        if dashboard_candidate:
            if table.get("_visual_ai_id"):
                dashboard_candidate["visual_id"] = clean_text(table.get("_visual_ai_id"), 80)
                dashboard_candidate["ai_rationale"] = clean_text(table.get("_visual_ai_rationale"), 160)
                dashboard_candidate["ai_selected"] = bool(table.get("_visual_ai_selected"))
            blocks.append(dashboard_candidate)
            kpi_rendered = True
            dashboard_count += 1
            if len(blocks) >= max_table_blocks:
                return blocks

        if should_render_table_summary(table_name, headers, rows) and not kpi_rendered and not truth["blocked_summary"]:
            _summary_subtitle = "Analisis documental" if profile["table_kind"] == "narrative" else "Vision consolidada"
            summary_lines = build_table_summary_lines(table_name, headers, rows)
            ai_message = clean_text(table.get("_visual_ai_message"), 160)
            if ai_message:
                summary_lines = unique_texts([ai_message] + summary_lines, limit=5, min_len=18)
            blocks.append({
                "type": "table_summary",
                "title": display_title,
                "subtitle": _summary_subtitle,
                "lines": summary_lines,
                "visual_id": clean_text(table.get("_visual_ai_id"), 80),
                "ai_rationale": clean_text(table.get("_visual_ai_rationale"), 160),
                "ai_selected": bool(table.get("_visual_ai_selected")),
                "source_basis": build_source_basis_text(basis["dimension_label"], basis["metric_label"], basis["aggregation"]),
                "traceability": build_traceability(source_sheet, basis["aggregation"], truth["score"], row_count=len(rows), col_count=len(headers), numeric_coverage=truth["numeric_coverage"]),
            })
            if len(blocks) >= max_table_blocks:
                return blocks

        # Narrative focus lines only when not hitting budget and table has real narrative insights
        if profile["table_kind"] == "narrative" and not detail_candidate:
            narrative_lines = build_narrative_focus_lines(table_name, headers, rows)
            if narrative_lines:
                blocks.append({
                    "type": "text",
                    "title": display_title,
                    "subtitle": "Puntos de atencion prioritarios",
                    "lines": expand_lines_for_slide(narrative_lines, max_len=92, max_items=6),
                    "visual_id": clean_text(table.get("_visual_ai_id"), 80),
                    "ai_rationale": clean_text(table.get("_visual_ai_rationale"), 160),
                    "ai_selected": bool(table.get("_visual_ai_selected")),
                })
                if len(blocks) >= max_table_blocks:
                    return blocks

        if not detail_candidate:
            continue

        for subset_headers, subset_rows, label in split_table_columns(headers, rows):
            subset_headers, subset_rows = densify_table_matrix(subset_headers, subset_rows, min_rows=2)
            if len(subset_headers) < 2 or len(subset_rows) < 2:
                continue
            max_pages_for_table = resolve_table_page_limit(table_name, budget, has_chart_blocks)
            rows_per_page = resolve_rows_per_table_page(subset_headers, subset_rows)
            total_pages = max(1, math.ceil(len(subset_rows) / rows_per_page))
            if total_pages > 6:
                continue
            for page_index in range(min(total_pages, max_pages_for_table)):
                chunk = subset_rows[page_index * rows_per_page : (page_index + 1) * rows_per_page]
                if not chunk:
                    continue
                page_headers, chunk = densify_table_matrix(subset_headers, chunk, min_rows=2)
                if len(page_headers) < 2 or len(chunk) < 2:
                    continue
                page_label = f"Pagina {page_index + 1}/{total_pages}"
                suffix = f"{label} · {page_label}" if label else page_label
                insight_tabla = data.get("resumen_ejecutivo_ia", {}).get("insight_tabla")
                if insight_tabla and page_index == 0:
                    suffix = f"{insight_tabla} | {suffix}"
                table_basis = derive_table_basis(page_headers, chunk)
                page_truth = evaluate_table_truth(page_headers, chunk, basis=table_basis, source_sheet=source_sheet)
                if page_truth["blocked_detail"]:
                    continue
                blocks.append({
                    "type": "table",
                    "title": display_title,
                    "subtitle": suffix,
                    "headers": page_headers,
                    "rows": chunk,
                    "visual_id": clean_text(table.get("_visual_ai_id"), 80),
                    "ai_rationale": clean_text(table.get("_visual_ai_rationale"), 160),
                    "ai_selected": bool(table.get("_visual_ai_selected")),
                    "source_basis": build_source_basis_text(table_basis["dimension_label"], table_basis["metric_label"], table_basis["aggregation"]),
                    "traceability": build_traceability(source_sheet, table_basis["aggregation"], page_truth["score"], row_count=len(chunk), col_count=len(page_headers), numeric_coverage=page_truth["numeric_coverage"]),
                })
                if len(blocks) >= max_table_blocks:
                    return blocks
    return blocks


def build_visual_plan_maps(data):
    plan = data.get("visual_plan_ia") or {}
    chart_map = {}
    for order, item in enumerate(safe_list(plan.get("charts"))):
        if not isinstance(item, dict):
            continue
        visual_id = clean_text(item.get("id"), 80)
        if not visual_id:
            continue
        chart_map[visual_id] = {
            "order": order,
            "message": clean_text(item.get("mensaje_clave"), 160),
        }

    table_map = {}
    for order, item in enumerate(safe_list(plan.get("tables"))):
        if not isinstance(item, dict):
            continue
        visual_id = clean_text(item.get("id"), 80)
        mode = clean_text(item.get("modo"), 20).lower()
        if not visual_id:
            continue
        table_map[visual_id] = {
            "order": order,
            "mode": mode,
            "message": clean_text(item.get("mensaje_clave"), 160),
        }
    return chart_map, table_map


def score_block_against_visual_plan(block, chart_map, table_map):
    visual_id = clean_text(block.get("visual_id"), 80)
    if not visual_id:
        return 0
    block_type = clean_text(block.get("type")).lower()
    score = 0
    if block_type == "chart" and visual_id in chart_map:
        score += 40 - chart_map[visual_id]["order"] * 4
    elif visual_id in table_map:
        mode = table_map[visual_id]["mode"]
        score += 30 - table_map[visual_id]["order"] * 3
        if mode == "summary" and block_type == "table_summary":
            score += 12
        elif mode == "detail" and block_type == "table":
            score += 12
        elif mode == "omit":
            score -= 20
        elif block_type in {"text", "kpi_dashboard"}:
            score += 2
    if block.get("ai_selected"):
        score += 6
    return score


def score_block_against_prompt(block, preferences):
    if not preferences:
        return 0
    haystack = " ".join(
        clean_text(item)
        for item in (
            block.get("title"),
            block.get("subtitle"),
            block.get("source_basis"),
            " ".join(safe_list(block.get("lines"))),
        )
        if clean_text(item)
    ).lower()
    score = 0
    for token in preferences.get("priority_tokens") or []:
        if token == "hallazgos" and any(keyword in haystack for keyword in ("hallazgo", "auditoria", "riesgo")):
            score += 7
        elif token == "costos" and any(keyword in haystack for keyword in ("costo", "valor", "monto", "gasto", "centro de costo")):
            score += 6
        elif token == "estado" and "estado" in haystack:
            score += 5
        elif token == "coso" and any(keyword in haystack for keyword in ("coso", "control interno", "control")):
            score += 6
        elif token == "riesgos" and any(keyword in haystack for keyword in ("riesgo", "alerta", "anomalia", "anomal", "outlier")):
            score += 5
    for term in preferences.get("focus_terms") or []:
        if len(term) < 4:
            continue
        if term in haystack:
            score += 4
    if preferences.get("prefer_short_tables") and block.get("type") == "table":
        row_count = len(block.get("rows") or [])
        if row_count <= 8:
            score += 4
        elif row_count <= 12:
            score += 2
        else:
            score -= 3
    if block.get("type") == "kpi_dashboard":
        score += 2
    return score


def order_blocks_by_prompt(blocks, preferences):
    if not preferences:
        return list(blocks)
    return sorted(
        blocks,
        key=lambda block: score_block_against_prompt(block, preferences),
        reverse=True,
    )


def enforce_chart_preferences(chart_blocks, preferences):
    adjusted = []
    for block in chart_blocks:
        current = dict(block)
        if preferences and preferences.get("prefer_bar_charts"):
            chart = dict(current.get("chart") or {})
            chart_type = clean_text(chart.get("tipo")).lower()
            if chart_type in {"pie", "doughnut"}:
                chart["tipo"] = "bar"
                current["chart"] = chart
        adjusted.append(current)
    return adjusted


def append_unique_blocks(target, blocks, seen_signatures, limit=None):
    for block in blocks:
        signature = (
            block.get("type"),
            clean_text(block.get("title"), 80).casefold(),
            clean_text(block.get("subtitle"), 80).casefold(),
        )
        if signature in seen_signatures:
            continue
        seen_signatures.add(signature)
        target.append(block)
        if limit and len(target) >= limit:
            break


def apply_design_slide_to_block(block, design_slide):
    current = dict(block)
    if not isinstance(design_slide, dict):
        return current
    title = clean_text(design_slide.get("title"), 120)
    notes = clean_text(design_slide.get("design_notes"), 180)
    if title:
        current["title"] = title
    if notes:
        current["subtitle"] = notes
    if current.get("type") == "chart":
        insight_parts = [
            clean_text(design_slide.get("insight_headline"), 80),
            clean_text(design_slide.get("insight_body"), 220),
        ]
        insight_text = " ".join(part for part in insight_parts if part).strip()
        if insight_text:
            current["insight"] = insight_text
    elif current.get("type") in {"table", "table_summary", "kpi_dashboard"}:
        footnote = clean_text(design_slide.get("footnote"), 160)
        insight_text = clean_text(design_slide.get("insight_text"), 200)
        if footnote:
            current["subtitle"] = footnote
        elif insight_text:
            current["subtitle"] = insight_text
    elif current.get("type") == "text" and design_slide.get("summary_bullets"):
        bullets = safe_list(design_slide.get("summary_bullets"))
        if bullets:
            current["lines"] = [clean_text(item, 140) for item in bullets[:5] if clean_text(item, 140)]
    cta_text = clean_text(design_slide.get("cta_text"), 160)
    if cta_text and current.get("type") == "text":
        lines = [clean_text(item, 140) for item in safe_list(current.get("lines")) if clean_text(item, 140)]
        if cta_text not in lines:
            lines.append(cta_text)
        current["lines"] = lines[:5]
    current["design_notes"] = notes
    return current


def build_prompt_fallback_design_plan(data, prompt_preferences, excel_path):
    request_context = get_request_context(data)
    raw_prompt = clean_text((request_context or {}).get("prompt"), 180)
    if not raw_prompt:
        return {}

    focus_phrase = build_prompt_focus_phrase(data, excel_path, max_len=58)
    theme = get_theme_context(data)
    prompt_tokens = prompt_preferences.get("priority_tokens") or []
    focus_terms = prompt_preferences.get("focus_terms") or []
    narrative_lines = unique_texts(data.get("conclusiones"), limit=6, min_len=18)
    summary_lines = build_management_highlights(data, excel_path, limit=4)
    storyline = safe_list((data.get("visual_plan_ia") or {}).get("storyline"))[:3]
    chart_requested = (prompt_preferences.get("chart_slides") or 0) > 0 or any(token in prompt_tokens for token in ("costos", "estado"))
    table_requested = (prompt_preferences.get("table_slides") or 0) > 0 or "hallazgos" in prompt_tokens
    visual_ready = is_visual_curation_ready(data)
    chart_available = visual_ready and bool(safe_list(data.get("graficas_automaticas")))
    table_available = visual_ready and bool(data.get("muestra_tabla") or data.get("otras_tablas") or data.get("genericas"))

    focus_label = focus_terms[0] if focus_terms else (prompt_tokens[0] if prompt_tokens else compact_line(infer_excel_topic(data, excel_path), 40))
    if "hallazgos" in prompt_tokens or "riesgos" in prompt_tokens:
        if focus_label in {"hallazgos", "riesgos"}:
            mid_title = "Riesgos y hallazgos clave"
            close_title = "Acciones de control prioritarias"
        else:
            mid_title = f"Riesgos y hallazgos sobre {focus_label}"
            close_title = f"Acciones de control para {focus_label}"
    elif "costos" in prompt_tokens or "estado" in prompt_tokens:
        if focus_label in {"costos", "estado"}:
            mid_title = "Concentraciones y montos clave"
            close_title = "Decisiones financieras prioritarias"
        else:
            mid_title = f"Concentraciones y montos en {focus_label}"
            close_title = f"Decisiones financieras sobre {focus_label}"
    else:
        mid_title = f"Lectura priorizada de {focus_label}"
        close_title = f"Decision siguiente sobre {focus_label}"

    summary_bullets = storyline or narrative_lines[:3] or summary_lines[:3] or [focus_phrase]
    slides = [
        {
            "slide_number": 1,
            "type": "title_slide",
            "title": compact_line(f"Decision ejecutiva: {focus_phrase}", 110),
            "title_text": compact_line(f"Decision ejecutiva: {focus_phrase}", 110),
            "subtitle_text": compact_line(f"Analisis ejecutivo de {basename_label(excel_path)}", 130),
            "design_notes": "La apertura se ancla al prompt del usuario para evitar una portada generica.",
        },
        {
            "slide_number": 2,
            "type": "text_bullets",
            "title": "Contexto y foco ejecutivo",
            "design_notes": "Esta slide traduce el prompt a una lectura ejecutiva concreta antes de mostrar evidencia.",
            "summary_bullets": (summary_lines or narrative_lines or [focus_phrase])[:4],
        },
        {
            "slide_number": 3,
            "type": "text_bullets",
            "title": compact_line(mid_title, 110),
            "design_notes": "La narrativa intermedia cambia segun los terminos clave del prompt y los hallazgos del Excel.",
            "summary_bullets": (narrative_lines or summary_lines or storyline or [focus_phrase])[:4],
        },
    ]

    if chart_requested and chart_available:
        slides.append({
            "slide_number": len(slides) + 1,
            "type": "chart_with_insight",
            "title": compact_line(f"Visual clave para {focus_phrase}", 110),
            "insight_headline": compact_line(summary_bullets[0], 60),
            "insight_body": compact_line((summary_lines or narrative_lines or [focus_phrase])[0], 220),
            "design_notes": "Se reserva una visual solo si el prompt la pide y la curaduria visual esta realmente disponible.",
        })

    if table_requested and table_available:
        slides.append({
            "slide_number": len(slides) + 1,
            "type": "data_table",
            "title": compact_line(f"Detalle puntual para {focus_phrase}", 110),
            "footnote": "Tabla incluida solo cuando el prompt exige detalle y existe soporte valido en el Excel.",
            "design_notes": "La tabla deja de ser automatica y se usa solo como soporte del foco pedido.",
        })

    slides.append({
        "slide_number": len(slides) + 1,
        "type": "closing_slide",
        "title": compact_line(close_title, 110),
        "cta_text": compact_line(f"Profundizar la decision sobre {focus_phrase} con base en la evidencia priorizada.", 150),
        "summary_bullets": summary_bullets[:4],
        "design_notes": "El cierre retoma el pedido del usuario y lo convierte en accion concreta.",
    })

    return {
        "source": "prompt_fallback",
        "presentation_meta": {
            "title": compact_line(f"Decision ejecutiva: {focus_phrase}", 110),
            "subtitle": compact_line(f"Analisis ejecutivo de {basename_label(excel_path)}", 150),
            "author": "Generado con Excel -> PowerPoint Pro",
            "date": clean_text(request_context.get("current_date"), 24),
            "total_slides": len(slides),
            "narrative_summary": compact_line(f"La presentacion prioriza {focus_phrase} con evidencia real del Excel.", 220),
        },
        "slides": slides[:12],
        "global_design": {
            "font_primary": "Calibri",
            "font_secondary": "Calibri Light",
            "footer_text": clean_text(f"{theme.get('name', 'Analitica Moderna')} - {request_context.get('current_date', '')}", 100),
        },
        "validation": {
            "total_slides_matches_array": True,
            "all_slides_have_required_fields": True,
            "no_slide_exceeds_5_bullets": True,
            "no_chart_exceeds_8_series": True,
            "json_is_valid": True,
        },
    }


def build_blocks_from_design_plan(design_plan, pools):
    slides = safe_list((design_plan or {}).get("slides"))
    if not slides:
        return []
    plan_source = clean_text((design_plan or {}).get("source"), 40)

    queue_map = {
        "chart": list(pools.get("chart", [])),
        "kpi_dashboard": list(pools.get("kpi_dashboard", [])),
        "table_summary": list(pools.get("table_summary", [])),
        "table": list(pools.get("table", [])),
        "text": list(pools.get("text", [])),
    }
    selected = []
    seen = set()

    def pick(queue_name):
        queue = queue_map.get(queue_name) or []
        while queue:
            candidate = queue.pop(0)
            signature = (
                candidate.get("type"),
                clean_text(candidate.get("title"), 90).casefold(),
                clean_text(candidate.get("subtitle"), 90).casefold(),
            )
            if signature in seen:
                continue
            seen.add(signature)
            return candidate
        return None

    def build_text_candidate(slide):
        slide_type = clean_text(slide.get("type"), 40)
        title = clean_text(slide.get("title"), 120) or "Narrativa ejecutiva"
        notes = clean_text(slide.get("design_notes"), 180)
        insight_text = clean_text(slide.get("insight_text"), 220)
        insight_headline = clean_text(slide.get("insight_headline"), 80)
        insight_body = clean_text(slide.get("insight_body"), 220)
        cta_text = clean_text(slide.get("cta_text"), 160)
        summary_bullets = [clean_text(item, 140) for item in safe_list(slide.get("summary_bullets")) if clean_text(item, 140)]

        lines = []
        if slide_type == "quote_highlight":
            lines = [insight_headline or title]
            if insight_body:
                lines.append(insight_body)
        elif slide_type == "section_divider":
            lines = [title]
            if notes:
                lines.append(notes)
        elif slide_type == "closing_slide":
            lines = summary_bullets[:4]
            if cta_text:
                lines.append(cta_text)
        else:
            lines = summary_bullets[:5]
            if insight_text:
                lines.insert(0, insight_text)
            if insight_headline and insight_headline not in lines:
                lines.insert(0, insight_headline)
            if insight_body and insight_body not in lines:
                lines.append(insight_body)
            if notes and notes not in lines:
                lines.append(notes)

        lines = [clean_text(item, 140) for item in lines if clean_text(item, 140)]
        if not lines:
            lines = [title]

        candidate = {
            "type": "text",
            "origin": "ai_design",
            "design_slide_type": slide_type,
            "title": title,
            "subtitle": notes or (
                "Narrativa guiada por prompt"
                if plan_source == "prompt_fallback" and slide_type != "closing_slide"
                else "Cierre guiado por prompt"
                if plan_source == "prompt_fallback"
                else "Narrativa guiada por Hermes"
                if slide_type != "closing_slide"
                else "Cierre guiado por Hermes"
            ),
            "lines": lines[:5],
            "design_notes": notes,
        }
        if cta_text:
            candidate["cta_text"] = cta_text
        return candidate

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_type = clean_text(slide.get("type"), 40)
        if slide_type in {"title_slide", "closing_slide"}:
            continue
        candidate = None
        if slide_type in {"chart_full", "chart_with_insight"}:
            candidate = pick("chart")
        elif slide_type == "executive_summary":
            candidate = pick("kpi_dashboard") or pick("text")
        elif slide_type == "data_table":
            candidate = pick("table") or pick("table_summary")
        elif slide_type in {"text_bullets", "two_columns"}:
            candidate = pick("text") or build_text_candidate(slide)
        elif slide_type in {"quote_highlight", "section_divider"}:
            candidate = build_text_candidate(slide)
        if candidate:
            selected.append(apply_design_slide_to_block(candidate, slide))

    return selected


def build_content_blocks(data, prompt_preferences=None):
    blocks = []
    design_plan = get_presentation_design_plan(data)
    design_slides = safe_list((design_plan or {}).get("slides"))
    chart_blocks = enforce_chart_preferences(build_chart_blocks(data), prompt_preferences)
    text_blocks = build_text_blocks(data)
    boardroom_blocks = build_boardroom_blocks(data, data.get("metadatos", {}).get("archivo") or "archivo")
    table_blocks = build_table_blocks(data)
    chart_plan, table_plan = build_visual_plan_maps(data)
    ai_status = get_ai_curation_status(data)
    visual_curation_ready = is_visual_curation_ready(data)
    if len(chart_blocks) >= 3 and len(text_blocks) > 1:
        sheet_text_blocks = [block for block in text_blocks if block.get("origin") == "sheet_text"]
        conclusion_blocks = [block for block in text_blocks if block.get("origin") != "sheet_text"]
        if sheet_text_blocks:
            text_blocks = sheet_text_blocks + conclusion_blocks[:1]
        else:
            text_blocks = text_blocks[:1]

    kpi_blocks = [block for block in table_blocks if block.get("type") == "kpi_dashboard"]
    summary_blocks = [block for block in table_blocks if block.get("type") == "table_summary"]
    detail_blocks = [block for block in table_blocks if block.get("type") == "table"]
    narrative_blocks = [block for block in text_blocks if block.get("origin") != "sheet_text"]
    sheet_text_blocks = [block for block in text_blocks if block.get("origin") == "sheet_text"]
    preferred_chart_blocks = list(chart_blocks)
    preferred_kpi_blocks = list(kpi_blocks)
    preferred_summary_blocks = list(summary_blocks)
    preferred_detail_blocks = list(detail_blocks)

    if STRICT_VISUAL_CURATION_MODE:
        if visual_curation_ready:
            selected_chart_ids = set(ai_status.get("selected_chart_ids") or [])
            selected_table_ids = set(ai_status.get("selected_table_ids") or [])
            chart_blocks = [block for block in chart_blocks if clean_text(block.get("visual_id"), 80) in selected_chart_ids]
            kpi_blocks = [block for block in kpi_blocks if clean_text(block.get("visual_id"), 80) in selected_table_ids]
            summary_blocks = [block for block in summary_blocks if clean_text(block.get("visual_id"), 80) in selected_table_ids]
            detail_blocks = [block for block in detail_blocks if clean_text(block.get("visual_id"), 80) in selected_table_ids]
            table_blocks = kpi_blocks + summary_blocks + detail_blocks
            preferred_chart_blocks = list(chart_blocks)
            preferred_kpi_blocks = list(kpi_blocks)
            preferred_summary_blocks = list(summary_blocks)
            preferred_detail_blocks = list(detail_blocks)
        else:
            chart_blocks = []
            kpi_blocks = []
            summary_blocks = []
            detail_blocks = []
            table_blocks = []
            preferred_chart_blocks = []
            preferred_kpi_blocks = []
            preferred_summary_blocks = []
            preferred_detail_blocks = []
            boardroom_blocks = [block for block in boardroom_blocks if block.get("type") == "text"]

    if chart_plan:
        chart_blocks = sorted(
            chart_blocks,
            key=lambda block: score_block_against_visual_plan(block, chart_plan, table_plan),
            reverse=True,
        )
        selected_chart_blocks = [block for block in chart_blocks if clean_text(block.get("visual_id"), 80) in chart_plan]
        if selected_chart_blocks:
            chart_blocks = selected_chart_blocks + [block for block in chart_blocks if block not in selected_chart_blocks]
            preferred_chart_blocks = selected_chart_blocks

    if table_plan:
        ordered_table_blocks = sorted(
            table_blocks,
            key=lambda block: score_block_against_visual_plan(block, chart_plan, table_plan),
            reverse=True,
        )
        kpi_blocks = [block for block in ordered_table_blocks if block.get("type") == "kpi_dashboard"]
        summary_blocks = [block for block in ordered_table_blocks if block.get("type") == "table_summary"]
        detail_blocks = [block for block in ordered_table_blocks if block.get("type") == "table"]
        preferred_kpi_blocks = [
            block for block in kpi_blocks
            if clean_text(block.get("visual_id"), 80) in table_plan
            and table_plan[clean_text(block.get("visual_id"), 80)]["mode"] != "omit"
        ] or kpi_blocks
        preferred_summary_blocks = [
            block for block in summary_blocks
            if clean_text(block.get("visual_id"), 80) in table_plan
            and table_plan[clean_text(block.get("visual_id"), 80)]["mode"] == "summary"
        ] or summary_blocks
        preferred_detail_blocks = [
            block for block in detail_blocks
            if clean_text(block.get("visual_id"), 80) in table_plan
            and table_plan[clean_text(block.get("visual_id"), 80)]["mode"] == "detail"
        ] or detail_blocks

    use_strict_design_plan = bool(design_slides)

    if use_strict_design_plan:
        design_guided_blocks = build_blocks_from_design_plan(
            design_plan,
            {
                "chart": preferred_chart_blocks or chart_blocks,
                "kpi_dashboard": preferred_kpi_blocks or kpi_blocks,
                "table_summary": preferred_summary_blocks or summary_blocks,
                "table": preferred_detail_blocks or detail_blocks,
                "text": narrative_blocks + sheet_text_blocks,
            },
        )
        if design_guided_blocks:
            return design_guided_blocks[: min(MAX_CONTENT_BLOCKS, len(design_slides))]
    elif prompt_preferences and prompt_preferences.get("explicit_structure"):
        boardroom_blocks = order_blocks_by_prompt(boardroom_blocks, prompt_preferences)
        kpi_blocks = order_blocks_by_prompt(kpi_blocks, prompt_preferences)
        summary_blocks = order_blocks_by_prompt(summary_blocks, prompt_preferences)
        detail_blocks = order_blocks_by_prompt(detail_blocks, prompt_preferences)
        narrative_blocks = order_blocks_by_prompt(narrative_blocks, prompt_preferences)
        sheet_text_blocks = order_blocks_by_prompt(sheet_text_blocks, prompt_preferences)
        chart_blocks = order_blocks_by_prompt(chart_blocks, prompt_preferences)

        seen_signatures = set()
        append_unique_blocks(blocks, boardroom_blocks[: prompt_preferences.get("max_boardroom_blocks", 2)], seen_signatures)
        append_unique_blocks(blocks, kpi_blocks[: prompt_preferences.get("kpi_slides", 0)], seen_signatures)
        append_unique_blocks(blocks, chart_blocks[: prompt_preferences.get("chart_slides", 0)], seen_signatures)

        table_goal = prompt_preferences.get("table_slides", 0)
        preferred_tables = summary_blocks if prompt_preferences.get("prefer_short_tables") else (detail_blocks + summary_blocks)
        if prompt_preferences.get("prefer_short_tables"):
            preferred_tables = summary_blocks + detail_blocks
        append_unique_blocks(blocks, preferred_tables[:table_goal], seen_signatures)

        append_unique_blocks(blocks, narrative_blocks[: prompt_preferences.get("conclusion_slides", 0)], seen_signatures)

        fallback_pool = (
            chart_blocks
            + kpi_blocks
            + summary_blocks
            + detail_blocks
            + narrative_blocks
            + sheet_text_blocks
        )
        append_unique_blocks(blocks, fallback_pool, seen_signatures, limit=prompt_preferences.get("max_total_content_blocks", MAX_CONTENT_BLOCKS))
    else:
        lead_chart_pool = preferred_chart_blocks or chart_blocks
        lead_kpi_pool = preferred_kpi_blocks or kpi_blocks
        lead_summary_pool = preferred_summary_blocks or summary_blocks
        lead_detail_pool = preferred_detail_blocks or detail_blocks
        primary_sheet_name = clean_text((data.get("resumen_generico") or {}).get("hoja_principal"), 80)
        primary_kpi_pool = [
            block for block in lead_kpi_pool
            if clean_text((block.get("traceability") or {}).get("source_sheet"), 80) == primary_sheet_name
        ] or lead_kpi_pool
        if lead_chart_pool:
            blocks.append(lead_chart_pool[0])
        elif primary_kpi_pool:
            blocks.append(primary_kpi_pool[0])
        elif lead_summary_pool:
            blocks.append(lead_summary_pool[0])

        if visual_curation_ready:
            if IS_BOARDROOM_MODE or boardroom_blocks:
                blocks.extend(boardroom_blocks[:2])
            blocks.extend(lead_chart_pool[1:2])
            blocks.extend(lead_kpi_pool[:2])
            blocks.extend(lead_summary_pool[:1])
            blocks.extend(narrative_blocks[:1])
            blocks.extend(sheet_text_blocks[:1])
            blocks.extend(lead_detail_pool[:1])
        else:
            boardroom_text_blocks = [block for block in boardroom_blocks if block.get("type") == "text"]
            if boardroom_text_blocks:
                blocks.extend(boardroom_text_blocks[:1])
            if lead_chart_pool:
                blocks.extend(lead_chart_pool[:1])
            if primary_kpi_pool:
                blocks.extend(primary_kpi_pool[:1])
            elif not lead_chart_pool:
                blocks.extend(lead_summary_pool[:1])
            blocks.extend(narrative_blocks[:1])
            if not narrative_blocks:
                blocks.extend(sheet_text_blocks[:1])

        deduped_blocks = []
        seen_signatures = set()
        for block in blocks:
            signature = (
                block.get("type"),
                clean_text(block.get("title"), 80).casefold(),
                clean_text(block.get("subtitle"), 80).casefold(),
            )
            if signature in seen_signatures:
                continue
            seen_signatures.add(signature)
            deduped_blocks.append(block)
        blocks = deduped_blocks

    if design_plan and not use_strict_design_plan:
        design_guided_blocks = build_blocks_from_design_plan(
            design_plan,
            {
                "chart": preferred_chart_blocks or chart_blocks,
                "kpi_dashboard": preferred_kpi_blocks or kpi_blocks,
                "table_summary": preferred_summary_blocks or summary_blocks,
                "table": preferred_detail_blocks or detail_blocks,
                "text": narrative_blocks + sheet_text_blocks,
            },
        )
        if design_guided_blocks:
            seen_signatures = set()
            final_blocks = []
            append_unique_blocks(final_blocks, design_guided_blocks, seen_signatures)
            append_unique_blocks(final_blocks, blocks, seen_signatures, limit=prompt_preferences.get("max_total_content_blocks", MAX_CONTENT_BLOCKS) if prompt_preferences else MAX_CONTENT_BLOCKS)
            blocks = final_blocks

    if not blocks:
        blocks.append({
            "type": "text",
            "title": "Contenido principal",
            "subtitle": "Curaduria visual",
            "lines": [
                "No se renderizan graficas ni tablas ejecutivas sin curaduria visual valida de Hermes.",
                "La presentacion conserva solo contenido textual confiable hasta recibir seleccion explicita de la IA.",
            ],
        })
    max_blocks = prompt_preferences.get("max_total_content_blocks", MAX_CONTENT_BLOCKS) if prompt_preferences else MAX_CONTENT_BLOCKS
    return blocks[:max_blocks]


def summarize_effective_slide_budget(data, content_blocks):
    budget = dict(data.get("presupuesto_slides") or {})
    effective_chart_count = len([item for item in content_blocks if item.get("type") == "chart"])
    effective_table_count = len([item for item in content_blocks if item.get("type") == "table"])
    effective_text_count = len([item for item in content_blocks if item.get("type") == "text"])
    effective_dashboard_count = len([item for item in content_blocks if item.get("type") == "kpi_dashboard"])

    budget["graficas"] = max(int(budget.get("graficas", 0) or 0), effective_chart_count)
    budget["tabla_principal"] = max(int(budget.get("tabla_principal", 0) or 0), min(2, effective_table_count))
    budget["hallazgos"] = max(int(budget.get("hallazgos", 0) or 0), min(4, effective_text_count))
    budget["genericas"] = max(int(budget.get("genericas", 0) or 0), min(4, effective_dashboard_count))
    return budget


def run_content_quality_checks(data, content_blocks):
    checked_blocks = []
    issues = []
    for index, block in enumerate(content_blocks, start=1):
        block_type = block.get("type")
        title = clean_text(block.get("title"), 80) or f"Bloque {index}"
        if block_type == "chart":
            chart = block.get("chart") or {}
            if not is_valid_chart(chart):
                issues.append({"block": index, "title": title, "issue": "grafica invalida"})
                continue
            truth = block.get("traceability") or {}
            if (truth.get("confidence_score") or 0) < 0.58:
                issues.append({"block": index, "title": title, "issue": "grafica descartada por baja confianza"})
                continue
        elif block_type == "table":
            if not block.get("headers") or not block.get("rows"):
                issues.append({"block": index, "title": title, "issue": "tabla sin estructura valida"})
                continue
        elif block_type == "table_summary":
            if not safe_list(block.get("lines")):
                issues.append({"block": index, "title": title, "issue": "resumen tabular vacio"})
                continue
        elif block_type == "kpi_dashboard":
            if not safe_list(block.get("kpis")) and not safe_list(block.get("ranking")):
                issues.append({"block": index, "title": title, "issue": "dashboard sin contenido util"})
                continue
        elif block_type == "text":
            if not safe_list(block.get("lines")):
                issues.append({"block": index, "title": title, "issue": "bloque textual vacio"})
                continue
        checked_blocks.append(block)
    if not checked_blocks:
        fallback_lines = [
            "No se detectaron bloques suficientemente confiables para poblar la presentacion final.",
            "El motor conserva la plantilla y evita mostrar contenido dudoso o inconsistente.",
        ]
        if STRICT_VISUAL_CURATION_MODE and not is_visual_curation_ready(data):
            fallback_lines = [
                "Hermes no entrego una curaduria visual valida para esta ejecucion.",
                "El modo estricto bloquea graficas y tablas ejecutivas hasta recibir seleccion explicita de la IA.",
            ]
        checked_blocks.append({
            "type": "text",
            "title": "Contenido principal",
            "subtitle": "Validacion de calidad",
            "lines": fallback_lines,
        })
    return checked_blocks[:MAX_CONTENT_BLOCKS], issues


def style_text_frame(text_frame, font_size=16, color=COLOR_TEXT, bold=False, align=PP_ALIGN.LEFT):
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def fill_cover(slide, excel_path, data=None):
    plan = get_presentation_design_plan(data or {})
    meta = (plan.get("presentation_meta") or {}) if isinstance(plan, dict) else {}
    title_slide = next((item for item in safe_list(plan.get("slides")) if item.get("type") == "title_slide"), {}) if isinstance(plan, dict) else {}
    title = clean_text(title_slide.get("title_text") or meta.get("title"), 90) or clean_text(Path(excel_path).stem.replace("_", " ").replace("-", " "), 60) or "Reporte Ejecutivo"
    subtitle = clean_text(title_slide.get("subtitle_text") or meta.get("subtitle"), 120)
    today = clean_text(meta.get("date"), 24) or datetime.now().strftime("%d/%m/%Y")
    placeholders = list(slide.placeholders)
    if len(placeholders) >= 1:
        placeholders[0].text = title
        style_text_frame(placeholders[0].text_frame, font_size=24, color=COLOR_WHITE, bold=True)
    if len(placeholders) >= 2:
        placeholders[1].text = today
        style_text_frame(placeholders[1].text_frame, font_size=12, color=COLOR_WHITE, bold=False, align=PP_ALIGN.RIGHT)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1.15), Inches(4.15), Inches(10.8), Inches(0.7))
        subtitle_box.text_frame.text = subtitle
        style_text_frame(subtitle_box.text_frame, font_size=16, color=COLOR_WHITE, bold=False, align=PP_ALIGN.CENTER)


def fill_closing_slide(slide, data):
    plan = get_presentation_design_plan(data)
    slides = safe_list((plan or {}).get("slides"))
    closing = next((item for item in slides if item.get("type") == "closing_slide"), {})
    title = "Gracias"
    farewell = "Gracias por su atencion."
    secondary_line = clean_text(closing.get("title"), 100)
    if not secondary_line or normalize_header_key(secondary_line) in {"gracias", "cierre ejecutivo"}:
        secondary_line = "Quedo a disposicion para profundizar en la evidencia presentada."

    overlay = slide.shapes.add_shape(1, Inches(0.9), Inches(1.1), Inches(11.5), Inches(4.95))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLOR_WHITE
    overlay.line.color.rgb = COLOR_LINE
    overlay.line.width = Pt(0.6)

    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.45), Inches(10.6), Inches(0.5))
    title_box.text_frame.text = title
    style_text_frame(title_box.text_frame, font_size=22, color=COLOR_BLUE, bold=True, align=PP_ALIGN.LEFT)

    summary_box = slide.shapes.add_textbox(Inches(1.25), Inches(2.2), Inches(10.2), Inches(0.9))
    summary_box.text_frame.word_wrap = True
    summary_box.text_frame.text = compact_line(farewell, 200)
    style_text_frame(summary_box.text_frame, font_size=18, color=COLOR_TEXT, bold=False, align=PP_ALIGN.CENTER)

    secondary_box = slide.shapes.add_textbox(Inches(1.45), Inches(3.35), Inches(9.8), Inches(0.6))
    secondary_box.text_frame.word_wrap = True
    secondary_box.text_frame.text = compact_line(secondary_line, 140)
    style_text_frame(secondary_box.text_frame, font_size=12, color=COLOR_BLUE, bold=False, align=PP_ALIGN.CENTER)


def fill_description_slide(slide, data, excel_path):
    cards = build_slide_2_cards(data, excel_path)
    badge_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]

    for index, card in enumerate(cards):
        if index < len(badge_shapes):
            badge_shapes[index].text = card["badge"]
            style_text_frame(badge_shapes[index].text_frame, font_size=20, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

        box = slide.shapes.add_textbox(card["x"], card["y"], card["w"], card["h"])
        box.text_frame.word_wrap = True
        box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        box.text_frame.margin_left = Pt(4)
        box.text_frame.margin_right = Pt(4)
        box.text_frame.margin_top = Pt(2)
        box.text_frame.margin_bottom = Pt(2)
        box.text_frame.clear()

        p1 = box.text_frame.paragraphs[0]
        p1.text = card["title"]
        p1.alignment = PP_ALIGN.LEFT
        p1.space_after = Pt(3)
        for run in p1.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(15.5)
            run.font.bold = True
            run.font.color.rgb = COLOR_BLUE

        for line_index, line in enumerate([clean_text(part) for part in str(card["body"]).splitlines() if clean_text(part)]):
            paragraph = box.text_frame.add_paragraph()
            paragraph.text = line
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(1 if line_index == 0 else 0)
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(10.6)
                run.font.bold = False
                run.font.color.rgb = COLOR_TEXT


def fill_priority_slide(slide, data, excel_path):
    text_box = slide.shapes.add_textbox(
        SLIDE3_ZONE["x"],
        SLIDE3_ZONE["y"],
        SLIDE3_ZONE["w"],
        SLIDE3_ZONE["h"],
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Pt(14)
    text_frame.margin_right = Pt(14)
    text_frame.margin_top = Pt(14)
    text_frame.margin_bottom = Pt(12)
    text_frame.clear()
    sections = build_slide_3_sections(data, excel_path)
    for section_index, (header, bullets) in enumerate(sections):
        header_paragraph = text_frame.paragraphs[0] if section_index == 0 else text_frame.add_paragraph()
        header_paragraph.text = header
        header_paragraph.alignment = PP_ALIGN.LEFT
        header_paragraph.space_after = Pt(4)
        for run in header_paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = COLOR_ORANGE

        for bullet in bullets[:SLIDE3_MAX_BULLETS]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"- {clean_text(bullet)}"
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(2)
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(10.2)
                run.font.bold = False
                run.font.color.rgb = COLOR_WHITE


def add_content_title(slide, title, subtitle):
    # Franja vertical decorativa izquierda (acento de color)
    accent_stripe = slide.shapes.add_shape(
        1,
        SLIDE4["title_x"] - Inches(0.12),
        SLIDE4["title_y"] - Inches(0.02),
        Inches(0.048),
        Inches(0.56),
    )
    accent_stripe.fill.solid()
    accent_stripe.fill.fore_color.rgb = RGBColor(0, 87, 184)
    accent_stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(
        SLIDE4["title_x"],
        SLIDE4["title_y"],
        SLIDE4["title_w"],
        SLIDE4["title_h"],
    )
    title_frame = title_box.text_frame
    title_frame.text = clean_text(title)
    style_text_frame(title_frame, font_size=22, color=RGBColor(8, 40, 90), bold=True)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            SLIDE4["title_x"],
            SLIDE4["title_y"] + Inches(0.36),
            SLIDE4["title_w"],
            Inches(0.22),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = clean_text(subtitle)
        style_text_frame(subtitle_frame, font_size=9.5, color=RGBColor(90, 110, 140), bold=False)



def add_content_footer(slide, excel_path, index, total):
    footer = slide.shapes.add_textbox(
        SLIDE4["footer_x"],
        SLIDE4["footer_y"],
        SLIDE4["footer_w"],
        SLIDE4["footer_h"],
    )
    footer.text_frame.text = f"Fuente: {clean_text(Path(excel_path).name)} | Contenido {index}/{total}"
    style_text_frame(footer.text_frame, font_size=8, color=COLOR_SOFT, bold=False)


def render_chart_image(chart, output_path):
    chart_type, pairs = normalize_chart_pairs(chart)
    labels = [label for label, _ in pairs]
    values = [value for _, value in pairs]
    layout = build_chart_layout_profile(labels, values)
    wrapped_labels = [wrap_chart_label(item, max_len=layout["label_max_len"], max_lines=2) for item in labels]
    metric_kind = infer_metric_kind(chart.get("metric_label") or chart.get("titulo"), values)
    # Premium color palette — vibrant, corporate-consultive
    PALETTE_PREMIUM = [
        "#0057B8",  # Azul corporativo profundo
        "#00B0CC",  # Cyan-teal
        "#FFB300",  # Dorado ejecutivo
        "#26C676",  # Verde éxito
        "#EF5350",  # Rojo alerta
        "#7C4DFF",  # Violeta premium
    ]
    palette = PALETTE_PREMIUM

    plt.ioff()
    BG_LIGHT = "#F8FAFB"  # Fondo blanco-azulado premium
    figure_height = 4.6 if chart_type in {"pie", "doughnut"} else (4.7 if layout["count"] <= 4 else 5.0)
    fig = plt.figure(figsize=(8.4, figure_height), dpi=180, facecolor=BG_LIGHT)
    title = clean_text(chart.get("titulo"), 72) or "Grafica"

    if chart_type in {"pie", "doughnut"}:
        grid = fig.add_gridspec(1, 2, width_ratios=[1.05, 0.95], wspace=0.02)
        pie_ax = fig.add_subplot(grid[0, 0])
        legend_ax = fig.add_subplot(grid[0, 1])
        legend_ax.axis("off")
        wedge_width = 0.42 if chart_type == "doughnut" else 0.56
        wedges, _ = pie_ax.pie(
            values,
            labels=None,
            startangle=92,
            colors=palette[: len(values)],
            wedgeprops={"width": wedge_width, "linewidth": 2, "edgecolor": "white"},
        )
        pie_ax.set(aspect="equal")
        total_value = sum(values)
        total_label = format_value(total_value, metric_kind, compact=True)
        total_len = len(str(total_label))
        total_font = 16
        total_y = -0.08
        if chart_type == "doughnut":
            if total_len >= 13:
                total_font = 11.5
                total_y = -0.04
            elif total_len >= 10:
                total_font = 13
                total_y = -0.05
            elif total_len >= 8:
                total_font = 14
                total_y = -0.06
        pie_ax.text(0, 0.12, "Total", ha="center", va="center", fontsize=9.5, color="#5B6573", fontweight="bold")
        pie_ax.text(0, total_y, total_label, ha="center", va="center", fontsize=total_font, color="#0B5CAB", fontweight="bold")
        legend_ax.set_xlim(0, 1)
        legend_ax.set_ylim(0, 1)
        legend_step = 0.18 if len(pairs) <= 4 else 0.16
        legend_font = 10.2 if len(pairs) <= 4 else 9.6
        for index, (label, value) in enumerate(pairs[:5]):
            share = (value / total_value) * 100 if total_value else 0
            y = 0.82 - (index * legend_step)
            legend_ax.add_patch(Rectangle((0.02, y - 0.035), 0.04, 0.07, facecolor=palette[index], edgecolor="none"))
            legend_ax.text(0.10, y + 0.02, clean_text(label, 30), fontsize=legend_font, color="#243041", fontweight="bold", va="center")
            legend_ax.text(
                0.10,
                y - 0.05,
                f"{format_value(value, metric_kind, compact=False)} · {share:.0f}%",
                fontsize=max(8.6, legend_font - 0.8),
                color="#5B6573",
                va="center",
            )
        pie_ax.set_title(title, loc="left", pad=12, color="#0057B8", fontsize=13.5, fontweight="bold")
    elif chart_type == "line":
        ax = fig.add_subplot(111)
        positions = list(range(len(values)))
        ax.plot(positions, values, color="#0B5CAB", marker="o", linewidth=2.8, markersize=6)
        ax.fill_between(positions, values, color="#0B5CAB", alpha=0.12)
        peak_index = max(range(len(values)), key=lambda index: values[index])
        last_index = len(values) - 1
        ax.scatter([peak_index], [values[peak_index]], color="#F2994A", s=55, zorder=3)
        ax.scatter([last_index], [values[last_index]], color="#5EA33B", s=55, zorder=3)
        min_value = min(values) if values else 0
        max_value = max(values) if values else 0
        span = max_value - min_value
        if span <= 1e-9:
            padding = max(abs(max_value) * 0.16, 1)
        else:
            padding = max(span * 0.18, abs(max_value) * 0.06, 0.5)
        ax.set_ylim(min_value - padding * 0.45, max_value + padding)
        peak_label = format_chart_callout_value(values[peak_index], metric_kind, prefer_compact=layout["prefer_compact_callout"])
        last_label = format_chart_callout_value(values[last_index], metric_kind, prefer_compact=layout["prefer_compact_callout"])
        peak_offset_x = 0
        last_offset_x = 0
        last_above = abs(values[last_index] - values[peak_index]) <= max(padding * 0.35, 0.5)
        if last_index != peak_index and abs(last_index - peak_index) <= 1 and abs(values[last_index] - values[peak_index]) <= max(padding * 0.55, 1):
            peak_offset_x = -16
            last_offset_x = 16
            last_above = False
        draw_line_annotation(ax, peak_index, values[peak_index], peak_label, above=True, offset_x=peak_offset_x)
        if last_index != peak_index:
            draw_line_annotation(
                ax,
                last_index,
                values[last_index],
                last_label,
                above=last_above,
                offset_x=last_offset_x,
            )
        ax.set_xticks(positions)
        ax.set_xticklabels(wrapped_labels, fontsize=layout["label_font"], color="#243041")
        ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=5, min_n_ticks=4))
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda value, _: format_axis_value(value, metric_kind)))
        ax.grid(True, axis="y", linestyle="--", alpha=0.22)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#D9E1EA")
        ax.spines["bottom"].set_color("#D9E1EA")
        ax.tick_params(axis="x", length=0)
        ax.tick_params(axis="y", labelsize=layout["axis_font"], colors="#5B6573")
        ax.set_title(title, loc="left", pad=12, color="#0057B8", fontsize=13.5, fontweight="bold")
    else:
        ax = fig.add_subplot(111)
        ax.set_facecolor(BG_LIGHT)
        positions = list(range(len(values)))
        # Barras premium: primera barra en azul profundo, resto en gradiente cyan
        bar_colors = [PALETTE_PREMIUM[0]] + [PALETTE_PREMIUM[1]] * max(0, len(values) - 1)
        bars = ax.barh(positions, values, color=bar_colors[: len(values)], height=0.58,
                       edgecolor='white', linewidth=0.8)
        ax.set_yticks(positions)
        ax.set_yticklabels(wrapped_labels, fontsize=layout["label_font"], color="#243041")
        ax.invert_yaxis()
        ax.xaxis.set_major_locator(mticker.MaxNLocator(nbins=5, min_n_ticks=4))
        ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda value, _: format_axis_value(value, metric_kind)))
        ax.grid(True, axis="x", linestyle="--", alpha=0.22)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_visible(False)
        ax.spines["bottom"].set_color("#D9E1EA")
        ax.tick_params(axis="x", labelsize=layout["axis_font"], colors="#5B6573")
        ax.tick_params(axis="y", length=0)
        max_value = max(values) if values else 0
        label_lengths = [len(format_chart_callout_value(value, metric_kind, prefer_compact=layout["prefer_compact_callout"])) for value in values]
        label_factor = 0.2 if any(length >= 12 for length in label_lengths) else 0.16
        ax.set_xlim(0, max_value * (1 + label_factor) if max_value else 1)
        for bar, value in zip(bars, values):
            draw_bar_value_label(
                ax,
                bar,
                value,
                max_value,
                format_chart_callout_value(value, metric_kind, prefer_compact=layout["prefer_compact_callout"]),
            )
        ax.set_title(title, loc="left", pad=12, color="#0057B8", fontsize=13.5, fontweight="bold",
                     fontfamily='DejaVu Sans')

    fig.patch.set_facecolor(BG_LIGHT)
    fig.subplots_adjust(left=layout["left_margin"], right=0.965, top=0.86, bottom=layout["bottom_margin"])
    plt.savefig(output_path, facecolor=BG_LIGHT, format="png")
    plt.close("all")


def render_chart_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("source_basis") or block.get("insight") or "Visual generado con datos reales del Excel.")
    summary = summarize_chart_payload(block["chart"])
    detail_lines = build_chart_detail_lines(block["chart"], limit=4)
    trace_caption = build_traceability_caption(block.get("traceability"))
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        image_path = Path(tmp.name)
    try:
        render_chart_image(block["chart"], image_path)
        slide.shapes.add_picture(
            str(image_path),
            SLIDE4["content_x"],
            SLIDE4["content_y"],
            width=Inches(7.25),
            height=Inches(4.95),
        )
    finally:
        if image_path.exists():
            image_path.unlink(missing_ok=True)

    panel_x = SLIDE4["content_x"] + Inches(7.48)
    # Fondo del panel premium con borde sutil
    panel = slide.shapes.add_shape(1, panel_x, SLIDE4["content_y"], Inches(3.47), Inches(4.95))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 253)  # Azul muy claro premium
    panel.line.color.rgb = RGBColor(210, 220, 235)
    panel.line.width = Pt(0.5)

    # Acento superior del panel (barra horizontal azul)
    panel_accent = slide.shapes.add_shape(1, panel_x, SLIDE4["content_y"], Inches(3.47), Inches(0.055))
    panel_accent.fill.solid()
    panel_accent.fill.fore_color.rgb = RGBColor(0, 87, 184)  # Azul profundo
    panel_accent.line.fill.background()

    # Colores de acento por tipo de métrica
    stat_colors = [
        RGBColor(0, 87, 184),    # Azul — categoría
        RGBColor(255, 179, 0),   # Dorado — valor
        RGBColor(38, 198, 118),  # Verde — total
        RGBColor(0, 176, 204),   # Teal — participación
    ]
    stat_specs = [
        ("Categoria principal", summary["top_label"], stat_colors[0]),
        ("Valor destacado",     summary["top_value"],  stat_colors[1]),
        ("Total analizado",     summary["total_value"], stat_colors[2]),
        ("Peso relativo",       summary["share"],       stat_colors[3]),
    ]
    for index, (label, value, color) in enumerate(stat_specs):
        y = SLIDE4["content_y"] + Inches(0.20) + Inches(0.83) * index

        # Card fondo blanco con sombra suave
        stat = slide.shapes.add_shape(1, panel_x + Inches(0.14), y, Inches(3.18), Inches(0.70))
        stat.fill.solid()
        stat.fill.fore_color.rgb = COLOR_WHITE
        stat.line.color.rgb = RGBColor(225, 232, 242)
        stat.line.width = Pt(0.4)

        # Franja de color a la izquierda del card
        stripe = slide.shapes.add_shape(1, panel_x + Inches(0.14), y, Inches(0.055), Inches(0.70))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()

        label_box = slide.shapes.add_textbox(panel_x + Inches(0.28), y + Inches(0.08), Inches(2.78), Inches(0.18))
        label_box.text_frame.text = label
        style_text_frame(label_box.text_frame, font_size=9.0, color=COLOR_SOFT, bold=False)

        value_box = slide.shapes.add_textbox(panel_x + Inches(0.28), y + Inches(0.29), Inches(2.78), Inches(0.30))
        value_box.text_frame.text = clean_text(value)
        style_text_frame(value_box.text_frame, font_size=15.5, color=color, bold=True)

    # ── Panel de insight con badge IA ───────────────────────────────────────
    insight = clean_text(block.get("insight"))  # Definir antes del panel
    insight_panel_y = SLIDE4["content_y"] + Inches(3.55)
    insight_panel = slide.shapes.add_shape(
        1, panel_x + Inches(0.14), insight_panel_y, Inches(3.18), Inches(1.25)
    )
    insight_panel.fill.solid()
    insight_panel.fill.fore_color.rgb = RGBColor(0, 57, 120)  # Azul oscuro premium
    insight_panel.line.fill.background()

    # Badge IA dorado
    badge = slide.shapes.add_textbox(
        panel_x + Inches(0.22), insight_panel_y + Inches(0.07), Inches(0.80), Inches(0.18)
    )
    badge.text_frame.text = "AI"
    style_text_frame(badge.text_frame, font_size=8.5, color=RGBColor(255, 193, 7), bold=True)

    # Texto de insight en blanco
    insight_box = slide.shapes.add_textbox(
        panel_x + Inches(0.22), insight_panel_y + Inches(0.26), Inches(2.90), Inches(0.92)
    )
    insight_box.text_frame.word_wrap = True
    insight_text = clean_text(insight or "La concentracion observada representa una oportunidad de optimizacion detectada automaticamente.")
    insight_box.text_frame.text = insight_text
    style_text_frame(insight_box.text_frame, font_size=9.0, color=COLOR_WHITE, bold=False)

    # Sección de distribución detallada
    note = slide.shapes.add_textbox(
        panel_x + Inches(0.18),
        SLIDE4["content_y"] + Inches(4.82) - Inches(0.42),
        Inches(3.05),
        Inches(0.38),
    )
    note_frame = note.text_frame
    note_frame.word_wrap = True
    note_frame.margin_left = Pt(4)
    note_frame.margin_right = Pt(3)
    note_frame.margin_top = Pt(2)
    note_frame.margin_bottom = Pt(2)
    note_frame.clear()
    detail_header = note_frame.paragraphs[0]
    detail_header.text = "Distribucion observada"
    for run in detail_header.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(9.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 87, 184)


    if trace_caption:
        trace_box = slide.shapes.add_textbox(
            panel_x + Inches(0.18),
            SLIDE4["content_y"] + Inches(4.73),
            Inches(3.05),
            Inches(0.16),
        )
        trace_box.text_frame.text = clean_text(trace_caption, 120)
        style_text_frame(trace_box.text_frame, font_size=7.2, color=COLOR_SOFT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def set_cell_text(cell, text, header=False):
    cell.text = clean_text(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = True
    cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    cell.text_frame.margin_left = Pt(4)
    cell.text_frame.margin_right = Pt(4)
    cell.text_frame.margin_top = Pt(2)
    cell.text_frame.margin_bottom = Pt(2)
    align = PP_ALIGN.CENTER if header else (PP_ALIGN.RIGHT if try_number(text) is not None else PP_ALIGN.LEFT)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(10 if header else 9)
            run.font.bold = header
            run.font.color.rgb = COLOR_WHITE if header else COLOR_TEXT
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_BLUE if header else COLOR_WHITE


def render_table_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    source_rows = block["rows"]
    source_headers = block["headers"]
    basis = derive_table_basis(source_headers, source_rows)
    display_order = []
    if basis.get("label_index") is not None:
        display_order.append(basis["label_index"])
    if basis.get("numeric_index") is not None and basis["numeric_index"] not in display_order:
        display_order.append(basis["numeric_index"])
    display_order.extend(index for index in range(len(source_headers)) if index not in display_order)
    headers = [source_headers[index] for index in display_order]
    rows = [[row[index] if index < len(row) else "" for index in display_order] for row in source_rows]
    remapped_basis = {
        "label_index": display_order.index(basis["label_index"]) if basis.get("label_index") is not None else None,
        "numeric_index": display_order.index(basis["numeric_index"]) if basis.get("numeric_index") is not None else None,
        "dimension_label": basis.get("dimension_label"),
        "metric_label": basis.get("metric_label"),
        "aggregation": basis.get("aggregation"),
    }
    num_rows = len(rows) + 1
    num_cols = len(headers)
    profile = build_table_signal_profile(headers, rows)
    font_size = 8.0 if (profile["avg_text_len"] >= 24 or num_cols >= 6) else 8.8
    trace_caption = build_traceability_caption(block.get("traceability"))
    meta_y = SLIDE4["content_y"]
    meta_shape = slide.shapes.add_shape(1, SLIDE4["content_x"], meta_y, SLIDE4["content_w"], Inches(0.68))
    meta_shape.fill.solid()
    meta_shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    meta_shape.line.color.rgb = COLOR_LINE

    meta_left = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.18), meta_y + Inches(0.08), Inches(7.25), Inches(0.22))
    meta_left.text_frame.text = clean_text(block.get("source_basis") or "Base real del Excel")
    style_text_frame(meta_left.text_frame, font_size=9.2, color=COLOR_BLUE, bold=True)

    fields_text = build_table_meta_summary(headers, rows, basis=remapped_basis)
    meta_right = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(7.6), meta_y + Inches(0.08), Inches(3.15), Inches(0.22))
    meta_right.text_frame.text = fields_text
    style_text_frame(meta_right.text_frame, font_size=8.8, color=COLOR_SOFT, bold=False, align=PP_ALIGN.RIGHT)

    if trace_caption:
        meta_trace = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.18), meta_y + Inches(0.34), SLIDE4["content_w"] - Inches(0.36), Inches(0.18))
        meta_trace.text_frame.text = clean_text(trace_caption, 140)
        style_text_frame(meta_trace.text_frame, font_size=7.6, color=COLOR_SOFT, bold=False)

    table_height = 4.43
    table_shape = slide.shapes.add_table(
        num_rows,
        num_cols,
        SLIDE4["content_x"],
        SLIDE4["content_y"] + Inches(0.77),
        SLIDE4["content_w"],
        Inches(table_height),
    )
    table = table_shape.table
    header_height = Inches(0.46)
    body_height = Inches(max(0.32, (table_height - 0.46) / max(1, len(rows))))
    table.rows[0].height = header_height

    weights = []
    for col_index, header in enumerate(headers):
        sample_values = [clean_text(row[col_index]) for row in rows[:4] if col_index < len(row)]
        max_len = max([len(clean_text(header))] + [len(value) for value in sample_values], default=10)
        weights.append(max(1.0, min(3.7, max_len / 12)))
    total_weight = sum(weights) or 1
    for col_index, weight in enumerate(weights):
        table.columns[col_index].width = Inches(10.95 * (weight / total_weight))

    numeric_index = remapped_basis.get("numeric_index")
    label_index = remapped_basis.get("label_index")
    numeric_values = [
        try_number(row[numeric_index])
        for row in rows
        if numeric_index is not None and numeric_index < len(row)
    ]
    numeric_values = [value for value in numeric_values if value is not None]
    metric_kind = infer_metric_kind(remapped_basis.get("metric_label"), numeric_values)
    max_numeric_value = max((abs(value) for value in numeric_values), default=0)

    for header_index, header in enumerate(headers):
        set_cell_text(table.cell(0, header_index), header, header=True)
        if header_index == numeric_index:
            table.cell(0, header_index).fill.fore_color.rgb = COLOR_ORANGE

    for row_index, row in enumerate(rows, start=1):
        table.rows[row_index].height = body_height
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            numeric_value = try_number(value)
            display_value = value
            if col_index == numeric_index and numeric_value is not None:
                display_value = format_value(numeric_value, metric_kind, compact=False)
            display_value = clean_text(display_value)
            set_cell_text(cell, display_value, header=False)
            base_fill = COLOR_LIGHT if row_index % 2 == 0 else COLOR_WHITE
            if col_index == label_index:
                cell.fill.fore_color.rgb = blend_rgb(base_fill, COLOR_BLUE, 0.05)
            elif col_index == numeric_index and numeric_value is not None and max_numeric_value > 0:
                intensity = max(0.12, min(0.42, abs(numeric_value) / max_numeric_value * 0.42))
                cell.fill.fore_color.rgb = blend_rgb(base_fill, COLOR_ORANGE, intensity)
            else:
                cell.fill.fore_color.rgb = base_fill
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    if col_index == label_index:
                        run.font.bold = True
                        run.font.color.rgb = COLOR_BLUE
                    if col_index == numeric_index and numeric_value is not None:
                        run.font.bold = True

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_table_summary_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    left = SLIDE4["content_x"]
    top = SLIDE4["content_y"]
    width = SLIDE4["content_w"]
    trace_caption = build_traceability_caption(block.get("traceability"))

    basis_chip = slide.shapes.add_shape(1, left, top, Inches(4.3), Inches(0.34))
    basis_chip.fill.solid()
    basis_chip.fill.fore_color.rgb = RGBColor(245, 247, 250)
    basis_chip.line.color.rgb = COLOR_LINE
    basis_text = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.05), Inches(4.02), Inches(0.18))
    basis_text.text_frame.text = clean_text(block.get("source_basis") or "Sintesis basada en datos reales del Excel")
    style_text_frame(basis_text.text_frame, font_size=8.8, color=COLOR_BLUE, bold=True)

    if trace_caption:
        trace_box = slide.shapes.add_textbox(left + Inches(4.5), top + Inches(0.05), width - Inches(4.55), Inches(0.18))
        trace_box.text_frame.text = clean_text(trace_caption, 96)
        style_text_frame(trace_box.text_frame, font_size=7.6, color=COLOR_SOFT, bold=False, align=PP_ALIGN.RIGHT)

    accent = slide.shapes.add_shape(1, left, top + Inches(0.52), Inches(0.18), Inches(4.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ORANGE
    accent.line.color.rgb = COLOR_ORANGE

    summary_box = slide.shapes.add_textbox(
        left + Inches(0.35),
        top + Inches(0.58),
        width - Inches(0.45),
        Inches(4.15),
    )
    text_frame = summary_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Pt(6)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(4)
    text_frame.margin_bottom = Pt(4)
    text_frame.clear()

    header = text_frame.paragraphs[0]
    header.text = "Analisis de la tabla"
    header.alignment = PP_ALIGN.LEFT
    for run in header.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = COLOR_BLUE

    for line in block.get("lines") or []:
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"- {clean_text(line)}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(7)
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(15)
            run.font.color.rgb = COLOR_TEXT

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_kpi_dashboard_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("source_basis") or block.get("subtitle"))
    trace_caption = build_traceability_caption(block.get("traceability"))

    card_y = SLIDE4["content_y"]
    card_h = Inches(1.08)
    card_w = Inches(3.28)
    card_gap = Inches(0.25)
    card_x_positions = [
        SLIDE4["content_x"],
        SLIDE4["content_x"] + card_w + card_gap,
        SLIDE4["content_x"] + (card_w + card_gap) * 2,
    ]
    fills = [COLOR_BLUE, COLOR_ORANGE, RGBColor(67, 104, 43)]

    for index, kpi in enumerate((block.get("kpis") or [])[:3]):
        shape = slide.shapes.add_shape(1, card_x_positions[index], card_y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = COLOR_LINE

        stripe = slide.shapes.add_shape(1, card_x_positions[index], card_y, card_w, Inches(0.09))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = fills[index % len(fills)]
        stripe.line.color.rgb = fills[index % len(fills)]

        label_box = slide.shapes.add_textbox(card_x_positions[index] + Inches(0.16), card_y + Inches(0.14), card_w - Inches(0.3), Inches(0.25))
        label_box.text_frame.text = clean_text(kpi.get("label"), 28)
        style_text_frame(label_box.text_frame, font_size=10.2, color=COLOR_SOFT, bold=False)

        value_box = slide.shapes.add_textbox(card_x_positions[index] + Inches(0.16), card_y + Inches(0.44), card_w - Inches(0.3), Inches(0.42))
        value_box.text_frame.text = clean_text(kpi.get("value"))
        style_text_frame(value_box.text_frame, font_size=21.5, color=fills[index % len(fills)], bold=True)

    rank_box_x = SLIDE4["content_x"]
    rank_box_y = card_y + Inches(1.42)
    rank_box_w = Inches(6.25)
    rank_box_h = Inches(3.35)
    rank_shape = slide.shapes.add_shape(1, rank_box_x, rank_box_y, rank_box_w, rank_box_h)
    rank_shape.fill.solid()
    rank_shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    rank_shape.line.color.rgb = COLOR_LINE

    rank_title = slide.shapes.add_textbox(rank_box_x + Inches(0.18), rank_box_y + Inches(0.12), rank_box_w - Inches(0.3), Inches(0.25))
    rank_title.text_frame.text = "Concentracion de valor"
    style_text_frame(rank_title.text_frame, font_size=12.5, color=COLOR_BLUE, bold=True)

    ranking = block.get("ranking") or []
    max_value = max([item.get("value", 0) for item in ranking], default=1) or 1
    for index, item in enumerate(ranking[:MAX_DASHBOARD_BARS]):
        row_y = rank_box_y + Inches(0.55) + Inches(0.62) * index
        label = clean_text(item.get("label"))
        label_box = slide.shapes.add_textbox(rank_box_x + Inches(0.18), row_y - Inches(0.02), Inches(2.15), Inches(0.42))
        label_box.text_frame.text = label
        label_box.text_frame.word_wrap = True
        label_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        style_text_frame(label_box.text_frame, font_size=9.6 if len(label) > 26 else 10.2, color=COLOR_TEXT, bold=False)

        bar_w = Inches(2.95) * (item.get("value", 0) / max_value)
        bar = slide.shapes.add_shape(1, rank_box_x + Inches(2.45), row_y + Inches(0.03), max(bar_w, Inches(0.18)), Inches(0.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ORANGE if index == 0 else COLOR_BLUE
        bar.line.color.rgb = COLOR_ORANGE if index == 0 else COLOR_BLUE

        value_box = slide.shapes.add_textbox(rank_box_x + Inches(5.38), row_y - Inches(0.02), Inches(0.7), Inches(0.28))
        value_box.text_frame.text = clean_text(item.get("display"))
        style_text_frame(value_box.text_frame, font_size=10.2, color=COLOR_BLUE, bold=True, align=PP_ALIGN.RIGHT)

    insight_box_x = rank_box_x + rank_box_w + Inches(0.32)
    insight_box = slide.shapes.add_shape(1, insight_box_x, rank_box_y, Inches(4.38), rank_box_h)
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(252, 248, 242)
    insight_box.line.color.rgb = RGBColor(234, 205, 177)

    insight_title = slide.shapes.add_textbox(insight_box_x + Inches(0.18), rank_box_y + Inches(0.12), Inches(3.95), Inches(0.25))
    insight_title.text_frame.text = "Hallazgos ejecutivos"
    style_text_frame(insight_title.text_frame, font_size=12.5, color=COLOR_BLUE, bold=True)

    insight_text = slide.shapes.add_textbox(insight_box_x + Inches(0.18), rank_box_y + Inches(0.48), Inches(3.95), Inches(2.58))
    insight_frame = insight_text.text_frame
    insight_frame.word_wrap = True
    insight_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    insight_frame.clear()
    for idx, line in enumerate((block.get("insights") or [])[:3]):
        paragraph = insight_frame.paragraphs[0] if idx == 0 else insight_frame.add_paragraph()
        paragraph.text = f"- {clean_text(line)}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8)
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(11.2)
            run.font.color.rgb = COLOR_TEXT

    if trace_caption:
        trace_box = slide.shapes.add_textbox(insight_box_x + Inches(0.18), rank_box_y + Inches(2.92), Inches(3.95), Inches(0.24))
        trace_box.text_frame.text = clean_text(trace_caption, 110)
        style_text_frame(trace_box.text_frame, font_size=7.8, color=COLOR_SOFT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_boardroom_traffic_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    items = block.get("items") or []
    card_w = Inches(2.55)
    card_h = Inches(3.95)
    gap = Inches(0.22)
    start_x = SLIDE4["content_x"]
    top = SLIDE4["content_y"] + Inches(0.25)

    for index, item in enumerate(items[:4]):
        x = start_x + (card_w + gap) * index
        color = detect_status_color(item.get("status"))
        card = slide.shapes.add_shape(1, x, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT
        card.line.color.rgb = COLOR_LINE

        chip = slide.shapes.add_shape(1, x + Inches(0.18), top + Inches(0.18), Inches(0.56), Inches(0.28))
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.color.rgb = color

        label_box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(0.58), card_w - Inches(0.36), Inches(0.3))
        label_box.text_frame.text = clean_text(item.get("label"))
        style_text_frame(label_box.text_frame, font_size=12.5, color=COLOR_BLUE, bold=True)

        metric_box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(1.05), card_w - Inches(0.36), Inches(0.58))
        metric_box.text_frame.text = clean_text(item.get("metric"))
        style_text_frame(metric_box.text_frame, font_size=23, color=color, bold=True)

        detail_box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(1.9), card_w - Inches(0.36), Inches(1.55))
        detail_frame = detail_box.text_frame
        detail_frame.word_wrap = True
        detail_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        detail_frame.text = clean_text(item.get("detail"))
        style_text_frame(detail_frame, font_size=11.2, color=COLOR_TEXT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_risk_opportunity_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    sections = [
        ("Riesgos", block.get("risks") or [], COLOR_RED, SLIDE4["content_x"]),
        ("Oportunidades", block.get("opportunities") or [], COLOR_GREEN, SLIDE4["content_x"] + Inches(5.55)),
    ]
    top = SLIDE4["content_y"] + Inches(0.15)

    for section_title, cards, color, x in sections:
        panel = slide.shapes.add_shape(1, x, top, Inches(5.1), Inches(4.85))
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_LIGHT
        panel.line.color.rgb = COLOR_LINE

        header = slide.shapes.add_textbox(x + Inches(0.2), top + Inches(0.14), Inches(4.7), Inches(0.3))
        header.text_frame.text = section_title
        style_text_frame(header.text_frame, font_size=13, color=color, bold=True)

        for idx, card in enumerate(cards[:2]):
            card_y = top + Inches(0.55) + Inches(1.98) * idx
            stripe = slide.shapes.add_shape(1, x + Inches(0.18), card_y, Inches(0.12), Inches(1.62))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = color
            stripe.line.color.rgb = color

            title_box = slide.shapes.add_textbox(x + Inches(0.42), card_y + Inches(0.02), Inches(4.45), Inches(0.28))
            title_box.text_frame.text = clean_text(card.get("title"))
            style_text_frame(title_box.text_frame, font_size=11.5, color=COLOR_BLUE, bold=True)

            body_box = slide.shapes.add_textbox(x + Inches(0.42), card_y + Inches(0.34), Inches(4.35), Inches(0.56))
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            body_frame.text = clean_text(card.get("body"))
            style_text_frame(body_frame, font_size=10.5, color=COLOR_TEXT, bold=False)

            action_box = slide.shapes.add_textbox(x + Inches(0.42), card_y + Inches(0.98), Inches(4.35), Inches(0.5))
            action_frame = action_box.text_frame
            action_frame.word_wrap = True
            action_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            action_frame.text = f"Accion: {clean_text(card.get('action'))}"
            style_text_frame(action_frame, font_size=10.1, color=COLOR_SOFT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_before_after_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    left_box = slide.shapes.add_shape(1, SLIDE4["content_x"], SLIDE4["content_y"] + Inches(0.5), Inches(3.2), Inches(2.45))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_LIGHT
    left_box.line.color.rgb = COLOR_LINE

    right_box = slide.shapes.add_shape(1, SLIDE4["content_x"] + Inches(7.2), SLIDE4["content_y"] + Inches(0.5), Inches(3.2), Inches(2.45))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_LIGHT
    right_box.line.color.rgb = COLOR_LINE

    connector = slide.shapes.add_shape(1, SLIDE4["content_x"] + Inches(3.55), SLIDE4["content_y"] + Inches(1.55), Inches(3.3), Inches(0.18))
    connector.fill.solid()
    connector.fill.fore_color.rgb = COLOR_ORANGE
    connector.line.color.rgb = COLOR_ORANGE

    for x, label, value in (
        (SLIDE4["content_x"], block.get("before_label"), block.get("before_value")),
        (SLIDE4["content_x"] + Inches(7.2), block.get("after_label"), block.get("after_value")),
    ):
        label_box = slide.shapes.add_textbox(x + Inches(0.22), SLIDE4["content_y"] + Inches(0.72), Inches(2.75), Inches(0.28))
        label_box.text_frame.text = clean_text(label)
        style_text_frame(label_box.text_frame, font_size=12, color=COLOR_BLUE, bold=True)

        value_box = slide.shapes.add_textbox(x + Inches(0.22), SLIDE4["content_y"] + Inches(1.18), Inches(2.75), Inches(0.6))
        value_box.text_frame.text = clean_text(value)
        style_text_frame(value_box.text_frame, font_size=25, color=COLOR_ORANGE, bold=True)

    delta_color = COLOR_GREEN if "+" in clean_text(block.get("delta")) else COLOR_RED
    delta_box = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(4.25), SLIDE4["content_y"] + Inches(1.05), Inches(1.9), Inches(0.72))
    delta_box.text_frame.text = clean_text(block.get("delta"))
    style_text_frame(delta_box.text_frame, font_size=24, color=delta_color, bold=True, align=PP_ALIGN.CENTER)

    commentary = slide.shapes.add_textbox(SLIDE4["content_x"], SLIDE4["content_y"] + Inches(3.45), SLIDE4["content_w"], Inches(1.0))
    commentary_frame = commentary.text_frame
    commentary_frame.word_wrap = True
    commentary_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    commentary_frame.text = clean_text(block.get("commentary"))
    style_text_frame(commentary_frame, font_size=13, color=COLOR_TEXT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_top5_alerts_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))

    threshold = slide.shapes.add_textbox(SLIDE4["content_x"], SLIDE4["content_y"], SLIDE4["content_w"], Inches(0.24))
    threshold.text_frame.text = clean_text(block.get("threshold"))
    style_text_frame(threshold.text_frame, font_size=10.2, color=COLOR_SOFT, bold=False)

    items = block.get("items") or []
    max_count = len(items) or 1
    max_value = max([try_number(item.get("raw_value")) or 0 for item in items], default=1) or 1
    top = SLIDE4["content_y"] + Inches(0.55)

    for idx, item in enumerate(items[:MAX_BOARDROOM_ITEMS]):
        y = top + Inches(0.86) * idx
        color = detect_status_color(item.get("status"))
        bullet = slide.shapes.add_shape(1, SLIDE4["content_x"], y + Inches(0.06), Inches(0.18), Inches(0.18))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = color
        bullet.line.color.rgb = color

        label_box = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.28), y, Inches(2.7), Inches(0.28))
        label_box.text_frame.text = clean_text(item.get("label"))
        style_text_frame(label_box.text_frame, font_size=11.5, color=COLOR_TEXT, bold=True)

        value_box = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(8.55), y, Inches(1.2), Inches(0.28))
        value_box.text_frame.text = clean_text(item.get("value"))
        style_text_frame(value_box.text_frame, font_size=11.5, color=COLOR_BLUE, bold=True, align=PP_ALIGN.RIGHT)

        raw_value = try_number(item.get("raw_value")) or 0
        bar_w = Inches(5.1) * (raw_value / max_value) if max_value else Inches(0.2)
        bar = slide.shapes.add_shape(1, SLIDE4["content_x"] + Inches(3.0), y + Inches(0.05), max(bar_w, Inches(0.2)), Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color

        detail_box = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.28), y + Inches(0.28), Inches(9.45), Inches(0.24))
        detail_box.text_frame.text = clean_text(item.get("detail"))
        style_text_frame(detail_box.text_frame, font_size=9.4, color=COLOR_SOFT, bold=False)

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_text_slide(slide, block, excel_path, content_index, total_blocks):
    add_content_title(slide, block["title"], block.get("subtitle"))
    lines = expand_lines_for_slide(block.get("lines") or [], max_len=112, max_items=8)
    if len(lines) == 1:
        panel = slide.shapes.add_shape(1, SLIDE4["content_x"], SLIDE4["content_y"] + Inches(0.5), SLIDE4["content_w"], Inches(3.65))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(245, 247, 250)
        panel.line.color.rgb = COLOR_LINE

        panel_text = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.3), SLIDE4["content_y"] + Inches(0.78), SLIDE4["content_w"] - Inches(0.6), Inches(3.0))
        frame = panel_text.text_frame
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.text = clean_text(lines[0])
        style_text_frame(frame, font_size=18, color=COLOR_BLUE, bold=True, align=PP_ALIGN.LEFT)
    elif len(lines) <= 4:
        headline = slide.shapes.add_shape(1, SLIDE4["content_x"], SLIDE4["content_y"], SLIDE4["content_w"], Inches(0.78))
        headline.fill.solid()
        headline.fill.fore_color.rgb = RGBColor(245, 247, 250)
        headline.line.color.rgb = COLOR_LINE

        headline_text = slide.shapes.add_textbox(SLIDE4["content_x"] + Inches(0.18), SLIDE4["content_y"] + Inches(0.14), SLIDE4["content_w"] - Inches(0.36), Inches(0.45))
        headline_text.text_frame.text = clean_text(lines[0] if lines else block.get("subtitle") or block["title"])
        style_text_frame(headline_text.text_frame, font_size=16, color=COLOR_BLUE, bold=True)

        card_w = Inches(5.2)
        card_h = Inches(1.7)
        card_positions = [
            (SLIDE4["content_x"], SLIDE4["content_y"] + Inches(1.02)),
            (SLIDE4["content_x"] + Inches(5.45), SLIDE4["content_y"] + Inches(1.02)),
            (SLIDE4["content_x"], SLIDE4["content_y"] + Inches(2.95)),
            (SLIDE4["content_x"] + Inches(5.45), SLIDE4["content_y"] + Inches(2.95)),
        ]
        card_lines = lines[1:]
        for index, line in enumerate(card_lines):
            if index >= len(card_positions):
                break
            x, y = card_positions[index]
            card = slide.shapes.add_shape(1, x, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = COLOR_LIGHT if index % 2 == 0 else RGBColor(252, 248, 242)
            card.line.color.rgb = COLOR_LINE if index % 2 == 0 else RGBColor(234, 205, 177)

            card_text = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.18), card_w - Inches(0.36), card_h - Inches(0.24))
            card_frame = card_text.text_frame
            card_frame.word_wrap = True
            card_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            card_frame.text = f"- {clean_text(line)}"
            style_text_frame(card_frame, font_size=13.2, color=COLOR_TEXT, bold=False)
    else:
        box = slide.shapes.add_textbox(
            SLIDE4["content_x"],
            SLIDE4["content_y"],
            SLIDE4["content_w"],
            Inches(5.25),
        )
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = Pt(8)
        text_frame.margin_right = Pt(8)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(6)
        text_frame.clear()

        for line_index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.text = f"- {clean_text(line)}"
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(13.2)
                run.font.color.rgb = COLOR_TEXT

    add_content_footer(slide, excel_path, content_index, total_blocks)


def render_content_block(slide, block, excel_path, content_index, total_blocks):
    if block["type"] == "chart":
        render_chart_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "boardroom_traffic":
        render_boardroom_traffic_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "risk_opportunity_cards":
        render_risk_opportunity_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "before_after_compare":
        render_before_after_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "top5_alerts":
        render_top5_alerts_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "kpi_dashboard":
        render_kpi_dashboard_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "table_summary":
        render_table_summary_slide(slide, block, excel_path, content_index, total_blocks)
    elif block["type"] == "table":
        render_table_slide(slide, block, excel_path, content_index, total_blocks)
    else:
        render_text_slide(slide, block, excel_path, content_index, total_blocks)


def move_slide_to_end(prs, slide_index):
    slide_list = prs.slides._sldIdLst
    target = list(slide_list)[slide_index]
    slide_list.remove(target)
    slide_list.append(target)


def build_audit_record(excel_path, output_path, data, content_blocks, quality_issues=None):
    effective_budget = summarize_effective_slide_budget(data, content_blocks)
    return {
        "generatedAt": datetime.now(UTC).isoformat().replace("+00:00", "Z"),
        "inputFile": str(Path(excel_path).resolve()),
        "inputLabel": clean_text(Path(excel_path).name),
        "templateFile": str(TEMPLATE_PATH.resolve()),
        "outputFile": str(Path(output_path).resolve()),
        "outputExists": Path(output_path).exists(),
        "presentationMode": PRESENTATION_VISUAL_MODE,
        "structure": {
            "slide1": "cover",
            "slide2": "description",
            "slide3": "blue-zone-priorities",
            "slide4+": "duplicated-template-content",
            "finalSlide": "closing",
        },
        "sourceSummary": {
            "detectedSheets": safe_list(data.get("metadatos", {}).get("hojas_encontradas")),
            "primarySheet": clean_text((data.get("metadatos", {}) or {}).get("hoja_principal")),
            "chartCount": len(build_chart_blocks(data)),
            "boardroomBlocks": len([item for item in content_blocks if item["type"] in {"boardroom_traffic", "risk_opportunity_cards", "before_after_compare", "top5_alerts"}]),
            "kpiDashboardBlocks": len([item for item in content_blocks if item["type"] == "kpi_dashboard"]),
            "tableContentBlocks": len([item for item in content_blocks if item["type"] == "table"]),
            "tableSummaryBlocks": len([item for item in content_blocks if item["type"] == "table_summary"]),
            "textContentBlocks": len([item for item in content_blocks if item["type"] == "text"]),
            "conclusions": len(unique_texts(data.get("conclusiones"), limit=20, min_len=18)),
            "slideBudget": effective_budget,
            "minimumConfidenceApplied": {
                "chart": 0.58,
                "tableDetail": 0.56,
                "tableSummary": 0.4,
                "dashboard": 0.62,
            },
            "dataQuality": data.get("calidad_datos") or {},
        },
        "qualityChecks": {
            "issuesCount": len(safe_list(quality_issues)),
            "issues": safe_list(quality_issues),
        },
        "aiCurationStatus": get_ai_curation_status(data),
        "presentationRequest": get_request_context(data),
        "designPlanStatus": {
            "available": bool(get_presentation_design_plan(data)),
            "slideCount": len(safe_list((get_presentation_design_plan(data) or {}).get("slides"))),
            "metaTitle": clean_text(((get_presentation_design_plan(data) or {}).get("presentation_meta") or {}).get("title"), 120),
            "source": clean_text((get_presentation_design_plan(data) or {}).get("source"), 40) or ("ai" if data.get("presentation_design_ia") else "none"),
        },
        "contentAudit": [
            {
                "type": item.get("type"),
                "title": clean_text(item.get("title"), 90),
                "sourceBasis": clean_text(item.get("source_basis"), 120),
                "traceability": item.get("traceability") or {},
            }
            for item in content_blocks
        ],
    }


def generate_presentation(excel_path, output_path, user_instructions=None):
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"No se encontro la plantilla requerida: {TEMPLATE_PATH}")

    prompt_preferences = parse_prompt_preferences(user_instructions)
    apply_prompt_palette(prompt_preferences)
    data = preparar_datos_para_slides(excel_path, user_instructions)
    if not isinstance(data.get("presentation_design_ia"), dict) and clean_text((get_request_context(data) or {}).get("prompt")):
        data["prompt_fallback_design_plan"] = build_prompt_fallback_design_plan(data, prompt_preferences, excel_path)
    prs = Presentation(str(TEMPLATE_PATH))
    if len(prs.slides) < 5:
        raise RuntimeError("La plantilla no contiene las 5 diapositivas base requeridas.")

    fill_cover(prs.slides[0], excel_path, data)
    fill_description_slide(prs.slides[1], data, excel_path)
    fill_priority_slide(prs.slides[2], data, excel_path)

    content_blocks, quality_issues = run_content_quality_checks(data, build_content_blocks(data, prompt_preferences=prompt_preferences))
    content_layout = prs.slides[3].slide_layout

    for index, block in enumerate(content_blocks, start=1):
        if index == 1:
            slide = prs.slides[3]
        else:
            slide = prs.slides.add_slide(content_layout)
        render_content_block(slide, block, excel_path, index, len(content_blocks))

    fill_closing_slide(prs.slides[4], data)
    move_slide_to_end(prs, 4)
    prs.save(output_path)

    audit_path = Path(output_path).with_suffix(".audit.json")
    audit_path.write_text(
        json.dumps(build_audit_record(excel_path, output_path, data, content_blocks, quality_issues=quality_issues), indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    return output_path


def main():
    excel_path = sys.argv[1] if len(sys.argv) > 1 else ""
    output_path = sys.argv[2] if len(sys.argv) > 2 else ""

    if not excel_path:
        raise RuntimeError("Debe suministrar un archivo Excel de entrada.")

    excel_resolved = Path(excel_path).resolve()
    if not excel_resolved.exists():
        raise FileNotFoundError(f"No existe el archivo Excel: {excel_resolved}")

    if not output_path:
        safe_name = basename_label(excel_resolved)
        output_path = str(excel_resolved.with_name(f"Presentacion_Plantilla_{safe_name}.pptx"))

    user_instructions = sys.argv[3] if len(sys.argv) > 3 else None

    final_output = generate_presentation(str(excel_resolved), output_path, user_instructions)
    print(f"Presentación creada en: {final_output}")
    print(f"Auditoría creada en: {Path(final_output).with_suffix('.audit.json')}")


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        raise
