"""Tables must NEVER overflow the slide.

A PowerPoint table row height is a *minimum*: rows auto-grow to fit
word-wrapped text, so a long cell (or many wrapping rows) can push the table
off the bottom of the slide. These tests render adversarial tables through the
real renderer and assert — via a conservative wrapped-height estimate — that
every table stays on the slide (and within the content band), while realistic
cell text is never truncated.
"""
import math
from pathlib import Path

import pytest
from pptx import Presentation

from socya_pipeline.renderer import render_pptx

TEMPLATE = Path(__file__).resolve().parent.parent.parent / "Plantilla_Presentacion_Socya (1) (1).pptx"

# Slide is 13.33 x 7.50 in. Content band ends at 6.30in; footer/pagination at 7.05in.
SLIDE_H = 7.50
CONTENT_BOTTOM = 6.30

# Conservative wrap model (matches the renderer's cell margins; slightly under
# the renderer's own safety factor so passing here is a real guarantee).
_K = 0.52          # avg glyph advance as fraction of em
_LINE = 1.2        # line-height factor
_VMARGIN = 0.08    # cell top+bottom margins
_HMARGIN = 0.16    # cell left+right margins


def _text_units(s):
    n = 0.0
    for ch in str(s):
        o = ord(ch)
        wide = (0x1100 <= o <= 0x115F or 0x2E80 <= o <= 0xA4CF or
                0xAC00 <= o <= 0xD7A3 or 0xFF00 <= o <= 0xFF60)
        n += 1.8 if wide else 1.0
    return n


def _cell_lines(text, col_w_in, font_pt):
    usable = max(0.12, col_w_in - _HMARGIN)
    char_w = max(0.012, font_pt / 72.0 * _K)
    cpl = max(1, int(usable / char_w))
    total = 0
    for seg in str(text).split("\n"):
        total += max(1, int(math.ceil(_text_units(seg) / cpl)))
    return max(1, total)


def _table_bottoms(prs):
    """Estimated rendered bottom (inches) of every table in the deck."""
    bottoms = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            top = (shape.top or 0) / 914400.0
            widths = [c.width / 914400.0 for c in tbl.columns]
            height = 0.0
            for r in range(len(tbl.rows)):
                row_h = 0.0
                for c in range(len(tbl.columns)):
                    cell = tbl.cell(r, c)
                    fpt = 10.0
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.size is not None:
                                fpt = run.font.size.pt
                                break
                    cw = widths[c] if c < len(widths) else (sum(widths) / max(1, len(widths)))
                    lines = _cell_lines(cell.text or "", cw, fpt)
                    row_h = max(row_h, lines * fpt / 72.0 * _LINE + _VMARGIN)
                height += row_h
            bottoms.append(top + height)
    return bottoms


def _all_cell_texts(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.columns)):
                        texts.append(tbl.cell(r, c).text or "")
    return texts


def _render(slides, tmp_path, name="t"):
    out = tmp_path / f"{name}.pptx"
    render_pptx(slides, presentation_meta={"title": "T", "subtitle": "S"},
                template_path=TEMPLATE, output_path=out)
    return Presentation(out)


pytestmark = pytest.mark.skipif(not TEMPLATE.exists(), reason="template not available")


def test_wide_name_table_never_overflows(tmp_path):
    """6 columns, 12 rows, with a long-name column (~33 chars). With equal
    column widths the names wrap to 2 lines and the table runs off the slide.
    The renderer must keep it on the slide (and in the content band)."""
    names = ["MARIA FERNANDA GONZALEZ RODRIGUEZ",
             "JUAN SEBASTIAN MARTINEZ HERNANDEZ"] * 6
    rows = [[names[i], "Valle del Cauca", "Aprobado", f"${(i + 1) * 1234567:,}",
             f"{(i * 7) % 100}%", "2024-06-28"] for i in range(12)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": ["Solicitante", "Ciudad", "Estado", "Monto Solicitado",
                        "Tasa", "Fecha"],
            "rows": rows,
            "numeric_col_indices": [3, 4],
        }},
    ]
    prs = _render(slides, tmp_path, "wide")
    bottoms = _table_bottoms(prs)
    assert bottoms, "no table rendered"
    assert max(bottoms) <= SLIDE_H - 0.2, f"table runs off slide: bottom={max(bottoms):.2f}in"
    assert max(bottoms) <= CONTENT_BOTTOM + 0.3, f"table leaves content band: bottom={max(bottoms):.2f}in"
    # Realistic full names must NOT be ellipsized.
    joined = " ".join(_all_cell_texts(prs))
    assert "…" not in joined and "..." not in joined


def test_huge_blob_cell_never_overflows(tmp_path):
    """A pathological multi-thousand-char cell must not push the table off the
    slide (it may be capped — that is the only allowed truncation)."""
    blob = "lorem ipsum dolor sit amet " * 800  # ~21,600 chars
    rows = [[f"Item {i}", blob] for i in range(5)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": ["Concepto", "Observacion"], "rows": rows}},
    ]
    prs = _render(slides, tmp_path, "blob")
    bottoms = _table_bottoms(prs)
    assert bottoms, "no table rendered"
    assert max(bottoms) <= SLIDE_H - 0.2, f"blob table runs off slide: bottom={max(bottoms):.2f}in"


def test_long_sentence_preserved_and_fits(tmp_path):
    """A realistic long sentence (~165 chars) must be shown in full (no
    ellipsis) AND the table must still fit the slide."""
    sentence = ("El procedimiento es operativo, pero requiere mejoras "
                "estructurales en el flujo de aprobacion y validacion documental "
                "para reducir tiempos de procesamiento.")
    rows = [[f"Caso {i}", sentence] for i in range(3)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": ["Caso", "Hallazgo"], "rows": rows}},
    ]
    prs = _render(slides, tmp_path, "sentence")
    bottoms = _table_bottoms(prs)
    assert max(bottoms) <= SLIDE_H - 0.2, f"sentence table runs off slide: bottom={max(bottoms):.2f}in"
    texts = _all_cell_texts(prs)
    assert any(sentence in t for t in texts), "long sentence was altered/truncated"
    assert all("…" not in t for t in texts), "realistic sentence must not be ellipsized"


def test_all_columns_wide_text_never_overflows(tmp_path):
    """Worst realistic case: every one of 6 columns holds long text, 12 rows.
    Proportional widths can't help (all wide), so the renderer must shrink the
    font (and, if needed, the row count) to keep it on the slide."""
    cell = "Departamento del Valle del Cauca - Municipio"
    rows = [[f"{cell} {i}"] * 6 for i in range(12)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": [f"Columna Larga {k}" for k in range(6)], "rows": rows}},
    ]
    prs = _render(slides, tmp_path, "allwide")
    bottoms = _table_bottoms(prs)
    assert max(bottoms) <= SLIDE_H - 0.2, f"all-wide table runs off slide: bottom={max(bottoms):.2f}in"


def test_unbroken_long_tokens_never_overflow(tmp_path):
    """Long unbroken tokens (IDs / URLs with no spaces) character-wrap; the
    table must still stay on the slide."""
    token = "ABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    rows = [[f"{token}-{i}", f"https://example.com/very/long/path/segment/{i}/{token}"]
            for i in range(12)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": ["Codigo", "Enlace"], "rows": rows}},
    ]
    prs = _render(slides, tmp_path, "tokens")
    bottoms = _table_bottoms(prs)
    assert max(bottoms) <= SLIDE_H - 0.2, f"token table runs off slide: bottom={max(bottoms):.2f}in"


def test_cjk_fullwidth_text_never_overflows(tmp_path):
    """Full-width CJK glyphs are ~1.8x wider than Latin; the fit model must
    account for that so CJK tables don't overflow."""
    cjk = "売上高の詳細な分析と地域別の内訳および前年比較データ"  # ~26 full-width chars
    rows = [[f"{cjk}{i}", "売上高", f"¥{(i + 1) * 1234567:,}"] for i in range(12)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "詳細", "data": {
            "headers": ["項目の説明", "区分", "金額"], "rows": rows,
            "numeric_col_indices": [2]}},
    ]
    prs = _render(slides, tmp_path, "cjk")
    bottoms = _table_bottoms(prs)
    assert max(bottoms) <= SLIDE_H - 0.2, f"CJK table runs off slide: bottom={max(bottoms):.2f}in"


def test_headers_and_totals_preserved(tmp_path):
    """Header row, body rows and a totals row must all survive the fit logic."""
    rows = [["Bogota", f"${(i + 1) * 1000:,}"] for i in range(8)]
    slides = [
        {"type": "title", "data": {"title": "T", "subtitle": "S"}},
        {"type": "table", "title": "Detalle", "data": {
            "headers": ["Ciudad", "Total"],
            "rows": rows,
            "totals_row": ["TOTAL", "$36,000"],
            "numeric_col_indices": [1],
        }},
    ]
    prs = _render(slides, tmp_path, "totals")
    # find the table
    tbl = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
    assert tbl is not None
    # header row present (bold), totals row present (bold), body rows present
    assert tbl.cell(0, 0).text == "CIUDAD"  # humanized + upper
    assert len(tbl.rows) >= 3  # header + >=1 body + totals
    last = len(tbl.rows) - 1
    assert tbl.cell(last, 0).text == "TOTAL"
    assert tbl.cell(last, 0).text_frame.paragraphs[0].runs[0].font.bold
    # at least some body rows rendered with data
    body_texts = [tbl.cell(r, 0).text for r in range(1, last)]
    assert any("Bogota" in t for t in body_texts)
