from pathlib import Path
import pytest
from pptx import Presentation
from socya_pipeline.renderer import render_pptx

TEMPLATE = Path(__file__).resolve().parent.parent.parent / "Plantilla_Presentacion_Socya (1) (1).pptx"

@pytest.mark.skipif(not TEMPLATE.exists(), reason="template not available")
def test_renders_minimal_deck(tmp_path):
    """Renderiza 1 portada + 2 contenido. Si el template trae portada (slide
    0) y cierre (Items 2 con ¡Gracias!), el renderer los preserva intactos
    e inyecta el contenido entre ellos. Output esperado:
      [Portada, kpi_row, table, Items 2_Gracias] = 4 slides.
    """
    rendered_slides = [
        {"type": "title", "data": {"title": "Test", "subtitle": "Sub"}},
        {"type": "kpi_row", "data": {"kpis": [{"label": "Total", "value": "12.4M"}]}},
        {"type": "table", "data": {"headers": ["Ciudad", "Total"],
                                     "rows": [["Bogotá", "100"], ["Medellín", "80"]]}},
    ]
    out = tmp_path / "test.pptx"
    render_pptx(rendered_slides,
                presentation_meta={"title": "T", "subtitle": "S"},
                template_path=TEMPLATE, output_path=out)
    assert out.exists() and out.stat().st_size > 1000
    prs = Presentation(out)
    # 3 input slides (title + 2 content) + 1 closing auto-inyectado del
    # template que SIEMPRE va al final = 4 slides totales.
    assert len(prs.slides) == 4
    # La primera slide debe ser la Portada del template (con su picture).
    assert prs.slides[0].slide_layout.name == "Portada"
    # La última slide debe ser Items 2 (¡Gracias! preservado del template).
    assert prs.slides[-1].slide_layout.name == "Items 2"
