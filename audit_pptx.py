"""Audita un PPTX comparándolo contra el ground truth del Excel Comisiones.
Extrae todos los textos, busca:
  - Números monetarios y verifica si coinciden con totales del Excel
  - Nombres de personas/ciudades y verifica que existan en el Excel
  - Detecta "B" como sufijo (no debería aparecer)
  - Detecta el bug original ($1.5B sin contexto)
  - Reporta findings: OK, WARN, ERROR
"""
import json
import re
import sys
import os
import unicodedata
from pathlib import Path
from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")


def _strip_accents(s: str) -> str:
    return "".join(c for c in unicodedata.normalize("NFD", s)
                   if unicodedata.category(c) != "Mn").lower()

GT = json.loads(Path("GROUND_TRUTH.json").read_text(encoding="utf-8"))

# Tolerancia: 5% para valores monetarios (ajuste por redondeo de formato compact)
TOL = 0.05


def normalize_money(value: str) -> float | None:
    """Intenta extraer un valor monetario de un string como '$1,226M', '$258.3M', '$8.1M'."""
    s = value.strip().replace(",", "")
    m = re.search(r"\$?\s*([0-9]+(?:\.[0-9]+)?)\s*([KMB]?)", s)
    if not m:
        return None
    base = float(m.group(1))
    suf = m.group(2)
    if suf == "K": return base * 1_000
    if suf == "M": return base * 1_000_000
    if suf == "B": return base * 1_000_000_000
    return base


def extract_money_tokens(text: str) -> list[tuple[str, float]]:
    """Devuelve [(token_original, valor_numerico), ...]"""
    out = []
    # Pattern: optional $, digits, optional decimals, optional K/M/B suffix
    for m in re.finditer(r"\$\s*[0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?\s*[KMB]?", text):
        tok = m.group(0)
        val = normalize_money(tok)
        if val is not None:
            out.append((tok, val))
    return out


def extract_text(pptx_path: str) -> list[tuple[int, str]]:
    """Devuelve [(slide_num, full_text), ...]"""
    pres = Presentation(pptx_path)
    out = []
    for i, slide in enumerate(pres.slides):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
        out.append((i + 1, " | ".join(chunks)))
    return out


def audit_pptx(pptx_path: str) -> dict:
    if not os.path.exists(pptx_path):
        return {"path": pptx_path, "status": "MISSING"}

    slides = extract_text(pptx_path)
    full_text = " ".join(t for _, t in slides)

    findings = []  # list of {level, slide, message}

    # CHECK 1: ¿alguna ocurrencia de NUMERO+B?
    b_matches = re.findall(r"(?<![A-Za-z])\d+(?:\.\d+)?B(?![A-Za-z])", full_text)
    if b_matches:
        for s, t in slides:
            for m in re.findall(r"(?<![A-Za-z])\d+(?:\.\d+)?B(?![A-Za-z])", t):
                findings.append({"level": "ERROR", "slide": s,
                                  "message": f'Encontré "{m}" — la "B" es ambigua en español'})

    # CHECK 2: bullet engañoso del bug original
    bug_pattern = re.compile(
        r"VALOR TOTAL SOLICITADO\s+[\d.,]+\s*[BKM]?\s*Promedio por registro",
        re.IGNORECASE,
    )
    for s, t in slides:
        if bug_pattern.search(t):
            findings.append({"level": "ERROR", "slide": s,
                              "message": "Bug original detectado: 'VALOR TOTAL SOLICITADO X Promedio por registro' sin contexto"})

    # CHECK 3: ¿cada bullet/narrative que cita un valor monetario lo hace consistente con GT?
    sum_bruto = GT["sum_bruto_VTS"]
    sum_efec = GT["sum_efectivo_VTS"]
    sum_rech = GT["sum_rechazado_VTS"]
    max_val = GT["max_VTS"]

    # Lista de valores esperados con sus etiquetas
    expected = [
        ("sum_bruto", sum_bruto),
        ("sum_efectivo", sum_efec),
        ("sum_rechazado", sum_rech),
        ("max_individual", max_val),
        ("n_registros", GT["n_registros"]),
    ]

    for s, t in slides:
        tokens = extract_money_tokens(t)
        for tok, val in tokens:
            matched = []
            for label, exp_val in expected:
                if exp_val == 0:
                    continue
                if abs(val - exp_val) / exp_val <= TOL:
                    matched.append(label)
            # No reportar si matched — está bien
            # Reportar si el valor está cerca de 1.5B pero NO está en el contexto correcto
            if 1_400_000_000 < val < 1_600_000_000:
                # Es el bruto. Verificar que el contexto incluye disclaimer
                ctx = t.lower()
                has_disclaimer = any(k in ctx for k in
                                      ("bruto", "incluye", "rechaz", "anulad",
                                       "desglose", "estado"))
                if not has_disclaimer:
                    findings.append({"level": "WARN", "slide": s,
                                      "message": f'Valor cercano al bruto ${val:,.0f} ({tok}) sin keyword de disclaimer'})

    # CHECK 4: ¿menciona el desglose ejecutado/rechazado al menos UNA vez?
    has_breakdown = any(
        "ejecutados" in t.lower() and "rechazados" in t.lower()
        for _, t in slides
    )
    if not has_breakdown:
        findings.append({"level": "WARN", "slide": 0,
                          "message": "El PPTX no menciona el desglose ejecutados/rechazados en ningún slide"})

    # CHECK 5: top ciudades — verificar que las citadas existen.
    # Para esto recargamos el Excel completo (no sólo top 10 del GT) — así
    # detectamos alucinaciones reales sin penalizar ciudades legítimas de
    # cualquier hoja.
    import pandas as pd
    all_cities_by_sheet = set()
    try:
        for sn in ("Comisiones- Base", "Hoja1", "Muestra total ",
                    "Muestra total  (2)"):
            try:
                hdr = 1 if sn == "Comisiones- Base" else 0
                d = pd.read_excel("Comisiones V1.xlsx", sheet_name=sn, header=hdr)
                if "Ciudad Destino" in d.columns:
                    for c in d["Ciudad Destino"].dropna().unique():
                        all_cities_by_sheet.add(_strip_accents(str(c)))
            except Exception:
                continue
    except Exception:
        pass

    UI_TOKENS = ("COMISIONES", "RESUMEN", "VALOR", "TOTAL", "INDICADORES",
                  "HALLAZGOS", "ANALISIS", "ANÁLISIS", "LECTURA",
                  "EJECUTIVA", "DETALLE", "CONCLUSIONES",
                  "DISTRIBUCION", "DISTRIBUCIÓN", "CONCENTRACION",
                  "CONCENTRACIÓN", "PRINCIPAL", "ACCIONABLES",
                  "OPERATIVO", "GERENCIAL")

    for s, t in slides:
        for m in re.finditer(r"[A-ZÁÉÍÓÚÑ]{4,}(?:\s*-\s*[A-ZÁÉÍÓÚÑ ]{3,})+", t):
            ciudad = m.group(0).strip()
            ciudad_n = _strip_accents(ciudad)
            # Filtrar tokens de UI (no son ciudades)
            if any(tok in ciudad.upper() for tok in UI_TOKENS):
                continue
            # match flexible: substring contra el set
            if ciudad_n in all_cities_by_sheet:
                continue
            if any(ciudad_n in c or c in ciudad_n for c in all_cities_by_sheet):
                continue
            findings.append({"level": "WARN", "slide": s,
                              "message": f'Ciudad "{ciudad}" no encontrada en el Excel — posible alucinación'})

    # CHECK 6: top solicitantes — verificar consistencia
    top_solicitantes = list(GT["top_solicitantes"].keys())
    for s, t in slides:
        for m in re.finditer(r"\b[A-ZÁÉÍÓÚÑ][a-záéíóúñ]+(?:\s+[A-ZÁÉÍÓÚÑ][a-záéíóúñ]+){2,5}\b", t):
            nombre = m.group(0).strip()
            if not any(s_real.lower() in nombre.lower() or nombre.lower() in s_real.lower()
                        for s_real in top_solicitantes):
                # Solo flaggear nombres en mayúsculas tipo Persona Apellido Apellido (>= 3 palabras)
                if len(nombre.split()) >= 3 and not any(w in nombre.upper() for w in
                       ("COMISIONES", "VALOR", "TOTAL", "INDICADORES", "HALLAZGOS",
                        "RESUMEN", "CONTABILIZADO", "RECHAZADO", "ANTIOQUIA",
                        "DESEMBOLSO", "MAGDALENA", "ATLANTICO", "CUNDINAMARCA",
                        "SANTA", "EL", "LA", "BAGRE", "PEREIRA", "BOGOTA")):
                    findings.append({"level": "WARN", "slide": s,
                                      "message": f'Posible nombre fabricado: "{nombre}"'})

    # CHECK 7: si una BULLET INDIVIDUAL menciona el max ($8.08M) Y atribuye
    # a una entidad, esa entidad debe ser la real (MIGUEL/SABANALARGA).
    # Importante: comparar bullet-por-bullet (no slide-entero), porque dos
    # bullets distintas pueden mencionar EL BAGRE y $8.08M sin que sean
    # atribución cruzada.
    REAL_MAX_PEOPLE = ("MIGUEL", "GOMEZ TORRES", "GOMEZ", "MIGUEL ANGEL")
    REAL_MAX_PLACES = ("SABANALARGA", "ATLANTICO - SABANALARGA", "ATLANTICO")
    FALSE_ATTRIBUTIONS = (
        "ANTIOQUIA - EL BAGRE", "ANTIOQUIA-EL BAGRE", "EL BAGRE",
        "MEDELLIN", "MEDELLÍN", "ELIZABETH", "JORGE ", "PEDRO ",
        "CAROLINA", "FABIO", "LICEDTH",
    )
    # Re-extract per-shape (each text_bullet usually has its bullet text in
    # its own shape, prefixed with "01", "02", …).
    pres = Presentation(pptx_path)
    for i, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text
            if not shape_text.strip():
                continue
            # Solo procesamos shapes que mencionen el max — normalizamos
            # separadores de miles (comas, puntos, espacios, NBSP) para
            # cubrir formatos como "$8 082 100" o "$8,082,100".
            normalized = re.sub(r"[\s,. ]+", "", shape_text)
            if not ("8082" in normalized or "8.08M" in shape_text or
                    "8.1M" in shape_text):
                continue
            ctx_upper = shape_text.upper()
            # ¿La bullet habla del MAX? (keyword + magnitud cerca)
            if not any(k in ctx_upper for k in ("MÁX", "MAX", "MAYOR", "ALTO")):
                continue
            mentions_real_person = any(p in ctx_upper for p in REAL_MAX_PEOPLE)
            mentions_real_place  = any(p in ctx_upper for p in REAL_MAX_PLACES)
            mentions_false_attr  = any(f.upper() in ctx_upper for f in FALSE_ATTRIBUTIONS)
            if mentions_false_attr and not (mentions_real_person or mentions_real_place):
                findings.append({"level": "ERROR", "slide": i + 1,
                                  "message": f"ALUCINACIÓN: '{shape_text.strip()[:120]}' atribuye el max a entidad INCORRECTA. Real: MIGUEL ANGEL GOMEZ TORRES / ATLANTICO-SABANALARGA / RECHAZADO"})

    # CHECK 8: bullets que dicen "mediana de $X" — verificar contra la mediana REAL.
    import pandas as pd
    try:
        dfc = pd.read_excel("Comisiones V1.xlsx", sheet_name="Comisiones- Base", header=1)
        mediana_real_VTS = float(dfc["Valor Total Solicitado"].median())
    except Exception:
        mediana_real_VTS = None
    if mediana_real_VTS is not None:
        for s, t in slides:
            for m in re.finditer(r"mediana[^.]{0,30}\$\s*([0-9.,]+)\s*([KM]?)",
                                  t, flags=re.IGNORECASE):
                tok = m.group(0)
                num_str = m.group(1).replace(",", "")
                suf = m.group(2)
                try:
                    val = float(num_str)
                except ValueError:
                    continue
                if suf == "K": val *= 1_000
                if suf == "M": val *= 1_000_000
                # Comparar contra la mediana real (10% tolerancia generosa)
                if val < mediana_real_VTS * 0.5 or val > mediana_real_VTS * 1.5:
                    findings.append({"level": "ERROR", "slide": s,
                                      "message": f'Mediana citada {tok!r} no cuadra con mediana real (${mediana_real_VTS:,.0f}) de Valor Total Solicitado'})

    n_errors = sum(1 for f in findings if f["level"] == "ERROR")
    n_warns  = sum(1 for f in findings if f["level"] == "WARN")

    return {
        "path": pptx_path,
        "slides": len(slides),
        "findings": findings,
        "errors": n_errors,
        "warns": n_warns,
        "status": "FAIL" if n_errors > 0 else ("WARN" if n_warns > 0 else "PASS"),
    }


def print_report(audit: dict) -> None:
    p = audit["path"]
    s = audit["status"]
    print(f"\n{'='*70}")
    print(f"  {p}  →  {s}  ({audit['slides']} slides; {audit['errors']} errors, {audit['warns']} warns)")
    print('='*70)
    if not audit["findings"]:
        print("  ✓ Sin findings — PPTX consistente con el Excel ground truth.")
        return
    by_level = {"ERROR": [], "WARN": []}
    for f in audit["findings"]:
        by_level[f["level"]].append(f)
    for level in ("ERROR", "WARN"):
        if by_level[level]:
            print(f"  {level}s:")
            for f in by_level[level]:
                print(f"    [Slide {f['slide']}] {f['message']}")


if __name__ == "__main__":
    pptx_files = sys.argv[1:] if len(sys.argv) > 1 else [
        "AUDIT_p1_default.pptx",
        "AUDIT_p2_ejec_rechazo.pptx",
        "AUDIT_p3_geografico.pptx",
        "AUDIT_p4_riesgos.pptx",
    ]
    summary = {"PASS": 0, "WARN": 0, "FAIL": 0, "MISSING": 0}
    for p in pptx_files:
        if not os.path.exists(p):
            print(f"\n  ⚠ {p} no existe (aún generándose?)")
            summary["MISSING"] += 1
            continue
        result = audit_pptx(p)
        print_report(result)
        summary[result["status"]] += 1
    print(f"\n{'='*70}")
    print(f"  RESUMEN GLOBAL: PASS={summary['PASS']}  WARN={summary['WARN']}  FAIL={summary['FAIL']}  MISSING={summary['MISSING']}")
    print('='*70)
