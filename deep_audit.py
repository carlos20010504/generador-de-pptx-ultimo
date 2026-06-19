"""Auditoría profunda slide-por-slide. Para cada PPTX:
  - Lista cada slide con su tipo, título, y un resumen del contenido.
  - Extrae números (con K/M suffix) y nombres propios.
  - Cross-check cada número contra estadísticas del Excel.
  - Marca FACT vs SUSPECT vs UNVERIFIED.
"""
import sys, re, json, unicodedata
from pathlib import Path
from pptx import Presentation
import pandas as pd

sys.stdout.reconfigure(encoding="utf-8")

GT = json.loads(Path("GROUND_TRUTH.json").read_text(encoding="utf-8"))

# Use the pipeline's own parser — it does proper header detection,
# dtype inference, and produces the same SheetData the validator sees.
from socya_pipeline.parser import parse_workbook
_WB = parse_workbook("Comisiones V1.xlsx")
ALL_SHEETS_DATA = {s.name: s for s in _WB.sheets}

# Cargar dfs (header auto via parser) — usamos los stats de ColumnData
# que ya están bien computados.
def _load_df_for_sheet(sheet_name: str):
    """Best-effort load — for column-level value lookups."""
    try:
        df = pd.read_excel("Comisiones V1.xlsx", sheet_name=sheet_name, header=0)
        if sum(1 for c in df.columns if str(c).startswith("Unnamed")) / max(1, len(df.columns)) > 0.5:
            df = pd.read_excel("Comisiones V1.xlsx", sheet_name=sheet_name, header=1)
        return df
    except Exception:
        return None

ALL_SHEETS = {(sn, 0): _load_df_for_sheet(sn) for sn in [s.name for s in _WB.sheets] if _load_df_for_sheet(sn) is not None}


def _strip(s):
    return "".join(c for c in unicodedata.normalize("NFD", str(s))
                   if unicodedata.category(c) != "Mn").lower().strip()


def normalize_money(text):
    """Find all money tokens in text → list of (raw_token, value_float).
    Handles 1,852 / 1.852 / 1 852 / 1852 / 1.5M / $258.3M / $8 082 100."""
    out = []
    pattern = (
        r"\$?\s*"
        r"((?:[0-9]{1,3}(?:[,. ][0-9]{3})+|[0-9]+)"  # int with optional thousand-seps
        r"(?:\.[0-9]+)?"                              # optional decimal
        r")"
        r"\s*([KM]?)\b"
    )
    for m in re.finditer(pattern, text):
        raw = m.group(0).strip()
        val_str = m.group(1).replace(" ", "").replace(" ", "")
        # If multiple "." (European thousands), strip them; else single "."
        # is treated as decimal point.
        if val_str.count(".") > 1:
            val_str = val_str.replace(".", "")
        else:
            val_str = val_str.replace(",", "")
        try:
            v = float(val_str)
        except ValueError:
            continue
        suf = m.group(2)
        if suf == "K": v *= 1_000
        if suf == "M": v *= 1_000_000
        if v < 10:
            continue
        out.append((raw, v))
    return out


_CANDIDATE_CACHE = None


def _build_candidates() -> list:
    """Pre-compute every numeric statistic (sum/max/min/mean/count) and every
    cell value across every sheet — a number cited in a PPTX is "known" if
    it matches any of these within 5%. Cached because rebuilding it for
    each check is wasteful."""
    global _CANDIDATE_CACHE
    if _CANDIDATE_CACHE is not None:
        return _CANDIDATE_CACHE
    cand = []
    # First add the parser's own column stats — those use the correct
    # header detection and the same column names the validator sees.
    for sheet_data in _WB.sheets:
        for col in sheet_data.columns:
            for stat_name in ("sum", "max", "min", "mean"):
                v = getattr(col, stat_name, None)
                if v is None:
                    continue
                try:
                    cand.append((f"{stat_name}_{sheet_data.name}_{col.name}", float(v)))
                except (TypeError, ValueError):
                    continue
            cand.append((f"nunique_{sheet_data.name}_{col.name}",
                         float(col.n_unique)))
            for tv in (col.top_values or []):
                try:
                    label, count = tv[0], tv[1]
                    cand.append((f"top_{sheet_data.name}_{col.name}={label}",
                                 float(count)))
                except (TypeError, IndexError):
                    continue
        # Status breakdowns precomputed by the parser
        for col_name, bd in (sheet_data.status_breakdowns or {}).items():
            for k in ("positive_sum", "negative_sum", "pending_sum",
                       "unknown_sum", "gross_sum"):
                v = bd.get(k)
                if v is not None:
                    cand.append((f"{k}_{sheet_data.name}_{col_name}", float(v)))
            for k in ("positive_count", "negative_count", "pending_count",
                       "unknown_count", "gross_count"):
                v = bd.get(k)
                if v is not None:
                    cand.append((f"{k}_{sheet_data.name}_{col_name}", float(v)))

    # Also add raw per-sheet stats computed from df (for sheets the parser
    # may have skipped due to header issues).
    for (sn, _), df in ALL_SHEETS.items():
        # Row counts and column counts
        cand.append((f"rows_{sn}", float(len(df))))
        # Bucket counts and sums for currency columns with Estado
        from socya_pipeline.insights import (
            classify_status_value, is_status_filter_column,
        )
        status_col = None
        for c in df.columns:
            top = df[c].value_counts().head(8).reset_index().values.tolist()
            if is_status_filter_column(str(c), top):
                status_col = c
                break
        # Iterate numeric columns
        for c in df.columns:
            try:
                series = pd.to_numeric(df[c], errors="coerce").dropna()
            except Exception:
                continue
            if series.empty:
                continue
            cand.append((f"sum_{sn}_{c}", float(series.sum())))
            cand.append((f"max_{sn}_{c}", float(series.max())))
            cand.append((f"min_{sn}_{c}", float(series.min())))
            cand.append((f"mean_{sn}_{c}", float(series.mean())))
            cand.append((f"median_{sn}_{c}", float(series.median())))
            cand.append((f"count_{sn}_{c}", float(series.shape[0])))
            # If we found a status column on this sheet, add per-bucket sums.
            if status_col is not None and status_col in df.columns:
                buckets = {"positive": 0.0, "negative": 0.0, "pending": 0.0}
                bcounts = {"positive": 0, "negative": 0, "pending": 0}
                for st_val, sub in df.groupby(status_col, dropna=False):
                    bk = classify_status_value(st_val)
                    if bk not in buckets:
                        continue
                    try:
                        s2 = pd.to_numeric(sub[c], errors="coerce").dropna()
                        buckets[bk] += float(s2.sum())
                        bcounts[bk] += int(s2.shape[0])
                    except Exception:
                        continue
                for bk, v in buckets.items():
                    if v > 0:
                        cand.append((f"{bk}_sum_{sn}_{c}", v))
                for bk, v in bcounts.items():
                    if v > 0:
                        cand.append((f"{bk}_count_{sn}_{c}", float(v)))
        # Per-category counts for categorical cols
        for c in df.columns:
            try:
                vc = df[c].value_counts().head(20)
                for label, count in vc.items():
                    cand.append((f"cat_count_{sn}_{c}={label}", float(count)))
            except Exception:
                continue
    _CANDIDATE_CACHE = cand
    return cand


def check_number_against_excel(value: float, tol: float = 0.05) -> str:
    """Return a label if value matches a known stat or cell, else 'unknown'."""
    if 0 <= value <= 100 and value == int(value):
        return "percent_or_small_count"
    cands = _build_candidates()
    for label, ref in cands:
        if ref == 0:
            continue
        if abs(value - ref) / max(abs(ref), 1) <= tol:
            return label
    # Year-like 2020-2026
    if 2000 <= value <= 2030:
        return "year"
    return "unknown"


def audit_slide(idx: int, slide) -> dict:
    """Produce a structured audit of one slide."""
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                chunks.append(t)
    full_text = "\n".join(chunks)
    numbers = normalize_money(full_text)
    classified = {}
    for raw, v in numbers:
        lbl = check_number_against_excel(v)
        classified.setdefault(lbl, []).append((raw, v))
    return {"idx": idx, "text": full_text, "numbers": numbers, "classified": classified}


def deep_audit(pptx_path: str) -> dict:
    pres = Presentation(pptx_path)
    audits = [audit_slide(i + 1, s) for i, s in enumerate(pres.slides)]
    # Count problems
    unknown_count = sum(len(a["classified"].get("unknown", [])) for a in audits)
    return {"path": pptx_path, "slide_count": len(audits), "audits": audits,
            "unknown_count": unknown_count}


def print_compact(report: dict) -> None:
    print(f"\n{'='*78}")
    print(f"  {report['path']} — {report['slide_count']} slides, {report['unknown_count']} números sin match")
    print('='*78)
    for a in report["audits"]:
        first_line = (a["text"].split("\n")[0] if a["text"] else "(empty)")[:80]
        n = len(a["numbers"])
        unk = len(a["classified"].get("unknown", []))
        flag = "⚠" if unk > 0 else " "
        print(f"  {flag} Slide {a['idx']:2d}  [{n:2d} nums, {unk} unknown]  {first_line}")
        if unk > 0:
            for raw, v in a["classified"]["unknown"]:
                # Compactar tabla data: ignorar si parece valor de tabla
                # (suelen ser muchos)
                pass
            # Mostrar solo los primeros 5 unknowns
            for raw, v in a["classified"]["unknown"][:5]:
                print(f"        ? {raw}  ({v:,.0f})")
            if len(a["classified"]["unknown"]) > 5:
                print(f"        ... y {len(a['classified']['unknown']) - 5} más")


if __name__ == "__main__":
    files = sys.argv[1:] if len(sys.argv) > 1 else [
        "AUDIT3_p1_default.pptx",
        "AUDIT3_p2_ejec_rechazo.pptx",
        "AUDIT3_p3_geografico.pptx",
        "AUDIT3_p4_riesgos.pptx",
    ]
    for f in files:
        r = deep_audit(f)
        print_compact(r)
